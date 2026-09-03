import os
import asyncio
import json
import random
import re
import time
import colorsys
from collections import deque
from datetime import datetime
from aiogram import Bot, Dispatcher, F
from aiogram.filters import Command
from aiogram.types import Message, FSInputFile, ReplyKeyboardMarkup, KeyboardButton, ReplyKeyboardRemove
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup
from aiogram.fsm.storage.memory import MemoryStorage
from openai import AsyncOpenAI
from docx import Document
from docx.shared import Pt as DocxPt, Cm, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_TAB_ALIGNMENT
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import openpyxl
from openpyxl.styles import Font as XlFont, PatternFill, Border, Side, Alignment, Protection
from openpyxl.utils import get_column_letter, column_index_from_string
from openpyxl.worksheet.datavalidation import DataValidation
from openpyxl.formatting.rule import ColorScaleRule
from openpyxl.chart import BarChart, PieChart, LineChart, Reference
from openpyxl.chart.label import DataLabelList
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from PIL import Image, ImageOps, ImageEnhance, ImageChops, ImageStat, ImageFilter

try:
    import cv2
    _FACE_CASCADE = cv2.CascadeClassifier(cv2.data.haarcascades + "haarcascade_frontalface_default.xml")
except Exception:
    # opencv-python может отсутствовать на сервере (тяжёлая зависимость) - в этом случае
    # размытие фона просто откатится на центр кадра вместо детекции лица, без падений.
    cv2 = None
    _FACE_CASCADE = None
import httpx

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
FONT_CANDIDATES = [
    (os.path.join(BASE_DIR, "fonts", "DejaVuSans.ttf"), os.path.join(BASE_DIR, "fonts", "DejaVuSans-Bold.ttf")),
    ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"),
    ("C:/Windows/Fonts/arial.ttf", "C:/Windows/Fonts/arialbd.ttf"),
]


def safe_filename(title: str, fallback: str = "Документ", author: str = None) -> str:
    """Делает из темы/заголовка документа безопасное имя файла:
    убирает символы, которые нельзя использовать в имени файла (Windows/Telegram),
    обрезает длину, чтобы не упереться в лимиты, и подставляет fallback, если
    заголовок пустой. Если передан author (реальное ФИО, не плейсхолдер вида
    "[ФИО студента]") - добавляет его через тире, чтобы файл не выглядел безлико."""
    import re as _re

    def _clean(s):
        s = _re.sub(r'[\\/:*?"<>|\n\r\t]', " ", s)
        return _re.sub(r"\s+", " ", s).strip()

    text = _clean((title or "").strip() or fallback)
    text = text[:70].strip() or fallback
    if author:
        a = _clean(author.strip())
        if a and not a.startswith("["):  # не подставляем невыполненный плейсхолдер
            a = a[:40]
            text = f"{text[:95 - len(a) - 3]} - {a}"
    text = text[:110].strip() or fallback
    return text


def wrap_lines(text, font, size, max_width):
    """Переносит текст по словам, а не по количеству символов,
    чтобы слова не рвались посередине. Слишком длинное слово
    без пробелов режется по ширине, иначе оно вылезет за край PDF."""
    words = (text or "").split()
    lines, cur = [], ""

    def _split_long(word):
        chunks, buf = [], ""
        for ch in word:
            trial = buf + ch
            if pdfmetrics.stringWidth(trial, font, size) <= max_width or not buf:
                buf = trial
            else:
                chunks.append(buf)
                buf = ch
        if buf:
            chunks.append(buf)
        return chunks or [word]

    for word in words:
        parts = _split_long(word) if pdfmetrics.stringWidth(word, font, size) > max_width else [word]
        for part in parts:
            test = f"{cur} {part}".strip()
            if pdfmetrics.stringWidth(test, font, size) <= max_width:
                cur = test
            else:
                if cur:
                    lines.append(cur)
                cur = part
    if cur:
        lines.append(cur)
    return lines


def register_pdf_fonts():
    """Регистрирует шрифт с поддержкой кириллицы для reportlab.
    Сначала пробует шрифт, который лежит рядом со скриптом (папка fonts/),
    затем — типичные системные пути. Если ничего не нашлось, возвращает
    базовые Helvetica-шрифты и печатает предупреждение, чтобы не было
    тихой поломки кириллицы в PDF.
    Повторный вызов безопасен и дешёв: функция сама проверяет, зарегистрирован ли
    шрифт, и не пытается зарегистрировать его повторно (на некоторых версиях
    reportlab повторная регистрация того же имени могла вести себя не предсказуемо -
    проверка ничего не стоит, поэтому оставлена как подстраховка на будущее)."""
    registered = set(pdfmetrics.getRegisteredFontNames())
    if "CyrRegular" in registered and "CyrBold" in registered:
        return "CyrRegular", "CyrBold"
    for regular, bold in FONT_CANDIDATES:
        if os.path.exists(regular) and os.path.exists(bold):
            try:
                if "CyrRegular" not in registered:
                    pdfmetrics.registerFont(TTFont("CyrRegular", regular))
                if "CyrBold" not in registered:
                    pdfmetrics.registerFont(TTFont("CyrBold", bold))
                return "CyrRegular", "CyrBold"
            except Exception as e:
                print("Не удалось зарегистрировать шрифт", regular, ":", e)
    print("ВНИМАНИЕ: шрифт с кириллицей не найден — русский текст в PDF будет нечитаемым. "
          "Положи DejaVuSans.ttf и DejaVuSans-Bold.ttf в папку fonts/ рядом со скриптом.")
    return "Helvetica", "Helvetica-Bold"


BOT_TOKEN = os.getenv("BOT_TOKEN")
XAI_API_KEY = os.getenv("XAI_API_KEY")
REPLICATE_API_TOKEN = os.getenv("REPLICATE_API_TOKEN")
ADMIN_IDS = [909828109]

if not BOT_TOKEN:
    raise RuntimeError("BOT_TOKEN не задан в переменных окружения")
if not XAI_API_KEY:
    raise RuntimeError("XAI_API_KEY не задан в переменных окружения")
if not REPLICATE_API_TOKEN:
    print("ВНИМАНИЕ: REPLICATE_API_TOKEN не задан — генерация изображений будет недоступна")

client = AsyncOpenAI(api_key=XAI_API_KEY, base_url="https://api.x.ai/v1")
bot = Bot(token=BOT_TOKEN)
dp = Dispatcher(storage=MemoryStorage())
USERS_FILE = os.path.join(BASE_DIR, "users.json")
PLAN_LIMITS = {"premium": 15}


def load_users():
    """Читает users.json после рестарта, чтобы не обнулялись язык, тариф и история."""
    try:
        with open(USERS_FILE, "r", encoding="utf-8") as f:
            raw = json.load(f)
        out = {}
        for k, v in raw.items():
            try:
                uid = int(k)
            except (TypeError, ValueError):
                continue
            if not isinstance(v, dict):
                continue
            v["busy"] = False
            v["_counted"] = False
            out[uid] = v
        return out
    except FileNotFoundError:
        return {}
    except Exception as e:
        print("Не удалось прочитать users.json:", e)
        return {}


def _build_users_payload():
    """Собирает снимок устойчивых полей всех пользователей - быстрая синхронная операция
    без обращения к диску, безопасно вызывать прямо из event loop."""
    payload = {}
    for uid, u in users_db.items():
        payload[str(uid)] = {
            "name": u.get("name", ""),
            "plan": u.get("plan", "premium"),
            "generations": int(u.get("generations") or 0),
            "history": list(u.get("history") or [])[-30:],
            "lang": u.get("lang") or "ru",
            "lang_chosen": bool(u.get("lang_chosen")),
        }
    return payload


def _write_users_payload(payload):
    """Блокирующая часть - реальная запись на диск (atomic replace). Вызывается только
    в отдельном потоке через asyncio.to_thread, никогда напрямую из event loop."""
    try:
        tmp = USERS_FILE + ".tmp"
        with open(tmp, "w", encoding="utf-8") as f:
            json.dump(payload, f, ensure_ascii=False, indent=0)
        os.replace(tmp, USERS_FILE)
    except Exception as e:
        print("Не удалось сохранить users.json:", e)


_save_users_lock = asyncio.Lock()


async def _save_users_async():
    # Снимок данных берётся сразу (до ожидания лока), чтобы каждая запись несла
    # самые свежие на момент своего вызова данные. Сам asyncio.Lock не столько
    # защищает от гонки (payload уже неизменяемый снимок), сколько гарантирует,
    # что параллельные вызовы допишутся на диск строго в том порядке, в котором
    # были запланированы - иначе поздний, но более свежий снимок мог бы попасть
    # на диск РАНЬШЕ, а следом его молча затёр бы более старый снимок.
    payload = _build_users_payload()
    async with _save_users_lock:
        await asyncio.to_thread(_write_users_payload, payload)


def save_users():
    """Планирует сохранение в фоновом потоке, не блокируя event loop бота -
    раньше запись всего словаря пользователей выполнялась синхронно прямо в
    обработчике сообщения и тормозила ВСЕХ пользователей бота на время записи,
    а не только того, чьи данные сохраняются. Вызывается только из async-кода
    (все места использования - хендлеры aiogram), поэтому запущенный event loop
    должен быть всегда; RuntimeError - подстраховка на случай вызова вне него,
    чтобы данные в любом случае не потерялись."""
    try:
        asyncio.create_task(_save_users_async())
    except RuntimeError:
        _write_users_payload(_build_users_payload())


users_db = load_users()

_bot_username_cache = {"value": None}


async def get_bot_username() -> str:
    """Юзернейм бота (без @), запрашивается у Telegram один раз при первом обращении
    и кешируется - чтобы не дёргать get_me() на каждую отправку файла."""
    if _bot_username_cache["value"] is None:
        me = await bot.get_me()
        _bot_username_cache["value"] = me.username or "bot"
    return _bot_username_cache["value"]


async def signature_line(lang: str) -> str:
    """Строка-подпись 'Рад был помочь! Ваш, @ИмяБота', добавляется в конец финального
    сообщения после отправки готового файла - как у SaveAsBot."""
    username = await get_bot_username()
    return tr("msg_signature", lang, bot=f"@{username}")


class Form(StatesGroup):
    waiting_mode = State()
    waiting_topic = State()
    waiting_user_text = State()
    waiting_theme = State()
    waiting_slides = State()
    waiting_confirm = State()
    waiting_pres_photo_choice = State()
    waiting_pres_photos = State()
    waiting_extra = State()
    waiting_style = State()
    waiting_word_category = State()
    waiting_word_mode = State()
    waiting_word_kind = State()
    waiting_word_topic = State()
    waiting_word_text = State()
    waiting_word_size = State()
    waiting_word_confirm = State()
    waiting_word_extra = State()
    waiting_excel_category = State()
    waiting_excel_kind = State()
    waiting_excel_mode = State()
    waiting_excel_topic = State()
    waiting_excel_data = State()
    waiting_excel_startup_data = State()
    waiting_excel_confirm = State()
    waiting_excel_extra = State()
    waiting_language = State()
    waiting_pres_content_lang = State()
    waiting_word_content_lang = State()
    waiting_excel_content_lang = State()


# ==================== ЯЗЫКИ / i18n ====================
# 7 языков интерфейса и генерации контента. Выбор языка хранится в users_db[uid]["lang"],
# по умолчанию "ru" (для уже существующих пользователей и новых, пока не выбрали другой).
LANGS = {
    "ru": {"flag": "🇷🇺", "name": "Русский", "prompt_name": "русском"},
    "en": {"flag": "🇬🇧", "name": "English", "prompt_name": "английском"},
    "de": {"flag": "🇩🇪", "name": "Deutsch", "prompt_name": "немецком"},
    "ar": {"flag": "🇸🇦", "name": "العربية", "prompt_name": "арабском"},
    "zh": {"flag": "🇨🇳", "name": "中文", "prompt_name": "китайском"},
    "es": {"flag": "🇪🇸", "name": "Español", "prompt_name": "испанском"},
    "fr": {"flag": "🇫🇷", "name": "Français", "prompt_name": "французском"},
}
RTL_LANGS = {"ar"}

# Переводы интерфейса бота (кнопки, служебные сообщения). Ключ - смысловой id строки,
# значение - словарь {код_языка: текст}. Если перевода для конкретного языка нет -
# t() подставит русский вариант, чтобы бот не падал и не показывал пустоту.
TR = {
    "btn_pres": {"ru": "📊 Сделать презентацию", "en": "📊 Create a presentation", "de": "📊 Präsentation erstellen",
                 "ar": "📊 إنشاء عرض تقديمي", "zh": "📊 制作演示文稿", "es": "📊 Crear una presentación", "fr": "📊 Créer une présentation"},
    "btn_word": {"ru": "📄 Сделать документ Word", "en": "📄 Create a Word document", "de": "📄 Word-Dokument erstellen",
                 "ar": "📄 إنشاء مستند Word", "zh": "📄 制作 Word 文档", "es": "📄 Crear un documento Word", "fr": "📄 Créer un document Word"},
    "btn_excel": {"ru": "📈 Сделать таблицу Excel", "en": "📈 Create an Excel table", "de": "📈 Excel-Tabelle erstellen",
                  "ar": "📈 إنشاء جدول Excel", "zh": "📈 制作 Excel 表格", "es": "📈 Crear una tabla Excel", "fr": "📈 Créer un tableau Excel"},
    "btn_history": {"ru": "📁 Моя история", "en": "📁 My history", "de": "📁 Mein Verlauf",
                    "ar": "📁 سجلي", "zh": "📁 我的历史", "es": "📁 Mi historial", "fr": "📁 Mon historique"},
    "btn_plan": {"ru": "ℹ️ Мой тариф", "en": "ℹ️ My plan", "de": "ℹ️ Mein Tarif",
                 "ar": "ℹ️ خطتي", "zh": "ℹ️ 我的套餐", "es": "ℹ️ Mi plan", "fr": "ℹ️ Mon forfait"},
    "btn_language": {"ru": "🌐 Язык", "en": "🌐 Language", "de": "🌐 Sprache",
                     "ar": "🌐 اللغة", "zh": "🌐 语言", "es": "🌐 Idioma", "fr": "🌐 Langue"},
    "btn_same_as_interface": {"ru": "Как в интерфейсе ({iface_lang})", "en": "Same as interface ({iface_lang})", "de": "Wie die Oberfläche ({iface_lang})",
                              "ar": "نفس لغة الواجهة ({iface_lang})", "zh": "与界面语言相同（{iface_lang}）", "es": "Igual que la interfaz ({iface_lang})", "fr": "Comme l'interface ({iface_lang})"},
    "btn_main_menu": {"ru": "🏠 Главное меню", "en": "🏠 Main menu", "de": "🏠 Hauptmenü",
                      "ar": "🏠 القائمة الرئيسية", "zh": "🏠 主菜单", "es": "🏠 Menú principal", "fr": "🏠 Menu principal"},
    "btn_ai_generate": {"ru": "✨ Сгенерировать с ИИ", "en": "✨ Generate with AI", "de": "✨ Mit KI generieren",
                        "ar": "✨ إنشاء بالذكاء الاصطناعي", "zh": "✨ 用AI生成", "es": "✨ Generar con IA", "fr": "✨ Générer avec l'IA"},
    "btn_own_text": {"ru": "📝 Вставить свой текст", "en": "📝 Paste your own text", "de": "📝 Eigenen Text einfügen",
                     "ar": "📝 إدراج النص الخاص بك", "zh": "📝 粘贴自己的文本", "es": "📝 Pegar tu propio texto", "fr": "📝 Coller votre propre texte"},
    "btn_own_data": {"ru": "📝 Ввести свои данные", "en": "📝 Enter your own data", "de": "📝 Eigene Daten eingeben",
                     "ar": "📝 إدخال بياناتك الخاصة", "zh": "📝 输入自己的数据", "es": "📝 Introducir tus propios datos", "fr": "📝 Saisir vos propres données"},
    "btn_template": {"ru": "📄 Скачать шаблон", "en": "📄 Download template", "de": "📄 Vorlage herunterladen",
                     "ar": "📄 تنزيل القالب", "zh": "📄 下载模板", "es": "📄 Descargar plantilla", "fr": "📄 Télécharger le modèle"},
    "btn_full_version": {"ru": "✅ Делать полную версию", "en": "✅ Make full version", "de": "✅ Vollversion erstellen",
                         "ar": "✅ إنشاء النسخة الكاملة", "zh": "✅ 生成完整版本", "es": "✅ Crear versión completa", "fr": "✅ Créer la version complète"},
    "btn_add_info": {"ru": "➕ Добавить или изменить информацию", "en": "➕ Add or change information", "de": "➕ Informationen hinzufügen/ändern",
                     "ar": "➕ إضافة أو تعديل المعلومات", "zh": "➕ 添加或修改信息", "es": "➕ Añadir o modificar información", "fr": "➕ Ajouter ou modifier des informations"},
    "btn_change_style": {"ru": "🎨 Изменить стиль", "en": "🎨 Change style", "de": "🎨 Stil ändern",
                         "ar": "🎨 تغيير النمط", "zh": "🎨 更改风格", "es": "🎨 Cambiar estilo", "fr": "🎨 Changer le style"},
    "btn_change_topic": {"ru": "✏️ Изменить тему", "en": "✏️ Change topic", "de": "✏️ Thema ändern",
                         "ar": "✏️ تغيير الموضوع", "zh": "✏️ 更改主题", "es": "✏️ Cambiar el tema", "fr": "✏️ Changer le sujet"},
    "btn_change_query": {"ru": "✏️ Изменить запрос", "en": "✏️ Change request", "de": "✏️ Anfrage ändern",
                         "ar": "✏️ تغيير الطلب", "zh": "✏️ 更改请求", "es": "✏️ Cambiar la solicitud", "fr": "✏️ Modifier la demande"},
    "btn_ai_photos": {"ru": "✨ Все фото через ИИ", "en": "✨ All photos via AI", "de": "✨ Alle Fotos per KI",
                      "ar": "✨ كل الصور بالذكاء الاصطناعي", "zh": "✨ 全部用AI生成图片", "es": "✨ Todas las fotos con IA", "fr": "✨ Toutes les photos via l'IA"},
    "btn_own_photos": {"ru": "📷 Добавить свои фото", "en": "📷 Add your own photos", "de": "📷 Eigene Fotos hinzufügen",
                       "ar": "📷 إضافة صورك الخاصة", "zh": "📷 添加自己的照片", "es": "📷 Añadir tus propias fotos", "fr": "📷 Ajouter vos propres photos"},
    "btn_photos_done": {"ru": "✅ Готово, собери презентацию", "en": "✅ Done, build the presentation", "de": "✅ Fertig, Präsentation erstellen",
                        "ar": "✅ تم، أنشئ العرض التقديمي", "zh": "✅ 完成，生成演示文稿", "es": "✅ Listo, crea la presentación", "fr": "✅ Terminé, créez la présentation"},
    "btn_build_doc": {"ru": "✅ Собрать документ", "en": "✅ Build document", "de": "✅ Dokument erstellen",
                      "ar": "✅ إنشاء المستند", "zh": "✅ 生成文档", "es": "✅ Crear documento", "fr": "✅ Créer le document"},
    "btn_build_table": {"ru": "✅ Собрать таблицу", "en": "✅ Build table", "de": "✅ Tabelle erstellen",
                        "ar": "✅ إنشاء الجدول", "zh": "✅ 生成表格", "es": "✅ Crear tabla", "fr": "✅ Créer le tableau"},
    "btn_more_docs": {"ru": "➕ Показать ещё документы", "en": "➕ Show more documents", "de": "➕ Weitere Dokumente anzeigen",
                      "ar": "➕ عرض المزيد من المستندات", "zh": "➕ 显示更多文档", "es": "➕ Mostrar más documentos", "fr": "➕ Afficher plus de documents"},
    "btn_back_categories": {"ru": "⬅️ Назад к категориям", "en": "⬅️ Back to categories", "de": "⬅️ Zurück zu Kategorien",
                            "ar": "⬅️ العودة إلى الفئات", "zh": "⬅️ 返回分类", "es": "⬅️ Volver a categorías", "fr": "⬅️ Retour aux catégories"},
    "btn_8slides": {"ru": "8️⃣ 8 слайдов", "en": "8️⃣ 8 slides", "de": "8️⃣ 8 Folien", "ar": "8️⃣ 8 شرائح", "zh": "8️⃣ 8张幻灯片", "es": "8️⃣ 8 diapositivas", "fr": "8️⃣ 8 diapositives"},
    "btn_12slides": {"ru": "1️⃣2️⃣ 12 слайдов", "en": "1️⃣2️⃣ 12 slides", "de": "1️⃣2️⃣ 12 Folien", "ar": "1️⃣2️⃣ 12 شريحة", "zh": "1️⃣2️⃣ 12张幻灯片", "es": "1️⃣2️⃣ 12 diapositivas", "fr": "1️⃣2️⃣ 12 diapositives"},
    "btn_16slides": {"ru": "1️⃣6️⃣ 16 слайдов", "en": "1️⃣6️⃣ 16 slides", "de": "1️⃣6️⃣ 16 Folien", "ar": "1️⃣6️⃣ 16 شريحة", "zh": "1️⃣6️⃣ 16张幻灯片", "es": "1️⃣6️⃣ 16 diapositivas", "fr": "1️⃣6️⃣ 16 diapositives"},
    "btn_shallow": {"ru": "Поверхностное раскрытие темы", "en": "Brief coverage of the topic", "de": "Oberflächliche Themenbehandlung",
                    "ar": "تغطية سطحية للموضوع", "zh": "简要阐述主题", "es": "Cobertura breve del tema", "fr": "Traitement succinct du sujet"},
    "btn_deep": {"ru": "Полное раскрытие темы", "en": "Full coverage of the topic", "de": "Vollständige Themenbehandlung",
                 "ar": "تغطية كاملة للموضوع", "zh": "全面阐述主题", "es": "Cobertura completa del tema", "fr": "Traitement complet du sujet"},
    "btn_keep_style": {"ru": "👍 Оставить этот стиль", "en": "👍 Keep this style", "de": "👍 Diesen Stil behalten",
                       "ar": "👍 الاحتفاظ بهذا النمط", "zh": "👍 保留此风格", "es": "👍 Mantener este estilo", "fr": "👍 Conserver ce style"},
    "msg_choose_lang": {"ru": "Выбери язык интерфейса и генерации:", "en": "Choose interface and generation language:",
                        "de": "Wähle die Sprache für Oberfläche und Generierung:", "ar": "اختر لغة الواجهة والإنشاء:",
                        "zh": "选择界面和生成内容的语言：", "es": "Elige el idioma de la interfaz y la generación:", "fr": "Choisissez la langue de l'interface et de la génération :"},
    "msg_lang_set": {"ru": "Готово! Дальше всё — интерфейс и документы — будет на этом языке. 👇", "en": "Done! From now on the interface and documents will be in this language. 👇",
                     "de": "Fertig! Ab jetzt sind Oberfläche und Dokumente in dieser Sprache. 👇", "ar": "تم! من الآن ستكون الواجهة والمستندات بهذه اللغة. 👇",
                     "zh": "完成！从现在起界面和文档都将使用该语言。👇", "es": "¡Listo! A partir de ahora la interfaz y los documentos estarán en este idioma. 👇",
                     "fr": "C'est fait ! Désormais l'interface et les documents seront dans cette langue. 👇"},
    "msg_welcome": {
        "ru": "Привет, {name} 👋\n\nЯ собираю красивые презентации, документы Word и таблицы Excel: текст, стиль и оформление — по твоей теме или из твоих данных.\n\nНажми кнопку ниже и начнём.",
        "en": "Hi, {name} 👋\n\nI put together polished presentations, Word documents and Excel tables: text, style and layout — from your topic or your own data.\n\nTap a button below to start.",
        "de": "Hallo, {name} 👋\n\nIch erstelle ansprechende Präsentationen, Word-Dokumente und Excel-Tabellen: Text, Stil und Layout — zu deinem Thema oder aus deinen Daten.\n\nTippe unten auf einen Button, um zu starten.",
        "ar": "مرحباً {name} 👋\n\nأقوم بإعداد عروض تقديمية ومستندات Word وجداول Excel أنيقة: النص والنمط والتصميم — بناءً على موضوعك أو بياناتك.\n\nاضغط الزر أدناه للبدء.",
        "zh": "你好，{name} 👋\n\n我可以帮你制作精美的演示文稿、Word文档和Excel表格：根据你的主题或数据生成文本、风格和排版。\n\n点击下方按钮开始吧。",
        "es": "Hola, {name} 👋\n\nCreo presentaciones, documentos Word y tablas Excel cuidados: texto, estilo y diseño, a partir de tu tema o tus propios datos.\n\nToca un botón abajo para empezar.",
        "fr": "Bonjour {name} 👋\n\nJe crée des présentations, des documents Word et des tableaux Excel soignés : texte, style et mise en page — à partir de votre sujet ou de vos propres données.\n\nAppuyez sur un bouton ci-dessous pour commencer.",
    },
    "msg_main_menu": {"ru": "Главное меню 👇", "en": "Main menu 👇", "de": "Hauptmenü 👇", "ar": "القائمة الرئيسية 👇",
                      "zh": "主菜单 👇", "es": "Menú principal 👇", "fr": "Menu principal 👇"},
    "msg_cancelled": {"ru": "Отменил текущее действие. Начнём заново 👇", "en": "Cancelled the current action. Let's start over 👇",
                      "de": "Aktuelle Aktion abgebrochen. Fangen wir neu an 👇", "ar": "تم إلغاء الإجراء الحالي. لنبدأ من جديد 👇",
                      "zh": "已取消当前操作。重新开始吧 👇", "es": "Acción actual cancelada. Empecemos de nuevo 👇", "fr": "Action en cours annulée. Recommençons 👇"},
    "msg_limit": {"ru": "Лимит генераций закончился.", "en": "Generation limit reached.", "de": "Generierungslimit erreicht.",
                  "ar": "انتهى حد الإنشاء.", "zh": "生成次数已用完。", "es": "Se agotó el límite de generaciones.", "fr": "Limite de générations atteinte."},
    "msg_table_ready": {"ru": "Таблица готова ✅", "en": "Table ready ✅", "de": "Tabelle fertig ✅", "ar": "الجدول جاهز ✅",
                        "zh": "表格已完成 ✅", "es": "Tabla lista ✅", "fr": "Tableau prêt ✅"},
    "msg_doc_ready": {"ru": "Документ готов ✅", "en": "Document ready ✅", "de": "Dokument fertig ✅", "ar": "المستند جاهز ✅",
                      "zh": "文档已完成 ✅", "es": "Documento listo ✅", "fr": "Document prêt ✅"},
    "msg_build_error_table": {"ru": "Не собрал таблицу. Попробуй ещё раз.", "en": "Couldn't build the table. Try again.",
                              "de": "Tabelle konnte nicht erstellt werden. Versuche es erneut.", "ar": "تعذر إنشاء الجدول. حاول مرة أخرى.",
                              "zh": "未能生成表格。请重试。", "es": "No se pudo crear la tabla. Inténtalo de nuevo.", "fr": "Impossible de créer le tableau. Réessayez."},
    "msg_build_error_doc": {"ru": "Не собрал документ. Попробуй ещё раз.", "en": "Couldn't build the document. Try again.",
                            "de": "Dokument konnte nicht erstellt werden. Versuche es erneut.", "ar": "تعذر إنشاء المستند. حاول مرة أخرى.",
                            "zh": "未能生成文档。请重试。", "es": "No se pudo crear el documento. Inténtalo de nuevo.", "fr": "Impossible de créer le document. Réessayez."},
    "msg_for_whom_table": {"ru": "Для кого таблица?", "en": "Who is the table for?", "de": "Für wen ist die Tabelle?",
                           "ar": "لمن هذا الجدول؟", "zh": "表格是给谁用的？", "es": "¿Para quién es la tabla?", "fr": "Pour qui est ce tableau ?"},
    "msg_for_whom_doc": {"ru": "Для кого документ?", "en": "Who is the document for?", "de": "Für wen ist das Dokument?",
                         "ar": "لمن هذا المستند؟", "zh": "文档是给谁用的？", "es": "¿Para quién es el documento?", "fr": "Pour qui est ce document ?"},
    "msg_pick_category": {"ru": "Выбери категорию кнопкой ниже.", "en": "Pick a category using the button below.", "de": "Wähle eine Kategorie mit dem Button unten.",
                          "ar": "اختر فئة من الأزرار أدناه.", "zh": "请用下方按钮选择分类。", "es": "Elige una categoría con el botón de abajo.", "fr": "Choisis une catégorie avec le bouton ci-dessous."},
    "msg_which_table": {"ru": "Какая таблица нужна?", "en": "Which table do you need?", "de": "Welche Tabelle brauchst du?",
                        "ar": "أي جدول تحتاج؟", "zh": "需要哪种表格？", "es": "¿Qué tabla necesitas?", "fr": "Quel tableau te faut-il ?"},
    "msg_which_doc": {"ru": "Какой документ нужен?", "en": "Which document do you need?", "de": "Welches Dokument brauchst du?",
                      "ar": "أي مستند تحتاج؟", "zh": "需要哪种文档？", "es": "¿Qué documento necesitas?", "fr": "Quel document te faut-il ?"},
    "msg_pick_table": {"ru": "Выбери таблицу кнопкой ниже.", "en": "Pick a table using the button below.", "de": "Wähle eine Tabelle mit dem Button unten.",
                       "ar": "اختر جدولاً من الأزرار أدناه.", "zh": "请用下方按钮选择表格。", "es": "Elige una tabla con el botón de abajo.", "fr": "Choisis un tableau avec le bouton ci-dessous."},
    "msg_pick_doc": {"ru": "Выбери документ кнопкой ниже.", "en": "Pick a document using the button below.", "de": "Wähle ein Dokument mit dem Button unten.",
                     "ar": "اختر مستنداً من الأزرار أدناه.", "zh": "请用下方按钮选择文档。", "es": "Elige un documento con el botón de abajo.", "fr": "Choisis un document avec le bouton ci-dessous."},
    "msg_more_docs": {"ru": "Ещё документы:", "en": "More documents:", "de": "Weitere Dokumente:",
                      "ar": "المزيد من المستندات:", "zh": "更多文档：", "es": "Más documentos:", "fr": "Plus de documents :"},
    "msg_build_table_q": {"ru": "Собрать таблицу?", "en": "Build the table?", "de": "Tabelle erstellen?",
                          "ar": "هل ننشئ الجدول؟", "zh": "生成表格吗？", "es": "¿Creo la tabla?", "fr": "On crée le tableau ?"},
    "msg_build_model_q": {"ru": "Собрать финмодель?", "en": "Build the financial model?", "de": "Finanzmodell erstellen?",
                          "ar": "هل ننشئ النموذج المالي؟", "zh": "生成财务模型吗？", "es": "¿Creo el modelo financiero?", "fr": "On crée le modèle financier ?"},
    "msg_resend_data": {"ru": "Пришли исходные данные заново.", "en": "Please send the source data again.", "de": "Sende die Ausgangsdaten erneut.",
                        "ar": "أرسل البيانات الأصلية مرة أخرى.", "zh": "请重新发送原始数据。", "es": "Envía de nuevo los datos originales.", "fr": "Renvoie les données d'origine."},
    "msg_what_add": {"ru": "Что добавить или уточнить?", "en": "What should I add or clarify?", "de": "Was soll ich hinzufügen oder ändern?",
                     "ar": "ما الذي تريد إضافته أو توضيحه؟", "zh": "需要添加或说明什么？", "es": "¿Qué añado o aclaro?", "fr": "Qu'est-ce qu'il faut ajouter ou préciser ?"},
    "msg_choose_action": {"ru": "Выбери действие кнопкой ниже.", "en": "Choose an action using the button below.", "de": "Wähle eine Aktion mit dem Button unten.",
                          "ar": "اختر إجراءً من الأزرار أدناه.", "zh": "请用下方按钮选择操作。", "es": "Elige una acción con el botón de abajo.", "fr": "Choisis une action avec le bouton ci-dessous."},
    "msg_new_text": {"ru": "Пришли новый текст:", "en": "Send the new text:", "de": "Sende den neuen Text:",
                     "ar": "أرسل النص الجديد:", "zh": "请发送新文本：", "es": "Envía el nuevo texto:", "fr": "Envoie le nouveau texte :"},
    "msg_new_topic": {"ru": "Напиши новую тему:", "en": "Write the new topic:", "de": "Schreibe das neue Thema:",
                      "ar": "اكتب الموضوع الجديد:", "zh": "请输入新主题：", "es": "Escribe el nuevo tema:", "fr": "Écris le nouveau sujet :"},
    "msg_extra_limit": {"ru": "Лимит добавлений исчерпан.", "en": "You've reached the edit limit.", "de": "Das Limit für Ergänzungen ist erreicht.",
                        "ar": "لقد استنفدت حد الإضافات.", "zh": "修改次数已用完。", "es": "Se agotó el límite de adiciones.", "fr": "Limite de modifications atteinte."},
    "msg_extra_prompt": {"ru": "Напиши, что добавить или изменить. Максимум 800 символов.", "en": "Write what to add or change. Max 800 characters.",
                         "de": "Schreibe, was hinzugefügt oder geändert werden soll. Max. 800 Zeichen.", "ar": "اكتب ما تريد إضافته أو تغييره. الحد الأقصى ٨٠٠ حرف.",
                         "zh": "请输入要添加或修改的内容，最多800字。", "es": "Escribe qué añadir o cambiar. Máximo 800 caracteres.", "fr": "Écris ce qu'il faut ajouter ou modifier. 800 caractères max."},
    "msg_no_template_study": {"ru": "Для учебных документов шаблон не предусмотрен.", "en": "There's no template for study documents.",
                              "de": "Für Studienarbeiten gibt es keine Vorlage.", "ar": "لا يوجد قالب للمستندات الدراسية.",
                              "zh": "学习类文档暂无模板。", "es": "No hay plantilla para documentos académicos.", "fr": "Il n'y a pas de modèle pour les documents scolaires."},
    "msg_building_template": {"ru": "Собираю шаблон…", "en": "Building the template…", "de": "Vorlage wird erstellt…",
                              "ar": "جاري إنشاء القالب…", "zh": "正在生成模板…", "es": "Creando la plantilla…", "fr": "Création du modèle…"},
    "msg_ready": {"ru": "Готово ✅", "en": "Done ✅", "de": "Fertig ✅", "ar": "تم ✅", "zh": "完成 ✅", "es": "Listo ✅", "fr": "Terminé ✅"},
    "msg_which_size": {"ru": "Какой объём?", "en": "How much detail?", "de": "Welcher Umfang?",
                       "ar": "ما الحجم المطلوب؟", "zh": "需要多少内容？", "es": "¿Qué extensión?", "fr": "Quelle ampleur ?"},
    "msg_building_draft": {"ru": "Собираю черновик…", "en": "Building the draft…", "de": "Entwurf wird erstellt…",
                           "ar": "جاري إعداد المسودة…", "zh": "正在生成草稿…", "es": "Creando el borrador…", "fr": "Création du brouillon…"},
    "msg_no_text": {"ru": "Не собрал текст. Нажми ещё раз «{btn}».", "en": "Couldn't build the text. Tap \"{btn}\" again.",
                    "de": "Text konnte nicht erstellt werden. Tippe erneut auf „{btn}“.", "ar": "تعذر إنشاء النص. اضغط «{btn}» مرة أخرى.",
                    "zh": "未能生成文本，请再次点击「{btn}」。", "es": "No se pudo crear el texto. Vuelve a tocar «{btn}».", "fr": "Impossible de créer le texte. Appuie à nouveau sur « {btn} »."},
    "msg_photo_wait": {"ru": "Фото пока не пришло. Пришли хотя бы одно или вернись к «Все фото через ИИ».",
                       "en": "No photo yet. Send at least one, or go back to \"All photos via AI\".",
                       "de": "Noch kein Foto erhalten. Sende mindestens eins oder gehe zurück zu „Alle Fotos per KI“.",
                       "ar": "لم تصل أي صورة بعد. أرسل صورة واحدة على الأقل أو ارجع إلى «كل الصور بالذكاء الاصطناعي».",
                       "zh": "还没有收到照片。请至少发送一张，或返回「全部用AI生成图片」。",
                       "es": "Aún no llegó ninguna foto. Envía al menos una o vuelve a «Todas las fotos con IA».",
                       "fr": "Aucune photo reçue pour l'instant. Envoie-en au moins une ou reviens à « Toutes les photos via l'IA »."},
    "msg_send_photos": {"ru": "Пришли фото 📷 или нажми «Готово, собери презентацию», когда закончишь.",
                        "en": "Send photos 📷 or tap \"Done, build the presentation\" when finished.",
                        "de": "Sende Fotos 📷 oder tippe auf „Fertig, Präsentation erstellen“, wenn du fertig bist.",
                        "ar": "أرسل الصور 📷 أو اضغط «تم، أنشئ العرض التقديمي» عند الانتهاء.",
                        "zh": "发送照片📷，完成后点击「完成，生成演示文稿」。",
                        "es": "Envía fotos 📷 o toca «Listo, crea la presentación» cuando termines.",
                        "fr": "Envoie des photos 📷 ou appuie sur « Terminé, créez la présentation » quand tu as fini."},
    "msg_topic_prompt": {"ru": "Напиши тему. Можно коротко, например: киты, крипта, школа.",
                         "en": "Write the topic. Can be short, e.g.: whales, crypto, school.",
                         "de": "Schreibe das Thema. Kurz ist ok, z. B.: Wale, Krypto, Schule.",
                         "ar": "اكتب الموضوع. يمكن أن يكون قصيراً، مثل: الحيتان، العملات الرقمية، المدرسة.",
                         "zh": "请输入主题，可以简短一些，例如：鲸鱼、加密货币、学校。",
                         "es": "Escribe el tema. Puede ser breve, p. ej.: ballenas, cripto, escuela.",
                         "fr": "Écris le sujet. Ça peut être court, ex. : baleines, crypto, école."},
    "msg_own_text_prompt": {"ru": "Пришли свой текст 📝\nМожно черновик — я поправлю ошибки и соберу структуру.",
                            "en": "Send your text 📝\nA draft is fine — I'll fix errors and build the structure.",
                            "de": "Sende deinen Text 📝\nEin Entwurf reicht — ich korrigiere Fehler und erstelle die Struktur.",
                            "ar": "أرسل نصك 📝\nيمكن أن تكون مسودة — سأصحح الأخطاء وأبني الهيكل.",
                            "zh": "请发送你的文本📝\n草稿也可以——我会修正错误并构建结构。",
                            "es": "Envía tu texto 📝\nUn borrador está bien — corregiré errores y armaré la estructura.",
                            "fr": "Envoie ton texte 📝\nUn brouillon suffit — je corrigerai les erreurs et je construirai la structure."},
    "msg_which_style": {"ru": "Выбери стиль оформления:", "en": "Choose a design style:", "de": "Wähle einen Gestaltungsstil:",
                        "ar": "اختر نمط التصميم:", "zh": "选择设计风格：", "es": "Elige un estilo de diseño:", "fr": "Choisis un style de design :"},
    "msg_pick_style": {"ru": "Выбери стиль кнопкой ниже.", "en": "Pick a style using the button below.", "de": "Wähle einen Stil mit dem Button unten.",
                       "ar": "اختر نمطاً من الأزرار أدناه.", "zh": "请用下方按钮选择风格。", "es": "Elige un estilo con el botón de abajo.", "fr": "Choisis un style avec le bouton ci-dessous."},
    "msg_how_many_slides": {"ru": "Сколько слайдов сделать?", "en": "How many slides?", "de": "Wie viele Folien?",
                            "ar": "كم عدد الشرائح؟", "zh": "需要多少张幻灯片？", "es": "¿Cuántas diapositivas?", "fr": "Combien de diapositives ?"},
    "msg_building_sample": {"ru": "Секунду, собираю пробный вариант…", "en": "One moment, building a sample…", "de": "Einen Moment, ich erstelle einen Testentwurf…",
                            "ar": "لحظة، جاري إعداد نسخة تجريبية…", "zh": "请稍等，正在生成样例…", "es": "Un momento, creando una versión de prueba…", "fr": "Un instant, je prépare un aperçu…"},
    "msg_choose_action_below": {"ru": "Выбери действие:", "en": "Choose an action:", "de": "Wähle eine Aktion:",
                                "ar": "اختر إجراءً:", "zh": "选择操作：", "es": "Elige una acción:", "fr": "Choisis une action :"},
    "msg_how_build_doc": {"ru": "Как собираем документ?\n\n✨ Сгенерировать с ИИ — я сам напишу текст.\n📝 Вставить свой текст — пришли материал, я поправлю и оформлю.",
                          "en": "How should we build the document?\n\n✨ Generate with AI — I'll write the text myself.\n📝 Paste your own text — send the material, I'll fix and format it.",
                          "de": "Wie erstellen wir das Dokument?\n\n✨ Mit KI generieren — ich schreibe den Text selbst.\n📝 Eigenen Text einfügen — sende das Material, ich korrigiere und formatiere es.",
                          "ar": "كيف ننشئ المستند؟\n\n✨ إنشاء بالذكاء الاصطناعي — سأكتب النص بنفسي.\n📝 إدراج نصك الخاص — أرسل المادة وسأصححها وأنسقها.",
                          "zh": "如何生成文档？\n\n✨ 用AI生成——由我来撰写文本。\n📝 粘贴自己的文本——发送素材，我会修正并排版。",
                          "es": "¿Cómo creamos el documento?\n\n✨ Generar con IA — yo escribiré el texto.\n📝 Pegar tu propio texto — envía el material, lo corregiré y daré formato.",
                          "fr": "Comment créons-nous le document ?\n\n✨ Générer avec l'IA — j'écrirai le texte moi-même.\n📝 Coller votre propre texte — envoyez le contenu, je le corrigerai et le mettrai en forme."},
    "msg_template_hint_suffix": {"ru": "\n📄 Скачать шаблон — пустой бланк с пропусками для заполнения от руки.",
                                 "en": "\n📄 Download template — a blank form with fields to fill in by hand.",
                                 "de": "\n📄 Vorlage herunterladen — ein leeres Formular mit Feldern zum handschriftlichen Ausfüllen.",
                                 "ar": "\n📄 تنزيل القالب — نموذج فارغ بحقول للتعبئة يدوياً.",
                                 "zh": "\n📄 下载模板——带填空的空白表单，可手动填写。",
                                 "es": "\n📄 Descargar plantilla — un formulario en blanco con espacios para rellenar a mano.",
                                 "fr": "\n📄 Télécharger le modèle — un formulaire vierge à remplir à la main."},
    "msg_too_long_500": {"ru": "Слишком длинно. Максимум 500 символов.", "en": "Too long. Max 500 characters.", "de": "Zu lang. Max. 500 Zeichen.",
                         "ar": "طويل جداً. الحد الأقصى ٥٠٠ حرف.", "zh": "太长了，最多500字。", "es": "Demasiado largo. Máximo 500 caracteres.", "fr": "Trop long. 500 caractères max."},
    "msg_too_long_800": {"ru": "Слишком длинно. Максимум 800 символов.", "en": "Too long. Max 800 characters.", "de": "Zu lang. Max. 800 Zeichen.",
                         "ar": "طويل جداً. الحد الأقصى ٨٠٠ حرف.", "zh": "太长了，最多800字。", "es": "Demasiado largo. Máximo 800 caracteres.", "fr": "Trop long. 800 caractères max."},
    "msg_too_long_4000": {"ru": "Слишком длинно. Сократи до 4000 символов.", "en": "Too long. Shorten to 4000 characters.", "de": "Zu lang. Kürze auf 4000 Zeichen.",
                          "ar": "طويل جداً. اختصره إلى ٤٠٠٠ حرف.", "zh": "太长了，请缩短至4000字以内。", "es": "Demasiado largo. Reduce a 4000 caracteres.", "fr": "Trop long. Réduis à 4000 caractères."},
    "msg_too_little_text": {"ru": "Текста мало. Пришли хотя бы несколько абзацев.", "en": "Not enough text. Send at least a few paragraphs.",
                            "de": "Zu wenig Text. Sende mindestens ein paar Absätze.", "ar": "النص قليل جداً. أرسل بضع فقرات على الأقل.",
                            "zh": "文本太少，请至少发送几段。", "es": "Muy poco texto. Envía al menos varios párrafos.", "fr": "Pas assez de texte. Envoie au moins quelques paragraphes."},
    "msg_too_little_text_doc": {"ru": "Текста мало. Пришли чуть больше деталей.", "en": "Not enough text. Send a bit more detail.",
                                "de": "Zu wenig Text. Sende etwas mehr Details.", "ar": "النص قليل جداً. أرسل مزيداً من التفاصيل.",
                                "zh": "文本太少，请补充更多细节。", "es": "Muy poco texto. Envía un poco más de detalle.", "fr": "Pas assez de texte. Envoie un peu plus de détails."},
    "msg_updating_draft": {"ru": "Обновляю черновик…", "en": "Updating the draft…", "de": "Entwurf wird aktualisiert…",
                           "ar": "جاري تحديث المسودة…", "zh": "正在更新草稿…", "es": "Actualizando el borrador…", "fr": "Mise à jour du brouillon…"},
    "msg_enough_photos": {"ru": "Этого достаточно — дальше можно не присылать, нажми «Готово, собери презентацию».",
                          "en": "That's enough — no need to send more, tap \"Done, build the presentation\".",
                          "de": "Das reicht — du musst nicht mehr senden, tippe auf „Fertig, Präsentation erstellen“.",
                          "ar": "هذا يكفي — لا حاجة لإرسال المزيد، اضغط «تم، أنشئ العرض التقديمي».",
                          "zh": "这些已经够了——无需再发送，点击「完成，生成演示文稿」。",
                          "es": "Con esto es suficiente — no hace falta enviar más, toca «Listo, crea la presentación».",
                          "fr": "C'est suffisant — inutile d'en envoyer plus, appuie sur « Terminé, créez la présentation »."},
    "msg_photo_download_fail": {"ru": "Не получилось скачать это фото, пришли ещё раз или пропусти его.",
                                "en": "Couldn't download that photo, send it again or skip it.",
                                "de": "Dieses Foto konnte nicht heruntergeladen werden, sende es erneut oder überspringe es.",
                                "ar": "تعذر تنزيل هذه الصورة، أرسلها مرة أخرى أو تخطّها.",
                                "zh": "无法下载这张照片，请重新发送或跳过。",
                                "es": "No se pudo descargar esa foto, envíala de nuevo u omítela.",
                                "fr": "Impossible de télécharger cette photo, renvoie-la ou passe-la."},
    "msg_building_pres_photos": {"ru": "Собираю презентацию с картинками. Это займёт 1–2 минуты.",
                                 "en": "Building the presentation with images. This will take 1–2 minutes.",
                                 "de": "Präsentation mit Bildern wird erstellt. Das dauert 1–2 Minuten.",
                                 "ar": "جاري إنشاء العرض التقديمي بالصور. سيستغرق ١-٢ دقيقة.",
                                 "zh": "正在生成带图片的演示文稿，这需要1-2分钟。",
                                 "es": "Creando la presentación con imágenes. Esto tardará 1-2 minutos.",
                                 "fr": "Création de la présentation avec images. Cela prendra 1 à 2 minutes."},
    "msg_building_excel": {"ru": "Собираю Excel…", "en": "Building the Excel file…", "de": "Excel-Datei wird erstellt…",
                           "ar": "جاري إنشاء ملف Excel…", "zh": "正在生成Excel文件…", "es": "Creando el archivo Excel…", "fr": "Création du fichier Excel…"},
    "msg_building_word": {"ru": "Собираю Word. Обычно это быстрее презентации.", "en": "Building the Word file. This is usually faster than a presentation.",
                          "de": "Word-Datei wird erstellt. Das geht meist schneller als eine Präsentation.",
                          "ar": "جاري إنشاء ملف Word. عادة ما يكون أسرع من العرض التقديمي.",
                          "zh": "正在生成Word文件，通常比演示文稿更快。",
                          "es": "Creando el archivo Word. Suele ser más rápido que una presentación.",
                          "fr": "Création du fichier Word. C'est généralement plus rapide qu'une présentation."},
    "msg_history_empty": {"ru": "История пустая.", "en": "History is empty.", "de": "Verlauf ist leer.",
                          "ar": "السجل فارغ.", "zh": "历史记录为空。", "es": "El historial está vacío.", "fr": "L'historique est vide."},
    "msg_history_title": {"ru": "📁 История:\n\n", "en": "📁 History:\n\n", "de": "📁 Verlauf:\n\n",
                          "ar": "📁 السجل:\n\n", "zh": "📁 历史记录：\n\n", "es": "📁 Historial:\n\n", "fr": "📁 Historique :\n\n"},
    "msg_start_pres_again": {"ru": "Сделать презентацию", "en": "Create a presentation", "de": "Präsentation erstellen",
                            "ar": "إنشاء عرض تقديمي", "zh": "制作演示文稿", "es": "Crear una presentación", "fr": "Créer une présentation"},
    "msg_how_build_pres": {"ru": "Как делаем презентацию?\n\n✨ Сгенерировать с ИИ — я сам придумаю текст, стиль и фото.\n📝 Вставить свой текст — пришли материал, я поправлю ошибки и соберу слайды.",
                           "en": "How should we build the presentation?\n\n✨ Generate with AI — I'll come up with the text, style and photos myself.\n📝 Paste your own text — send the material, I'll fix errors and build the slides.",
                           "de": "Wie erstellen wir die Präsentation?\n\n✨ Mit KI generieren — ich denke mir Text, Stil und Fotos selbst aus.\n📝 Eigenen Text einfügen — sende das Material, ich korrigiere Fehler und erstelle die Folien.",
                           "ar": "كيف ننشئ العرض التقديمي؟\n\n✨ إنشاء بالذكاء الاصطناعي — سأبتكر النص والنمط والصور بنفسي.\n📝 إدراج نصك الخاص — أرسل المادة وسأصحح الأخطاء وأبني الشرائح.",
                           "zh": "如何制作演示文稿？\n\n✨ 用AI生成——由我构思文本、风格和图片。\n📝 粘贴自己的文本——发送素材，我会修正错误并制作幻灯片。",
                           "es": "¿Cómo creamos la presentación?\n\n✨ Generar con IA — yo mismo pensaré el texto, el estilo y las fotos.\n📝 Pegar tu propio texto — envía el material, corregiré errores y armaré las diapositivas.",
                           "fr": "Comment créons-nous la présentation ?\n\n✨ Générer avec l'IA — j'imaginerai moi-même le texte, le style et les photos.\n📝 Coller votre propre texte — envoyez le contenu, je corrigerai les erreurs et je construirai les diapositives."},
    "msg_style_fits_topic": {"ru": "🎨 По теме подходит стиль: {style}.\nОставить его или выбрать другой?",
                             "en": "🎨 This style fits the topic: {style}.\nKeep it or choose another one?",
                             "de": "🎨 Zum Thema passt der Stil: {style}.\nBehalten oder einen anderen wählen?",
                             "ar": "🎨 النمط المناسب للموضوع: {style}.\nهل تحتفظ به أم تختار غيره؟",
                             "zh": "🎨 适合该主题的风格：{style}。\n保留还是选择其他？",
                             "es": "🎨 Para este tema encaja el estilo: {style}.\n¿Lo mantenemos o eliges otro?",
                             "fr": "🎨 Ce style convient au sujet : {style}.\nOn le garde ou tu en choisis un autre ?"},
    "msg_style_fits_text": {"ru": "🎨 По тексту подходит стиль: {style}.\nОставить его или выбрать другой?",
                            "en": "🎨 This style fits the text: {style}.\nKeep it or choose another one?",
                            "de": "🎨 Zum Text passt der Stil: {style}.\nBehalten oder einen anderen wählen?",
                            "ar": "🎨 النمط المناسب للنص: {style}.\nهل تحتفظ به أم تختار غيره؟",
                            "zh": "🎨 适合该文本的风格：{style}。\n保留还是选择其他？",
                            "es": "🎨 Para este texto encaja el estilo: {style}.\n¿Lo mantenemos o eliges otro?",
                            "fr": "🎨 Ce style convient au texte : {style}.\nOn le garde ou tu en choisis un autre ?"},
    "msg_grok_error": {"ru": "Не получилось получить ответ от нейросети. Попробуй ещё раз через минуту.",
                       "en": "Couldn't get a response from the AI. Try again in a minute.",
                       "de": "Keine Antwort von der KI erhalten. Versuche es in einer Minute erneut.",
                       "ar": "تعذر الحصول على رد من الذكاء الاصطناعي. حاول مرة أخرى بعد دقيقة.",
                       "zh": "未能获得AI的回复，请一分钟后再试。",
                       "es": "No se pudo obtener respuesta de la IA. Inténtalo de nuevo en un minuto.",
                       "fr": "Impossible d'obtenir une réponse de l'IA. Réessaie dans une minute."},
    "msg_draft_ready_pres": {"ru": "Черновик готов ✅\n\n{sample}\n\nСтиль: {style}\n\nЕсли всё ок — собираем полную версию.",
                             "en": "Draft ready ✅\n\n{sample}\n\nStyle: {style}\n\nIf it looks good, let's build the full version.",
                             "de": "Entwurf fertig ✅\n\n{sample}\n\nStil: {style}\n\nWenn alles passt, erstellen wir die vollständige Version.",
                             "ar": "المسودة جاهزة ✅\n\n{sample}\n\nالنمط: {style}\n\nإذا كان كل شيء جيداً، دعنا ننشئ النسخة الكاملة.",
                             "zh": "草稿已完成 ✅\n\n{sample}\n\n风格：{style}\n\n如果没问题，我们来生成完整版本。",
                             "es": "Borrador listo ✅\n\n{sample}\n\nEstilo: {style}\n\nSi todo está bien, creemos la versión completa.",
                             "fr": "Brouillon prêt ✅\n\n{sample}\n\nStyle : {style}\n\nSi tout va bien, on crée la version complète."},
    "msg_style_changed": {"ru": "Стиль изменён: {style}\n\nЧерновик:\n\n{sample}\n\nВыбери действие:",
                          "en": "Style changed: {style}\n\nDraft:\n\n{sample}\n\nChoose an action:",
                          "de": "Stil geändert: {style}\n\nEntwurf:\n\n{sample}\n\nWähle eine Aktion:",
                          "ar": "تم تغيير النمط: {style}\n\nالمسودة:\n\n{sample}\n\nاختر إجراءً:",
                          "zh": "风格已更改：{style}\n\n草稿：\n\n{sample}\n\n选择操作：",
                          "es": "Estilo cambiado: {style}\n\nBorrador:\n\n{sample}\n\nElige una acción:",
                          "fr": "Style modifié : {style}\n\nBrouillon :\n\n{sample}\n\nChoisis une action :"},
    "msg_draft_updated": {"ru": "Обновлённый черновик ✅\n\n{sample}\n\nВыбери действие:",
                          "en": "Updated draft ✅\n\n{sample}\n\nChoose an action:",
                          "de": "Aktualisierter Entwurf ✅\n\n{sample}\n\nWähle eine Aktion:",
                          "ar": "المسودة المحدّثة ✅\n\n{sample}\n\nاختر إجراءً:",
                          "zh": "更新后的草稿 ✅\n\n{sample}\n\n选择操作：",
                          "es": "Borrador actualizado ✅\n\n{sample}\n\nElige una acción:",
                          "fr": "Brouillon mis à jour ✅\n\n{sample}\n\nChoisis une action :"},
    "msg_draft_update_failed": {"ru": "Не получилось обновить черновик — нейросеть не ответила. Прошлый черновик остался как был.",
                                "en": "Couldn't update the draft — the AI didn't respond. The previous draft stayed as it was.",
                                "de": "Entwurf konnte nicht aktualisiert werden — die KI hat nicht geantwortet. Der vorherige Entwurf blieb unverändert.",
                                "ar": "تعذر تحديث المسودة — لم يرد الذكاء الاصطناعي. بقيت المسودة السابقة كما كانت.",
                                "zh": "未能更新草稿——AI没有响应。之前的草稿保持不变。",
                                "es": "No se pudo actualizar el borrador — la IA no respondió. El borrador anterior se mantuvo igual.",
                                "fr": "Impossible de mettre à jour le brouillon — l'IA n'a pas répondu. Le brouillon précédent est resté tel quel."},
    "msg_photo_or_own": {"ru": "Фото для слайдов сделать через ИИ или использовать свои?\nЕсли фото будет меньше, чем слайдов — оставшиеся бот дорисует сам.",
                         "en": "Should slide photos be AI-generated or your own?\nIf you send fewer photos than slides, the bot will generate the rest.",
                         "de": "Fotos für die Folien per KI erstellen oder eigene verwenden?\nWenn du weniger Fotos als Folien sendest, generiert der Bot den Rest selbst.",
                         "ar": "هل تُنشأ صور الشرائح بالذكاء الاصطناعي أم تستخدم صورك الخاصة؟\nإذا كانت الصور أقل من الشرائح، سيكمل البوت الباقي تلقائياً.",
                         "zh": "幻灯片图片用AI生成还是使用你自己的照片？\n如果照片数量少于幻灯片数，机器人会自动补全其余部分。",
                         "es": "¿Las fotos de las diapositivas se generan con IA o usas las tuyas?\nSi envías menos fotos que diapositivas, el bot generará el resto.",
                         "fr": "Les photos des diapositives doivent-elles être générées par l'IA ou tes propres photos ?\nSi tu envoies moins de photos que de diapositives, le bot complétera le reste."},
    "msg_send_photos_one_by_one": {
        "ru": "У тебя {slides} слайдов — самое то прислать до {slides} фото, по одному на слайд (недостающие я дорисую сам через ИИ). Можно пачкой (альбомом, Telegram отправляет до 10 за раз) или по одному. Когда закончишь — нажми «Готово, собери презентацию».",
        "en": "You picked {slides} slides — feel free to send up to {slides} photos, one per slide (I'll generate the rest with AI if you send fewer). All at once as an album (Telegram allows up to 10 per batch) or one by one, whatever's easier. When done, tap \"Done, build the presentation\".",
        "de": "Du hast {slides} Folien gewählt — sende am besten bis zu {slides} Fotos, eins pro Folie (den Rest erstelle ich per KI, falls du weniger sendest). Alle auf einmal als Album (Telegram erlaubt bis zu 10 pro Stapel) oder einzeln, wie du magst. Wenn fertig, tippe auf „Fertig, Präsentation erstellen“.",
        "ar": "اخترت {slides} شريحة — يمكنك إرسال حتى {slides} صورة، صورة لكل شريحة (سأكمل الباقي بالذكاء الاصطناعي إذا أرسلت أقل). دفعة واحدة كألبوم (يسمح Telegram بحتى 10 دفعة واحدة) أو واحدة تلو الأخرى، كما يناسبك. عند الانتهاء، اضغط «تم، أنشئ العرض التقديمي».",
        "zh": "你选择了{slides}张幻灯片——建议发送最多{slides}张照片，每张幻灯片配一张（如果照片不够，我会用AI补全）。可以整批作为相册一次发送（Telegram单次最多10张），也可以逐张发送，随你方便。完成后点击「完成，生成演示文稿」。",
        "es": "Elegiste {slides} diapositivas — lo ideal es enviar hasta {slides} fotos, una por diapositiva (completaré el resto con IA si envías menos). Todas juntas como álbum (Telegram permite hasta 10 por lote) o una por una, como prefieras. Cuando termines, toca «Listo, crea la presentación».",
        "fr": "Tu as choisi {slides} diapositives — envoie idéalement jusqu'à {slides} photos, une par diapositive (je génère le reste avec l'IA si tu en envoies moins). Toutes en une fois comme un album (Telegram autorise jusqu'à 10 par envoi) ou une par une, comme tu préfères. Une fois terminé, appuie sur « Terminé, créez la présentation »."},
    "msg_photo_received": {"ru": "Фото {n} получено 📷 Пришли ещё или нажми «Готово, собери презентацию».",
                           "en": "Photo {n} received 📷 Send more or tap \"Done, build the presentation\".",
                           "de": "Foto {n} erhalten 📷 Sende weitere oder tippe auf „Fertig, Präsentation erstellen“.",
                           "ar": "تم استلام الصورة {n} 📷 أرسل المزيد أو اضغط «تم، أنشئ العرض التقديمي».",
                           "zh": "已收到第{n}张照片📷，可继续发送，或点击「完成，生成演示文稿」。",
                           "es": "Foto {n} recibida 📷 Envía más o toca «Listo, crea la presentación».",
                           "fr": "Photo {n} reçue 📷 Envoie-en d'autres ou appuie sur « Terminé, créez la présentation »."},
    "msg_startup_data_prompt": {
        "ru": "Пришли исходные данные для расчёта одним сообщением:\n\n• стартовые вложения\n• ежемесячные расходы списком (аренда, зарплаты, реклама и т.д. — с суммами)\n• цена за единицу товара/услуги\n• ожидаемое количество продаж в месяц\n• на сколько месяцев считать\n\nНапример: вложения 500 000 ₽, аренда 40 000, зарплаты 90 000, реклама 25 000, цена 350 ₽, продажи 600 шт/мес, считать на 12 месяцев.",
        "en": "Send the source data for the calculation in one message:\n\n• starting investment\n• monthly expenses as a list (rent, salaries, ads, etc. — with amounts)\n• price per unit of product/service\n• expected number of sales per month\n• how many months to calculate for\n\nExample: investment $5,000, rent $400, salaries $900, ads $250, price $3.5, sales 600 units/month, calculate for 12 months.",
        "de": "Sende die Ausgangsdaten für die Berechnung in einer Nachricht:\n\n• Startinvestition\n• monatliche Ausgaben als Liste (Miete, Gehälter, Werbung usw. — mit Beträgen)\n• Preis pro Einheit Produkt/Dienstleistung\n• erwartete Verkaufszahl pro Monat\n• für wie viele Monate gerechnet werden soll\n\nBeispiel: Investition 5.000 €, Miete 400 €, Gehälter 900 €, Werbung 250 €, Preis 3,5 €, Verkäufe 600 Stk/Monat, Berechnung für 12 Monate.",
        "ar": "أرسل البيانات الأصلية للحساب في رسالة واحدة:\n\n• الاستثمار الأولي\n• المصاريف الشهرية كقائمة (الإيجار، الرواتب، الإعلانات، إلخ — مع المبالغ)\n• سعر الوحدة من المنتج/الخدمة\n• عدد المبيعات المتوقع شهرياً\n• عدد الأشهر للحساب\n\nمثال: استثمار 5000$، إيجار 400$، رواتب 900$، إعلانات 250$، سعر 3.5$، مبيعات 600 وحدة/شهر، الحساب لمدة 12 شهراً.",
        "zh": "请一次性发送以下计算所需数据：\n\n• 初始投资\n• 月度支出清单（租金、工资、广告等——附金额）\n• 每单位产品/服务价格\n• 预期月销量\n• 计算周期（月数）\n\n示例：投资5000美元，租金400，工资900，广告250，单价3.5，月销量600件，按12个月计算。",
        "es": "Envía los datos para el cálculo en un solo mensaje:\n\n• inversión inicial\n• gastos mensuales en lista (alquiler, salarios, publicidad, etc. — con montos)\n• precio por unidad de producto/servicio\n• ventas esperadas por mes\n• para cuántos meses calcular\n\nEjemplo: inversión $5000, alquiler $400, salarios $900, publicidad $250, precio $3.5, ventas 600 uds/mes, calcular para 12 meses.",
        "fr": "Envoie les données source pour le calcul en un seul message :\n\n• investissement initial\n• dépenses mensuelles sous forme de liste (loyer, salaires, publicité, etc. — avec montants)\n• prix par unité de produit/service\n• nombre de ventes attendu par mois\n• sur combien de mois calculer\n\nExemple : investissement 5000 €, loyer 400 €, salaires 900 €, publicité 250 €, prix 3,5 €, ventes 600 unités/mois, calculer sur 12 mois."},
    "msg_how_build_table": {
        "ru": "Как собираем таблицу?\n\n✨ Сгенерировать с ИИ — сам придумаю правдоподобные данные по теме.\n📝 Ввести свои данные — пришли реальные цифры, я оформлю их в таблицу с формулами.",
        "en": "How should we build the table?\n\n✨ Generate with AI — I'll come up with plausible data on the topic.\n📝 Enter your own data — send real numbers, I'll format them into a table with formulas.",
        "de": "Wie erstellen wir die Tabelle?\n\n✨ Mit KI generieren — ich denke mir plausible Daten zum Thema aus.\n📝 Eigene Daten eingeben — sende reale Zahlen, ich formatiere sie in eine Tabelle mit Formeln.",
        "ar": "كيف ننشئ الجدول؟\n\n✨ إنشاء بالذكاء الاصطناعي — سأبتكر بيانات معقولة حول الموضوع.\n📝 إدخال بياناتك الخاصة — أرسل أرقاماً حقيقية وسأنسقها في جدول بمعادلات.",
        "zh": "如何生成表格？\n\n✨ 用AI生成——由我构思与主题相符的合理数据。\n📝 输入自己的数据——发送真实数字，我会将其整理成带公式的表格。",
        "es": "¿Cómo creamos la tabla?\n\n✨ Generar con IA — inventaré datos plausibles sobre el tema.\n📝 Introducir tus propios datos — envía cifras reales, las daré formato en una tabla con fórmulas.",
        "fr": "Comment créons-nous le tableau ?\n\n✨ Générer avec l'IA — j'imaginerai des données plausibles sur le sujet.\n📝 Saisir tes propres données — envoie des chiffres réels, je les mettrai en forme dans un tableau avec des formules."},
    "msg_table_ready": {"ru": "Таблица готова ✅", "en": "Table ready ✅", "de": "Tabelle fertig ✅",
                        "ar": "الجدول جاهز ✅", "zh": "表格已完成 ✅", "es": "Tabla lista ✅", "fr": "Tableau prêt ✅"},
    "msg_check_source_numbers": {"ru": "\n\n⚠️ Проверь исходные цифры — модель только оформила то, что ты прислал, ничего не добавляла от себя.",
                                 "en": "\n\n⚠️ Check the source numbers — the model only formatted what you sent, it didn't add anything on its own.",
                                 "de": "\n\n⚠️ Überprüfe die Ausgangszahlen — das Modell hat nur formatiert, was du gesendet hast, ohne etwas hinzuzufügen.",
                                 "ar": "\n\n⚠️ تحقق من الأرقام الأصلية — النموذج قام فقط بتنسيق ما أرسلته دون إضافة أي شيء من عنده.",
                                 "zh": "\n\n⚠️ 请核对原始数字——模型只是整理了你提供的内容，没有自行添加任何数据。",
                                 "es": "\n\n⚠️ Revisa las cifras originales — el modelo solo dio formato a lo que enviaste, sin añadir nada por su cuenta.",
                                 "fr": "\n\n⚠️ Vérifie les chiffres d'origine — le modèle s'est contenté de mettre en forme ce que tu as envoyé, sans rien ajouter de lui-même."},
    "msg_excel_caption": {"ru": "📈 Excel — с формулами, можно редактировать", "en": "📈 Excel — with formulas, editable",
                          "de": "📈 Excel — mit Formeln, bearbeitbar", "ar": "📈 Excel — بمعادلات، قابل للتعديل",
                          "zh": "📈 Excel——带公式，可编辑", "es": "📈 Excel — con fórmulas, editable", "fr": "📈 Excel — avec formules, modifiable"},
    "msg_template_caption": {"ru": "📄 Пустой шаблон — заполни пропуски от руки или в Word и распечатай",
                             "en": "📄 Blank template — fill in the gaps by hand or in Word and print",
                             "de": "📄 Leere Vorlage — fülle die Lücken handschriftlich oder in Word aus und drucke sie",
                             "ar": "📄 قالب فارغ — املأ الفراغات يدوياً أو في Word ثم اطبعه",
                             "zh": "📄 空白模板——手动或在Word中填写空白处后打印",
                             "es": "📄 Plantilla en blanco — rellena los espacios a mano o en Word e imprime",
                             "fr": "📄 Modèle vierge — remplis les champs à la main ou dans Word puis imprime"},
    "msg_draft_ready_doc": {"ru": "Черновик готов ✅\n\n{sample}\n\nЕсли всё ок — собираем файл.",
                            "en": "Draft ready ✅\n\n{sample}\n\nIf it looks good, let's build the file.",
                            "de": "Entwurf fertig ✅\n\n{sample}\n\nWenn alles passt, erstellen wir die Datei.",
                            "ar": "المسودة جاهزة ✅\n\n{sample}\n\nإذا كان كل شيء جيداً، دعنا ننشئ الملف.",
                            "zh": "草稿已完成 ✅\n\n{sample}\n\n如果没问题，我们来生成文件。",
                            "es": "Borrador listo ✅\n\n{sample}\n\nSi todo está bien, creemos el archivo.",
                            "fr": "Brouillon prêt ✅\n\n{sample}\n\nSi tout va bien, on crée le fichier."},
    "msg_writing_more": {"ru": "Черновик получился короче, чем нужно — дописываю подробнее…",
                         "en": "The draft came out shorter than needed — expanding it…",
                         "de": "Der Entwurf ist kürzer als nötig ausgefallen — ich erweitere ihn…",
                         "ar": "المسودة أقصر من اللازم — جاري إطالتها…",
                         "zh": "草稿比预期短——正在补充完善…",
                         "es": "El borrador quedó más corto de lo necesario — ampliándolo…",
                         "fr": "Le brouillon est plus court que nécessaire — je le développe…"},
    "msg_writing_a_bit_more": {"ru": "Ещё чуть-чуть осталось — дописываю…",
                               "en": "Almost there — still expanding…",
                               "de": "Fast fertig — ich schreibe noch etwas…",
                               "ar": "اقتربنا — ما زلت أكتب…",
                               "zh": "快好了——继续补充中…",
                               "es": "Ya casi está — sigo ampliando…",
                               "fr": "Presque terminé — je continue à développer…"},
    "msg_word_caption": {"ru": "📄 Word — этот файл можно править", "en": "📄 Word — this file can be edited",
                         "de": "📄 Word — diese Datei kann bearbeitet werden", "ar": "📄 Word — يمكن تعديل هذا الملف",
                         "zh": "📄 Word——此文件可编辑", "es": "📄 Word — este archivo se puede editar", "fr": "📄 Word — ce fichier est modifiable"},
    "msg_pdf_caption": {"ru": "📄 PDF-копия", "en": "📄 PDF copy", "de": "📄 PDF-Kopie",
                        "ar": "📄 نسخة PDF", "zh": "📄 PDF副本", "es": "📄 Copia en PDF", "fr": "📄 Copie PDF"},
    "msg_doc_ready": {"ru": "Документ готов ✅", "en": "Document ready ✅", "de": "Dokument fertig ✅",
                      "ar": "المستند جاهز ✅", "zh": "文档已完成 ✅", "es": "Documento listo ✅", "fr": "Document prêt ✅"},
    "msg_check_gaps": {"ru": "\n\n⚠️ Проверь пропуски в тексте перед распечаткой.",
                       "en": "\n\n⚠️ Check the blanks in the text before printing.",
                       "de": "\n\n⚠️ Überprüfe die Lücken im Text vor dem Drucken.",
                       "ar": "\n\n⚠️ تحقق من الفراغات في النص قبل الطباعة.",
                       "zh": "\n\n⚠️ 打印前请检查文中的空白处。",
                       "es": "\n\n⚠️ Revisa los espacios en blanco del texto antes de imprimir.",
                       "fr": "\n\n⚠️ Vérifie les champs vides dans le texte avant d'imprimer."},
    "msg_didnt_understand": {"ru": "Не понял. Выбери действие кнопкой ниже, или напиши /cancel, чтобы начать заново.",
                             "en": "Didn't catch that. Choose an action with the button below, or type /cancel to start over.",
                             "de": "Das habe ich nicht verstanden. Wähle eine Aktion mit dem Button unten oder schreibe /cancel, um neu zu beginnen.",
                             "ar": "لم أفهم. اختر إجراءً من الزر أدناه، أو اكتب /cancel للبدء من جديد.",
                             "zh": "没能理解。请用下方按钮选择操作，或输入 /cancel 重新开始。",
                             "es": "No entendí. Elige una acción con el botón de abajo, o escribe /cancel para empezar de nuevo.",
                             "fr": "Je n'ai pas compris. Choisis une action avec le bouton ci-dessous, ou tape /cancel pour recommencer."},
    "msg_plan_info": {"ru": "ℹ️ Генераций: {used} из {limit}\nОсталось: {left}",
                      "en": "ℹ️ Generations: {used} of {limit}\nRemaining: {left}",
                      "de": "ℹ️ Generierungen: {used} von {limit}\nVerbleibend: {left}",
                      "ar": "ℹ️ التوليدات: {used} من {limit}\nالمتبقي: {left}",
                      "zh": "ℹ️ 生成次数：{used} / {limit}\n剩余：{left}",
                      "es": "ℹ️ Generaciones: {used} de {limit}\nRestantes: {left}",
                      "fr": "ℹ️ Générations : {used} sur {limit}\nRestant : {left}"},
    "msg_pptx_ready": {
        "ru": "Готово ✅\n\nОткрывай именно PPTX в PowerPoint, Keynote или Google Презентациях.\n\nЕсли на телефоне все фото одинаковые, это не ошибка файла. Так бывает в предпросмотре Telegram, WPS и встроенных «Документах». Открой тот же файл на другом устройстве или в нормальном редакторе презентаций.",
        "en": "Done ✅\n\nOpen the PPTX file specifically in PowerPoint, Keynote, or Google Slides.\n\nIf all the photos look the same on your phone, that's not a file error. This happens in Telegram's preview, WPS, and built-in \"Files\" apps. Open the same file on another device or in a proper presentation editor.",
        "de": "Fertig ✅\n\nÖffne die PPTX-Datei speziell in PowerPoint, Keynote oder Google Präsentationen.\n\nWenn auf dem Handy alle Fotos gleich aussehen, ist das kein Dateifehler. Das passiert bei der Telegram-Vorschau, WPS und integrierten „Dateien“-Apps. Öffne dieselbe Datei auf einem anderen Gerät oder in einem richtigen Präsentationseditor.",
        "ar": "تم ✅\n\nافتح ملف PPTX تحديداً في PowerPoint أو Keynote أو Google Slides.\n\nإذا بدت كل الصور متشابهة على الهاتف، فهذا ليس خطأً في الملف. يحدث هذا في معاينة Telegram وWPS وتطبيقات «الملفات» المدمجة. افتح نفس الملف على جهاز آخر أو في محرر عروض تقديمية حقيقي.",
        "zh": "完成 ✅\n\n请务必在 PowerPoint、Keynote 或 Google 幻灯片中打开该 PPTX 文件。\n\n如果手机上所有照片看起来都一样，这不是文件错误。这在 Telegram 预览、WPS 和内置「文件」应用中很常见。请在其他设备或正规演示文稿编辑器中打开同一文件。",
        "es": "Listo ✅\n\nAbre el archivo PPTX específicamente en PowerPoint, Keynote o Presentaciones de Google.\n\nSi en el teléfono todas las fotos se ven iguales, no es un error del archivo. Esto pasa en la vista previa de Telegram, WPS y las apps integradas de «Archivos». Abre el mismo archivo en otro dispositivo o en un editor de presentaciones real.",
        "fr": "Terminé ✅\n\nOuvre bien le fichier PPTX dans PowerPoint, Keynote ou Google Slides.\n\nSi toutes les photos se ressemblent sur ton téléphone, ce n'est pas une erreur du fichier. Cela arrive dans l'aperçu de Telegram, WPS et les applications « Fichiers » intégrées. Ouvre le même fichier sur un autre appareil ou dans un vrai éditeur de présentations."},
    "msg_pptx_caption": {"ru": "📊 PPTX — открывай этот файл", "en": "📊 PPTX — open this file", "de": "📊 PPTX — öffne diese Datei",
                         "ar": "📊 PPTX — افتح هذا الملف", "zh": "📊 PPTX——请打开此文件", "es": "📊 PPTX — abre este archivo", "fr": "📊 PPTX — ouvre ce fichier"},
    "msg_pptx_pdf_caption": {"ru": "📄 PDF — полная текстовая копия (без фото)", "en": "📄 PDF — full text copy (without photos)",
                             "de": "📄 PDF — vollständige Textkopie (ohne Fotos)", "ar": "📄 PDF — نسخة نصية كاملة (بدون صور)",
                             "zh": "📄 PDF——完整文字副本（不含图片）", "es": "📄 PDF — copia completa en texto (sin fotos)",
                             "fr": "📄 PDF — copie texte complète (sans photos)"},
    "msg_content_lang_prompt": {
        "ru": "На каком языке сделать содержимое? Это не обязательно тот же язык, что и у интерфейса.",
        "en": "What language should the content be in? It doesn't have to match the interface language.",
        "de": "In welcher Sprache soll der Inhalt sein? Das muss nicht die Sprache der Oberfläche sein.",
        "ar": "بأي لغة تريد المحتوى؟ لا يلزم أن تكون نفس لغة الواجهة.",
        "zh": "内容需要用什么语言？不一定要和界面语言相同。",
        "es": "¿En qué idioma debe estar el contenido? No tiene por qué coincidir con el idioma de la interfaz.",
        "fr": "Dans quelle langue doit être le contenu ? Ce n'est pas obligatoirement la langue de l'interface."},
    "msg_pick_content_lang": {"ru": "Выбери язык кнопкой ниже.", "en": "Pick a language using the button below.", "de": "Wähle eine Sprache mit dem Button unten.",
                              "ar": "اختر لغة من الأزرار أدناه.", "zh": "请用下方按钮选择语言。", "es": "Elige un idioma con el botón de abajo.", "fr": "Choisis une langue avec le bouton ci-dessous."},
    "msg_signature": {"ru": "Рад был помочь! Ваш, {bot}", "en": "Glad I could help! Yours, {bot}", "de": "Gern geholfen! Dein {bot}",
                      "ar": "سعيد بمساعدتك! صديقك، {bot}", "zh": "很高兴能帮到你！你的 {bot}", "es": "¡Encantado de ayudar! Tuyo, {bot}", "fr": "Ravi d'avoir pu t'aider ! Ton {bot}"},
}


def tr(key, lang, **kw):
    """Перевод интерфейсной строки по ключу. Если для языка нет перевода - откатывается
    на русский, чтобы бот никогда не показал пустоту/ошибку из-за пропуска в словаре."""
    d = TR.get(key, {})
    s = d.get(lang) or d.get("ru") or key
    return s.format(**kw) if kw else s


# Множества всех переводов конкретных служебных кнопок - чтобы F.text.in_(...) фильтры
# в хендлерах узнавали нажатие кнопки независимо от текущего языка пользователя
# (кнопка "🏠 Главное меню" на любом из 7 языков должна сработать одинаково).
ALL_MAIN_MENU_LABELS = set(TR["btn_main_menu"].values())
ALL_MORE_DOCS_LABELS = set(TR["btn_more_docs"].values())
ALL_BACK_CATEGORIES_LABELS = set(TR["btn_back_categories"].values())
ALL_BTN_PRES_LABELS = set(TR["btn_pres"].values())
ALL_BTN_WORD_LABELS = set(TR["btn_word"].values())
ALL_BTN_EXCEL_LABELS = set(TR["btn_excel"].values())
ALL_BTN_HISTORY_LABELS = set(TR["btn_history"].values())
ALL_BTN_PLAN_LABELS = set(TR["btn_plan"].values())
ALL_BTN_AI_GENERATE_LABELS = set(TR["btn_ai_generate"].values())
ALL_BTN_OWN_TEXT_LABELS = set(TR["btn_own_text"].values())
ALL_BTN_OWN_DATA_LABELS = set(TR["btn_own_data"].values())
ALL_BTN_CHANGE_TOPIC_LABELS = set(TR["btn_change_topic"].values())
ALL_BTN_CHANGE_STYLE_LABELS = set(TR["btn_change_style"].values())
ALL_BTN_ADD_INFO_LABELS = set(TR["btn_add_info"].values())
ALL_BTN_FULL_VERSION_LABELS = set(TR["btn_full_version"].values()) | {"делай", "да", "ок", "yes", "ok"}
ALL_BTN_AI_PHOTOS_LABELS = set(TR["btn_ai_photos"].values())
ALL_BTN_OWN_PHOTOS_LABELS = set(TR["btn_own_photos"].values())
ALL_BTN_PHOTOS_DONE_LABELS = set(TR["btn_photos_done"].values())
ALL_BTN_CHANGE_QUERY_LABELS = set(TR["btn_change_query"].values())
ALL_BTN_BUILD_TABLE_LABELS = set(TR["btn_build_table"].values()) | {"делай", "да", "ок", "yes", "ok"}
ALL_BTN_BUILD_DOC_LABELS = set(TR["btn_build_doc"].values()) | {"делай", "да", "ок", "yes", "ok"}
ALL_BTN_TEMPLATE_LABELS = set(TR["btn_template"].values())


def lang_kb():
    rows = []
    codes = list(LANGS.keys())
    for i in range(0, len(codes), 2):
        pair = codes[i:i + 2]
        rows.append([KeyboardButton(text=f"{LANGS[c]['flag']} {LANGS[c]['name']}") for c in pair])
    return ReplyKeyboardMarkup(keyboard=rows, resize_keyboard=True)


LANG_LABEL_TO_CODE = {f"{v['flag']} {v['name']}": k for k, v in LANGS.items()}


def content_lang_kb(interface_lang):
    """Клавиатура выбора языка КОНТЕНТА документа/презентации - отдельно от языка интерфейса,
    т.к. это разные вещи (например, интерфейс на английском, а презентация нужна на русском).
    Первая кнопка - быстрый вариант "как в интерфейсе", дальше все 7 языков."""
    rows = [[KeyboardButton(text=tr("btn_same_as_interface", interface_lang, iface_lang=LANGS[interface_lang]["name"]))]]
    codes = list(LANGS.keys())
    for i in range(0, len(codes), 2):
        pair = codes[i:i + 2]
        rows.append([KeyboardButton(text=f"{LANGS[c]['flag']} {LANGS[c]['name']}") for c in pair])
    rows.append([KeyboardButton(text=tr("btn_main_menu", interface_lang))])
    return ReplyKeyboardMarkup(keyboard=rows, resize_keyboard=True)


def resolve_content_lang(text: str, interface_lang: str):
    """Разбирает нажатие кнопки на клавиатуре content_lang_kb. Возвращает код языка
    или None, если текст не распознан ни как "как в интерфейсе", ни как конкретный язык."""
    if text == tr("btn_same_as_interface", interface_lang, iface_lang=LANGS[interface_lang]["name"]):
        return interface_lang
    return LANG_LABEL_TO_CODE.get(text)


THEMES = {
    # heading_font - шрифт ЗАГОЛОВКОВ слайдов, свой под характер темы (тело текста везде
    # остаётся Calibri - оно должно быть одинаково читаемым во всех 7 языках интерфейса,
    # включая арабский и китайский, где декоративные латинские шрифты всё равно не
    # применяются рендерером). Для строгих/деловых тем (business, minimal) и универсальной
    # (default) шрифт заголовка намеренно не меняем - там уместна одна типографика без игры.
    "nature": {"bg": (248, 246, 241), "ink": (22, 22, 24), "mid": (70, 70, 74), "mute": (130, 128, 124), "line": (22, 22, 24),
               "photo": "photorealistic nature photography, cinematic sunlight, no text, no watermark", "heading_font": "Georgia"},
    "business": {"bg": (16, 16, 18), "ink": (245, 245, 247), "mid": (196, 196, 200), "mute": (120, 120, 126), "line": (212, 175, 90),
                 "photo": "premium business photography, architecture, cinematic, no text, no watermark", "heading_font": "Calibri"},
    "tech": {"bg": (8, 16, 28), "ink": (240, 246, 255), "mid": (176, 196, 220), "mute": (110, 130, 155), "line": (90, 170, 230),
             "photo": "futuristic technology photography, cinematic, no text, no watermark", "heading_font": "Trebuchet MS"},
    "school": {"bg": (250, 249, 246), "ink": (28, 32, 40), "mid": (60, 64, 72), "mute": (120, 124, 132), "line": (40, 90, 180),
               "photo": "clear educational photo, bright, no text, no watermark", "heading_font": "Trebuchet MS"},
    "fashion": {"bg": (252, 250, 247), "ink": (18, 18, 18), "mid": (70, 66, 62), "mute": (140, 134, 128), "line": (18, 18, 18),
                "photo": "editorial fashion photography, magazine look, no text, no watermark", "heading_font": "Georgia"},
    "history": {"bg": (245, 237, 224), "ink": (48, 32, 20), "mid": (90, 70, 50), "mute": (140, 120, 95), "line": (140, 90, 40),
                "photo": "historical documentary photography, museums, archives, cinematic, no text", "heading_font": "Georgia"},
    "science": {"bg": (244, 248, 252), "ink": (16, 32, 56), "mid": (50, 70, 95), "mute": (110, 130, 150), "line": (20, 90, 160),
                "photo": "scientific photography, labs, space, macro details, cinematic, no text", "heading_font": "Trebuchet MS"},
    "sport": {"bg": (18, 18, 20), "ink": (250, 250, 252), "mid": (210, 210, 214), "mute": (140, 140, 146), "line": (230, 70, 40),
              "photo": "dynamic sports photography, motion, stadium light, cinematic, no text", "heading_font": "Trebuchet MS"},
    "travel": {"bg": (247, 243, 236), "ink": (32, 28, 24), "mid": (80, 70, 60), "mute": (130, 120, 110), "line": (180, 120, 60),
               "photo": "travel photography, cities and landscapes, golden hour, cinematic, no text", "heading_font": "Georgia"},
    "food": {"bg": (252, 248, 242), "ink": (40, 24, 16), "mid": (90, 60, 40), "mute": (140, 110, 90), "line": (180, 80, 40),
             "photo": "food photography, editorial restaurant style, no text, no watermark", "heading_font": "Georgia"},
    "art": {"bg": (20, 18, 22), "ink": (248, 244, 238), "mid": (200, 190, 180), "mute": (140, 130, 125), "line": (220, 180, 120),
            "photo": "art gallery photography, paintings, sculpture, cinematic, no text", "heading_font": "Georgia"},
    "eco": {"bg": (236, 244, 236), "ink": (20, 40, 24), "mid": (50, 80, 55), "mute": (100, 125, 105), "line": (40, 110, 60),
            "photo": "ecology photography, forests, clean energy, cinematic, no text", "heading_font": "Trebuchet MS"},
    "minimal": {"bg": (250, 250, 250), "ink": (18, 18, 18), "mid": (70, 70, 70), "mute": (140, 140, 140), "line": (18, 18, 18),
                "photo": "minimalist photography, clean composition, negative space, no text", "heading_font": "Calibri"},
    "default": {"bg": (248, 246, 241), "ink": (22, 22, 24), "mid": (70, 70, 74), "mute": (130, 128, 124), "line": (22, 22, 24),
                "photo": "cinematic photorealistic photo, no text, no watermark", "heading_font": "Calibri"}
}

THEME_LABELS_I18N = {
    "nature": {"ru": "Природа", "en": "Nature", "de": "Natur", "ar": "الطبيعة", "zh": "自然", "es": "Naturaleza", "fr": "Nature"},
    "business": {"ru": "Бизнес", "en": "Business", "de": "Business", "ar": "الأعمال", "zh": "商业", "es": "Negocios", "fr": "Affaires"},
    "tech": {"ru": "Технологии", "en": "Technology", "de": "Technologie", "ar": "التقنية", "zh": "科技", "es": "Tecnología", "fr": "Technologie"},
    "school": {"ru": "Учёба", "en": "Study", "de": "Studium", "ar": "الدراسة", "zh": "学习", "es": "Estudios", "fr": "Études"},
    "fashion": {"ru": "Мода", "en": "Fashion", "de": "Mode", "ar": "الموضة", "zh": "时尚", "es": "Moda", "fr": "Mode"},
    "history": {"ru": "История", "en": "History", "de": "Geschichte", "ar": "التاريخ", "zh": "历史", "es": "Historia", "fr": "Histoire"},
    "science": {"ru": "Наука", "en": "Science", "de": "Wissenschaft", "ar": "العلوم", "zh": "科学", "es": "Ciencia", "fr": "Science"},
    "sport": {"ru": "Спорт", "en": "Sport", "de": "Sport", "ar": "الرياضة", "zh": "体育", "es": "Deporte", "fr": "Sport"},
    "travel": {"ru": "Путешествия", "en": "Travel", "de": "Reisen", "ar": "السفر", "zh": "旅行", "es": "Viajes", "fr": "Voyage"},
    "food": {"ru": "Еда", "en": "Food", "de": "Essen", "ar": "الطعام", "zh": "美食", "es": "Comida", "fr": "Gastronomie"},
    "art": {"ru": "Искусство", "en": "Art", "de": "Kunst", "ar": "الفن", "zh": "艺术", "es": "Arte", "fr": "Art"},
    "eco": {"ru": "Экология", "en": "Ecology", "de": "Ökologie", "ar": "البيئة", "zh": "生态", "es": "Ecología", "fr": "Écologie"},
    "minimal": {"ru": "Минимализм", "en": "Minimalism", "de": "Minimalismus", "ar": "البساطة", "zh": "极简", "es": "Minimalismo", "fr": "Minimalisme"},
    "default": {"ru": "Универсальный", "en": "Universal", "de": "Universell", "ar": "عام", "zh": "通用", "es": "Universal", "fr": "Universel"},
}
THEME_LABELS = {k: v["ru"] for k, v in THEME_LABELS_I18N.items()}  # обратная совместимость

ANGLES = [
    "через неожиданный факт", "через историю одного примера", "через контраст до и после",
    "через вопрос к зрителю", "через три сильных тезиса", "в стиле National Geographic",
    "в стиле Apple keynote", "в стиле модного журнала", "как лекция сильного преподавателя",
    "как премиальный pitch deck"
]


def pick_theme(topic: str):
    t = (topic or "").lower()
    rules = [
        ("nature", ["животн", "природ", "океан", "кит", "лес", "моря", "птиц", "растен"]),
        ("business", ["бизнес", "компани", "продаж", "финанс", "инвест", "стартап", "рынок"]),
        ("tech", ["крипт", "техно", "ai", "ии", "робот", "софт", "нейро", "код", "гаджет"]),
        ("school", ["школ", "универ", "урок", "студент", "доклад", "класс"]),
        ("fashion", ["мод", "стиль", "бренд", "одежд"]),
        ("history", ["истори", "войн", "древн", "импери", "век"]),
        ("science", ["наук", "физик", "хими", "космос", "медицин", "биологи"]),
        ("sport", ["спорт", "футбол", "тренир", "олимп", "матч"]),
        ("travel", ["путешеств", "город", "страна", "туризм", "поездк"]),
        ("food", ["еда", "кухн", "рецепт", "ресторан", "блюд"]),
        ("art", ["искусств", "живопис", "музей", "театр", "музык"]),
        ("eco", ["эколог", "климат", "мусор", "переработ"]),
        ("minimal", ["минимал", "чисто", "просто"]),
    ]
    for name, keys in rules:
        if any(k in t for k in keys):
            return name, THEMES[name]
    return "default", THEMES["default"]


def get_user(uid):
    if uid not in users_db:
        users_db[uid] = {"name": "", "plan": "premium", "generations": 0, "history": [], "busy": False,
                          "lang": "ru", "lang_chosen": False}
        save_users()
    return users_db[uid]


def user_lang(uid):
    return get_user(uid).get("lang", "ru")


def can_generate(uid):
    u = get_user(uid)
    return u["generations"] < PLAN_LIMITS.get(u["plan"], 15)


# --- Защита от флуда и DDoS ------------------------------------------------
# Три уровня, все проверяются в start_job() — единой точке входа перед любой
# тяжёлой генерацией (презентация/Word/Excel/шаблон):
#   1) rate-limit по времени на одного пользователя — не чаще раза в
#      RATE_LIMIT_SECONDS, независимо от того, завершилась ли предыдущая
#      генерация (per-uid busy-флаг это не покрывал);
#   2) глобальный потолок одновременных генераций по всему боту сразу —
#      защищает от атаки с множества разных Telegram-аккаунтов;
#   3) алерт админу при аномальной нагрузке за скользящий час — суточного
#      лимита нет (по запросу), но при подозрительном всплеске бот сам
#      предупредит владельца, чтобы среагировать вручную до того, как
#      выгорит баланс на xAI.
RATE_LIMIT_SECONDS = 30
FAIL_COOLDOWN_SECONDS = 5  # короткая пауза даже при неудачной генерации - см. start_job()
MAX_CONCURRENT_GENERATIONS = 20
HOURLY_ALERT_THRESHOLD = 50
ALERT_COOLDOWN_SECONDS = 900  # не чаще раза в 15 минут, чтобы не заспамить одним и тем же

_last_request_time = {}          # uid -> monotonic-время последнего успешного старта генерации
_generation_timestamps = deque() # monotonic-метки всех начатых генераций за последний час
_active_generations = 0          # сколько генераций сейчас реально выполняется одновременно
_last_alert_time = float("-inf") # когда в последний раз слали алерт админу.
# ВАЖНО: именно -inf, а не 0.0 — time.monotonic() отсчитывается от произвольной
# точки (не от старта процесса и не от эпохи), и если бы тут был 0.0, проверка
# "now - _last_alert_time > ALERT_COOLDOWN_SECONDS" в первые ~15 минут после
# старта бота была бы ложно False (потому что now ещё меньше 900), и самый
# первый алерт об атаке — то есть ровно тот момент, когда бот только что
# перезапустили и он особенно уязвим — просто не дошёл бы до админа.


def _check_hourly_load():
    """Чистит метки старше часа и, если нагрузка аномально высокая, шлёт
    алерт админу (не чаще раза в ALERT_COOLDOWN_SECONDS)."""
    global _last_alert_time
    now = time.monotonic()
    while _generation_timestamps and now - _generation_timestamps[0] > 3600:
        _generation_timestamps.popleft()
    if len(_generation_timestamps) >= HOURLY_ALERT_THRESHOLD and now - _last_alert_time > ALERT_COOLDOWN_SECONDS:
        _last_alert_time = now
        count = len(_generation_timestamps)
        for admin_id in ADMIN_IDS:
            asyncio.create_task(bot.send_message(
                admin_id,
                f"⚠️ Аномальная нагрузка на бота: {count} генераций за последний час.\n"
                f"Похоже на флуд или атаку — стоит проверить."
            ))


_photo_locks: dict = {}


def _get_photo_lock(uid) -> asyncio.Lock:
    """Блокировка на пользователя для сборки user_photos. Telegram присылает альбом
    из нескольких фото как отдельные сообщения почти одновременно - без этой блокировки
    конкурентные обработчики читают один и тот же список ДО того, как другие успели
    его перезаписать, и "затирают" друг друга (часть фото из альбома молча терялась бы)."""
    lock = _photo_locks.get(uid)
    if lock is None:
        lock = asyncio.Lock()
        _photo_locks[uid] = lock
    return lock


def start_job(uid):
    """Помечает пользователя как занятого дорогой генерацией и проверяет
    защиту от флуда/DDoS (rate-limit + глобальный потолок).
    Возвращает (True, None) при успехе, иначе (False, "текст для пользователя")."""
    global _active_generations
    u = get_user(uid)

    if u.get("busy"):
        return False, "Уже собираю предыдущую версию, подожди немного 🙂"

    now = time.monotonic()
    last = _last_request_time.get(uid, 0)
    if now - last < RATE_LIMIT_SECONDS:
        wait = int(RATE_LIMIT_SECONDS - (now - last)) + 1
        return False, f"Слишком часто 🙂 Подожди ещё {wait} сек. и попробуй снова."

    if _active_generations >= MAX_CONCURRENT_GENERATIONS:
        return False, "Сейчас бот перегружен — очень много запросов одновременно. Попробуй, пожалуйста, через минуту 🙏"

    u["busy"] = True
    u["_counted"] = True
    # Полный 30-секундный лимит взводится только при УСПЕХЕ (note_success ниже
    # перезапишет его временем реального завершения) - но чтобы неудачные попытки
    # не превращались в лазейку для накрутки запросов к платному xAI API без всякой
    # паузы (упавшая генерация мгновенно освобождает busy и разрешает повтор),
    # здесь сразу взводится короткая пауза FAIL_COOLDOWN_SECONDS. Если сборка
    # успеет завершиться успешно раньше - note_success её перезапишет на полную.
    _last_request_time[uid] = now - (RATE_LIMIT_SECONDS - FAIL_COOLDOWN_SECONDS)
    _active_generations += 1
    _generation_timestamps.append(now)
    _check_hourly_load()
    return True, None


def finish_job(uid):
    """Снимает занятость пользователя. Идемпотентна: если по этому uid уже
    финишировали (например, /cancel сработал раньше, чем реально завершилась
    генерация), счётчик _active_generations не уйдёт в минус и не спишется
    дважды за одну и ту же генерацию."""
    global _active_generations
    u = get_user(uid)
    if u.get("_counted"):
        _active_generations = max(0, _active_generations - 1)
        u["_counted"] = False
    u["busy"] = False


def note_success(uid):
    """Фиксирует успешную выдачу файла: лимит 30 сек и запись на диск."""
    _last_request_time[uid] = time.monotonic()
    save_users()


# Telegram режет любое сообщение длиннее 4096 символов и просто не отправляет его
# (ошибка на стороне API), а не обрезает — из-за этого черновик документа при
# "полном раскрытии темы" (когда нейросеть пишет заметно больше, чем короткий план)
# иногда падал в общий обработчик ошибок с "Что-то пошло не так. Попробуй ещё раз".
# send_draft() при коротком черновике отправляет его текстом как раньше, а при
# длинном — красивым PDF-файлом с именем по теме, а не рвёт на несколько сообщений.
TELEGRAM_MAX_LEN = 4000  # с запасом от лимита в 4096


async def send_draft(m: Message, text: str, title: str = "Черновик", reply_markup=None, note: str = ""):
    """Отправляет черновик пользователю: если он укладывается в лимит Telegram —
    обычным сообщением, а если длиннее (типично для "полного раскрытия темы") —
    красивым PDF-файлом с именем по теме документа, чтобы не дробить текст на
    несколько сообщений и не терять читаемость."""
    if len(text) <= TELEGRAM_MAX_LEN:
        await m.answer(text, reply_markup=reply_markup)
        return
    uid = m.from_user.id
    pdf_path = f"/tmp/draft_{uid}.pdf"
    pdf = canvas.Canvas(pdf_path, pagesize=A4)
    w, h = A4
    fn, fb = register_pdf_fonts()
    pdf.setFont(fb, 14)
    pdf.drawString(40, h - 50, safe_filename(title)[:65])
    y = h - 80
    pdf.setFont(fn, 10)
    for para in text.split("\n"):
        if not para.strip():
            y -= 10
            continue
        for line in wrap_lines(para, fn, 10, w - 80):
            if y < 50:
                pdf.showPage()
                y = h - 50
                pdf.setFont(fn, 10)
            pdf.drawString(40, y, line)
            y -= 14
    pdf.save()
    fname = safe_filename(title, fallback="Черновик")
    caption = f"📄 Черновик получился объёмным, поэтому вот файлом{(' — ' + note) if note else ''}"
    await m.answer_document(FSInputFile(pdf_path, filename=f"{fname} - черновик.pdf"), caption=caption, reply_markup=reply_markup)
    try:
        os.remove(pdf_path)
    except OSError:
        pass


def _balanced_json_spans(text: str):
    """Находит все сбалансированные по фигурным скобкам блоки {...} в тексте,
    корректно игнорируя скобки внутри строк."""
    spans = []
    depth = 0
    in_string = False
    escape = False
    start = None
    for i, ch in enumerate(text):
        if in_string:
            if escape:
                escape = False
            elif ch == "\\":
                escape = True
            elif ch == '"':
                in_string = False
            continue
        if ch == '"':
            in_string = True
        elif ch == "{":
            if depth == 0:
                start = i
            depth += 1
        elif ch == "}":
            if depth > 0:
                depth -= 1
                if depth == 0 and start is not None:
                    spans.append((start, i + 1))
    return spans


def extract_json(raw: str):
    """Надёжно достаёт JSON-объект из ответа модели: убирает markdown-обёртку
    ```json ... ``` и перебирает все сбалансированные по скобкам блоки
    (от самого большого к самому маленькому), а не просто первую/последнюю
    скобку в тексте — так лишний текст модели вокруг JSON не ломает разбор."""
    if not raw:
        raise ValueError("Пустой ответ модели")
    text = raw.strip()
    if text.startswith("```"):
        text = text.strip("`")
        if text.lower().startswith("json"):
            text = text[4:]
    spans = _balanced_json_spans(text)
    if not spans:
        raise ValueError("В ответе нет JSON-объекта")
    for start, end in sorted(spans, key=lambda s: s[1] - s[0], reverse=True):
        try:
            return json.loads(text[start:end])
        except json.JSONDecodeError:
            continue
    raise ValueError("Не найден валидный JSON-блок")


GROK_ERROR_PREFIX = "__GROK_ERROR__"

GROK_LANG_HINT = {
    "ru": "", "en": "\n\nВАЖНО: ответь полностью на английском языке (English).",
    "de": "\n\nВАЖНО: ответь полностью на немецком языке (Deutsch).",
    "ar": "\n\nВАЖНО: ответь полностью на арабском языке (العربية).",
    "zh": "\n\nВАЖНО: ответь полностью на китайском языке (中文).",
    "es": "\n\nВАЖНО: ответь полностью на испанском языке (Español).",
    "fr": "\n\nВАЖНО: ответь полностью на французском языке (Français).",
}


def grok_lang_instruction(lang: str) -> str:
    """Суффикс для промпта Grok, требующий ответ на языке пользователя. Для ru — пусто."""
    return GROK_LANG_HINT.get(lang, "")


GROK_JSON_LANG_HINT = {
    "ru": "",
    "en": "\n\nВАЖНО: весь видимый пользователю текст в JSON (title, content, любые названия, подписи) пиши на английском языке (English). Поля image_prompt оставляй на английском в любом случае (это для генерации картинок).",
    "de": "\n\nВАЖНО: весь видимый пользователю текст в JSON (title, content, любые названия, подписи) пиши на немецком языке (Deutsch). Поля image_prompt оставляй на английском в любом случае (это для генерации картинок).",
    "ar": "\n\nВАЖНО: весь видимый пользователю текст в JSON (title, content, любые названия, подписи) пиши на арабском языке (العربية). Поля image_prompt оставляй на английском в любом случае (это для генерации картинок).",
    "zh": "\n\nВАЖНО: весь видимый пользователю текст в JSON (title, content, любые названия, подписи) пиши на китайском языке (中文). Поля image_prompt оставляй на английском в любом случае (это для генерации картинок).",
    "es": "\n\nВАЖНО: весь видимый пользователю текст в JSON (title, content, любые названия, подписи) пиши на испанском языке (Español). Поля image_prompt оставляй на английском в любом случае (это для генерации картинок).",
    "fr": "\n\nВАЖНО: весь видимый пользователю текст в JSON (title, content, любые названия, подписи) пиши на французском языке (Français). Поля image_prompt оставляй на английском в любом случае (это для генерации картинок).",
}


def grok_json_lang_instruction(lang: str) -> str:
    """Как grok_lang_instruction, но уточняет что image_prompt должен остаться на английском (для генерации картинок)."""
    return GROK_JSON_LANG_HINT.get(lang, "")


def content_gen_lang(data: dict, interface_lang: str, mode_key: str = "mode"):
    """Язык, на котором нужно генерировать контент через ИИ - отдельно от языка интерфейса
    (человек может пользоваться ботом на английском, а презентацию просить на русском).
    Возвращает None для режима "свой текст/данные" - там язык уже задан текстом пользователя,
    форсировать его незачем и даже вредно (можно случайно попросить модель перевести текст)."""
    if data.get(mode_key) == "ai":
        return data.get("content_lang") or interface_lang
    return None


async def ask_grok(prompt: str, max_tokens: int = 4000) -> str:
    try:
        r = await client.chat.completions.create(
            model="grok-3",
            messages=[
                {"role": "system", "content": "Ты арт-директор презентаций. Пиши как живой сильный автор, не как нейросеть, и так, чтобы текст не триггерил детекторы ИИ-генерации: без канцелярита и шаблонных фраз («в современном мире», «является», «следует отметить», «данный», «невозможно переоценить», «таким образом», «подводя итог»), без идеально симметричной структуры абзацев и слишком гладких переходов. Чередуй короткие и длинные предложения неравномерно. Заголовки живые. Каждый ответ уникален. Исправляй ошибки. Только русский."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.95,
            max_tokens=max_tokens
        )
        return r.choices[0].message.content
    except Exception as e:
        print("Grok API error:", e)
        return f"{GROK_ERROR_PREFIX}{e}"


def grok_failed(text: str) -> bool:
    return (text or "").startswith(GROK_ERROR_PREFIX)


async def generate_image(prompt: str, path: str) -> bool:
    try:
        headers = {
            "Authorization": f"Bearer {REPLICATE_API_TOKEN}",
            "Content-Type": "application/json"
        }
        async with httpx.AsyncClient(timeout=120) as http:
            create = await http.post(
                "https://api.replicate.com/v1/models/black-forest-labs/flux-schnell/predictions",
                headers=headers,
                json={"input": {"prompt": prompt, "aspect_ratio": "16:9", "output_format": "png"}}
            )
            print("Replicate create:", create.status_code)
            create.raise_for_status()
            get_url = create.json()["urls"]["get"]
            for _ in range(30):
                await asyncio.sleep(2)
                info = (await http.get(get_url, headers=headers)).json()
                if info.get("status") == "succeeded":
                    out = info.get("output")
                    url = out[0] if isinstance(out, list) else out
                    raw = await http.get(str(url))
                    with open(path, "wb") as f:
                        f.write(raw.content)
                    Image.open(path).convert("RGB").save(path, "PNG")
                    print("Image saved:", path)
                    return True
                if info.get("status") in ("failed", "canceled"):
                    print("Replicate failed:", info)
                    return False
        return False
    except Exception as e:
        print("Image generation error:", e)
        return False


def _photo_tint_color(colors: dict):
    """Цвет для тонирования фото на основе акцентного цвета темы (colors["line"]/["ink"]),
    но с яркостью, безопасной для soft-light блендинга. Акцентные цвета тем часто почти
    чёрные (это цвет текста/линий, не для фото) - если тонировать фото им напрямую, тёмный
    тон "давит" все фото в тень сильнее, чем работает коррекция яркости. Поэтому берём только
    оттенок (hue) темы, а яркость и насыщенность тонирующего слоя всегда держим в безопасном
    среднем диапазоне."""
    base = colors.get("line") or colors.get("ink") or (190, 190, 190)
    r, g, b = (c / 255 for c in base)
    h, s, _v = colorsys.rgb_to_hsv(r, g, b)
    s = min(s, 0.5)
    v = 0.64  # чуть выше нейтральной точки soft-light (128/255 ≈ 0.5) - лёгкое, не разрушающее тонирование
    r2, g2, b2 = colorsys.hsv_to_rgb(h, s, v)
    return (int(r2 * 255), int(g2 * 255), int(b2 * 255))


def average_luminance(paths):
    """Средняя яркость по набору фото - используется как общий ориентир, чтобы все
    личные фото пользователя в одной презентации выглядели как единый комплект по свету,
    а не вразнобой (одно тёмное, другое пересвеченное)."""
    vals = []
    for p in paths:
        try:
            im = Image.open(p).convert("L")
            vals.append(ImageStat.Stat(im).mean[0])
        except Exception:
            continue
    return sum(vals) / len(vals) if vals else None


def _detect_subject_region(im: Image.Image):
    """Ищет лицо человека на фото через OpenCV (Haar cascade) и возвращает зону резкости
    (cx, cy, r) в относительных координатах 0..1 от размера фото. Если лиц несколько -
    охватывает все; если лиц нет (пейзаж, еда, предметы) или OpenCV недоступен - None,
    и размытие фона тогда откатится на центр кадра (разумное допущение для большинства
    любительских фото, где объект съёмки обычно в центре)."""
    if _FACE_CASCADE is None:
        return None
    try:
        import numpy as np
        gray = np.asarray(im.convert("L"))
        min_size = max(20, int(min(im.size) * 0.06))
        faces = _FACE_CASCADE.detectMultiScale(gray, scaleFactor=1.1, minNeighbors=5, minSize=(min_size, min_size))
        if len(faces) == 0:
            return None
        x1 = min(x for x, y, w, h in faces)
        y1 = min(y for x, y, w, h in faces)
        x2 = max(x + w for x, y, w, h in faces)
        y2 = max(y + h for x, y, w, h in faces)
        cx, cy = (x1 + x2) / 2 / im.width, (y1 + y2) / 2 / im.height
        span = max(x2 - x1, y2 - y1) / max(im.width, im.height)
        r = min(0.42, max(0.24, span * 1.3))  # запас вокруг лица (тело/плечи), но не более трети кадра
        return cx, cy, r
    except Exception as e:
        print("Face detection error:", e)
        return None


def _apply_depth_blur(im: Image.Image, subject=None, blur_radius: float = 7.0) -> Image.Image:
    """Имитация малой глубины резкости (боке) как на портретном объективе: зона объекта
    (лицо человека, если найдено детекцией, иначе центр кадра) остаётся чёткой, периферия
    кадра плавно размывается. Классический профессиональный портретный приём, визуально
    сразу читается как "обработанное" фото, а не просто вставленное."""
    import numpy as np
    w, h = im.size
    blurred = im.filter(ImageFilter.GaussianBlur(radius=blur_radius))
    y, x = np.ogrid[:h, :w]
    if subject:
        cx_rel, cy_rel, r_rel = subject
        cx, cy = cx_rel * w, cy_rel * h
        r_px = r_rel * max(w, h)
    else:
        cx, cy = w / 2, h / 2
        r_px = 0.40 * max(w, h)
    dist = np.sqrt((x - cx) ** 2 + (y - cy) ** 2)
    # 0 внутри зоны резкости, плавно нарастает до 1 к периферии
    mask = np.clip((dist - r_px) / (0.32 * max(w, h)), 0, 1).astype("float32")
    arr_sharp = np.asarray(im).astype("float32")
    arr_blur = np.asarray(blurred).astype("float32")
    out = arr_sharp * (1 - mask[:, :, None]) + arr_blur * mask[:, :, None]
    return Image.fromarray(out.clip(0, 255).astype("uint8"))


def _apply_vignette(im: Image.Image, strength: float = 0.22) -> Image.Image:
    """Классическое фото-виньетирование - плавное затемнение к краям кадра, притягивает
    взгляд к центру/объекту. Один из самых заметных и при этом безопасных для лиц приёмов
    (центр кадра, где обычно и находится человек, не темнеет вовсе)."""
    import numpy as np
    w, h = im.size
    y, x = np.ogrid[:h, :w]
    cx, cy = w / 2, h / 2
    # расстояние до центра, нормированное так, что края кадра ~= 1.0
    dist = np.sqrt(((x - cx) / (w / 2)) ** 2 + ((y - cy) / (h / 2)) ** 2)
    mask = 1 - np.clip(dist - 0.4, 0, 1) / 0.9 * strength
    mask = np.clip(mask, 1 - strength, 1.0).astype("float32")
    arr = np.asarray(im).astype("float32")
    arr *= mask[:, :, None]
    return Image.fromarray(arr.clip(0, 255).astype("uint8"))


def enhance_user_photo(src_path: str, colors: dict, target_luminance=None) -> Image.Image:
    """Профессиональная доработка ЛИЧНОГО фото пользователя под стиль презентации:
    выравнивает контраст и резкость, подтягивает яркость к общему уровню остальных фото
    в этой же презентации, размывает фон вокруг объекта (эффект боке - как на портретном
    объективе, объект ищется детекцией лица), добавляет заметный, но не разрушающий
    цветовой тон в духе темы (как в Canva/Lightroom presets) и лёгкую виньетку - фото
    перестаёт выглядеть "вставленным как есть", но остаётся собой. Намеренно НЕ используется
    генеративный ИИ (img2img) - на личных фото людей это рискует исказить лица и черты,
    что хуже, чем не трогать фото вовсе."""
    im = Image.open(src_path).convert("RGB")

    # Ищем лицо ДО любых правок - на необработанном изображении детекция надёжнее.
    subject = _detect_subject_region(im)

    # 1. Автоконтраст - тянет тусклые/плоские фото к полному диапазону тонов
    im = ImageOps.autocontrast(im, cutoff=2)

    # 2. Повышение насыщенности и контраста - делает фото более "глянцевым" и заметно живее
    im = ImageEnhance.Color(im).enhance(1.2)
    im = ImageEnhance.Contrast(im).enhance(1.12)

    # 3. Подстройка яркости к общему уровню комплекта фото (если он известен)
    if target_luminance is not None:
        cur = ImageStat.Stat(im.convert("L")).mean[0]
        if cur > 0:
            factor = 1 + 0.5 * ((target_luminance - cur) / 255)
            factor = max(0.75, min(1.35, factor))
            im = ImageEnhance.Brightness(im).enhance(factor)

    # 4. Доводка резкости/чёткости - делаем ДО размытия фона, иначе шарпен на следующем
    # шаге частично "отменил" бы размытие периферии.
    im = im.filter(ImageFilter.UnsharpMask(radius=1.6, percent=90, threshold=2))

    # 5. Размытие фона вокруг объекта (боке) - объект (лицо/люди, если найдены) остаётся
    # чётким, периферия кадра мягко уходит в размытие, как при съёмке с открытой диафрагмой.
    im = _apply_depth_blur(im, subject=subject, blur_radius=10.0)

    # 6. Заметный цветовой тон в духе темы презентации (soft-light блендинг на 28%) -
    # объединяет разномастные фото пользователя визуально с остальными слайдами.
    tint_color = _photo_tint_color(colors)
    tint_layer = Image.new("RGB", im.size, tint_color)
    graded = ImageChops.soft_light(im, tint_layer)
    im = Image.blend(im, graded, 0.28)

    # 7. Виньетка - лёгкое затемнение к краям, придаёт фото "обработанный" вид с первого взгляда
    im = _apply_vignette(im, strength=0.22)

    return im


def prepare_user_photo(src_path: str, colors: dict, target_luminance=None) -> str:
    """Доработка + сохранение в новый временный файл, путь к которому можно передавать
    дальше в cover() как обычный источник изображения."""
    im = enhance_user_photo(src_path, colors, target_luminance)
    out_path = src_path.rsplit(".", 1)[0] + "_enhanced.png"
    im.save(out_path, "PNG")
    return out_path


def cover(src, dest, w, h):
    im = Image.open(src).convert("RGB")
    t = w / h
    iw, ih = im.size
    if iw / ih > t:
        nw = int(ih * t)
        x = (iw - nw) // 2
        im = im.crop((x, 0, x + nw, ih))
    else:
        nh = int(iw / t)
        y = (ih - nh) // 2
        im = im.crop((0, y, iw, y + nh))
    im.resize((w, h), Image.LANCZOS).save(dest, "PNG")


def _pptx_tint(color_tuple, factor=0.9):
    """Осветляет RGB-цвет темы к белому (0..1) - для мягкой заливки карточек в раскладке
    'cards', аналог _xl_tint для Excel."""
    r, g, b = color_tuple
    return (int(r + (255 - r) * factor), int(g + (255 - g) * factor), int(b + (255 - b) * factor))


def slide_background(slide, colors):
    """Фон слайда - мягкий диагональный градиент в сторону акцентного цвета темы, а не
    плоская сплошная заливка (раньше каждый слайд был одним и тем же прямоугольником цвета -
    визуально скучно и все презентации выглядели "плоско"). Разница между стоп-цветами
    держится небольшой (12%), чтобы не спорить по контрасту с текстом поверх фона."""
    bg = colors["bg"]
    accent = colors["line"]
    blend = 0.12
    bg2 = tuple(int(bg[i] + (accent[i] - bg[i]) * blend) for i in range(3))
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.line.fill.background()
    shape.fill.gradient()
    stops = shape.fill.gradient_stops
    stops[0].color.rgb = RGBColor(*bg)
    stops[0].position = 0.0
    stops[-1].color.rgb = RGBColor(*bg2)
    stops[-1].position = 1.0
    try:
        shape.fill.gradient_angle = 45
    except Exception:
        pass
    return shape


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(*color)
    s.line.fill.background()


def txt(slide, l, t, w, h, text, size, color, bold=False, font_name="Calibri"):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text or ""
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    p.font.name = font_name
    return box


def add_chart(slide, l, t, w, h, chart_data_dict, colors):
    """Рисует нативный график PowerPoint (bar/line/pie) на слайде вместо фото.
    chart_data_dict: {"chart_type": "bar"/"line"/"pie", "title": "...",
    "categories": [...], "series": [{"name":..., "values":[...]}]}.
    Используется только когда модель сама решила, что тема подразумевает цифры -
    см. инструкцию про ключ "chart" в промпте _build_presentation."""
    ctype = chart_data_dict.get("chart_type", "bar")
    cats = chart_data_dict.get("categories") or []
    series_list = chart_data_dict.get("series") or []
    if not cats or not series_list:
        return None

    cd = CategoryChartData()
    cd.categories = cats
    for s in series_list:
        cd.add_series(s.get("name", ""), s.get("values", []))

    xl_type = {
        "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }.get(ctype, XL_CHART_TYPE.COLUMN_CLUSTERED)

    graphic_frame = slide.shapes.add_chart(xl_type, Inches(l), Inches(t), Inches(w), Inches(h), cd)
    chart = graphic_frame.chart

    # чем уже панель под график - тем мельче шрифты, иначе легенда/подписи не влезают
    compact = w < 5.0
    title_sz, legend_sz, label_sz, tick_sz = (13, 9, 9, 9) if compact else (15, 11, 11, 10)

    palette = [colors["line"], colors["mid"], colors["mute"], colors["ink"]]

    if chart_data_dict.get("title"):
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_data_dict["title"]
        for run in chart.chart_title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(title_sz)
            run.font.bold = True
            run.font.color.rgb = RGBColor(*colors["ink"])
            run.font.name = "Calibri"
    else:
        chart.has_title = False

    is_pie = ctype == "pie"
    multi_series = len(series_list) > 1

    if is_pie or multi_series:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(legend_sz)
        chart.legend.font.color.rgb = RGBColor(*colors["mid"])
        chart.legend.font.name = "Calibri"
    else:
        chart.has_legend = False

    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(label_sz)
    dl.font.color.rgb = RGBColor(*colors["ink"])
    dl.font.name = "Calibri"
    if is_pie:
        dl.number_format = '0%'
        dl.number_format_is_linked = False
        dl.show_percentage = True
        dl.show_value = False
    else:
        dl.show_value = True

    for i, series in enumerate(chart.series):
        color = palette[i % len(palette)]
        if is_pie:
            for j, point in enumerate(series.points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = RGBColor(*palette[j % len(palette)])
        elif ctype == "line":
            series.format.line.color.rgb = RGBColor(*color)
            series.format.line.width = Pt(2.5)
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor(*color)

    if not is_pie:
        cat_ax = chart.category_axis
        val_ax = chart.value_axis
        for ax in (cat_ax, val_ax):
            ax.tick_labels.font.size = Pt(tick_sz)
            ax.tick_labels.font.color.rgb = RGBColor(*colors["mute"])
            ax.tick_labels.font.name = "Calibri"
            ax.format.line.color.rgb = RGBColor(*colors["mute"])
        val_ax.has_major_gridlines = False

    return graphic_frame




def main_kb(lang="ru"):
    return ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text=tr("btn_pres", lang))],
        [KeyboardButton(text=tr("btn_word", lang))],
        [KeyboardButton(text=tr("btn_excel", lang))],
        [KeyboardButton(text=tr("btn_history", lang)), KeyboardButton(text=tr("btn_plan", lang))],
        [KeyboardButton(text=tr("btn_language", lang))]
    ], resize_keyboard=True)


def cancel_kb(lang="ru"):
    """Клавиатура-заглушка с одной кнопкой возврата в меню - используется в шагах,
    где пользователь должен прислать свой текст (тему, данные, правки), но всё равно
    должен иметь возможность выйти в главное меню в один тап, а не только через /cancel."""
    return ReplyKeyboardMarkup(keyboard=[[KeyboardButton(text=tr("btn_main_menu", lang))]], resize_keyboard=True)


def mode_kb(show_template=False, lang="ru"):
    rows = [
        [KeyboardButton(text=tr("btn_ai_generate", lang))],
        [KeyboardButton(text=tr("btn_own_text", lang))],
    ]
    if show_template:
        rows.append([KeyboardButton(text=tr("btn_template", lang))])
    rows.append([KeyboardButton(text=tr("btn_main_menu", lang))])
    return ReplyKeyboardMarkup(keyboard=rows, resize_keyboard=True)


def slides_kb(lang="ru"):
    return ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text=tr("btn_8slides", lang)), KeyboardButton(text=tr("btn_12slides", lang))],
        [KeyboardButton(text=tr("btn_16slides", lang))],
        [KeyboardButton(text=tr("btn_main_menu", lang))]
    ], resize_keyboard=True)


def confirm_kb(lang="ru"):
    return ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text=tr("btn_full_version", lang))],
        [KeyboardButton(text=tr("btn_add_info", lang))],
        [KeyboardButton(text=tr("btn_change_style", lang))],
        [KeyboardButton(text=tr("btn_change_topic", lang))],
        [KeyboardButton(text=tr("btn_main_menu", lang))]
    ], resize_keyboard=True)


def photo_source_kb(lang="ru"):
    return ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text=tr("btn_ai_photos", lang))],
        [KeyboardButton(text=tr("btn_own_photos", lang))],
        [KeyboardButton(text=tr("btn_main_menu", lang))]
    ], resize_keyboard=True)


def photos_done_kb(lang="ru"):
    return ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text=tr("btn_photos_done", lang))],
        [KeyboardButton(text=tr("btn_main_menu", lang))]
    ], resize_keyboard=True)


# Полный банк из 30 видов документов, разбитых на 3 категории (физлица/юрлица/учёба).
# В каждой категории есть "основные" (видны сразу) и "дополнительные" (открываются
# кнопкой "Показать ещё документы"). Итого 8 основных на весь бот + 22 дополнительных.
STUDY_KINDS = {"doc", "referat", "report", "essay", "notes", "coursework"}

# Инструкция, снижающая вероятность срабатывания детекторов ИИ-текста (не антиплагиата —
# антиплагиат по совпадениям и так не страдает, т.к. текст пишется с нуля). Применяется
# только к учебным документам и презентациям, где стиль текста имеет значение;
# юридические документы намеренно формальные и канцелярские по своей природе, их не трогаем.
ANTI_AI_DETECTOR_STYLE = (
    "Пиши так, чтобы текст не выглядел сгенерированным нейросетью и не триггерил "
    "детекторы ИИ-текста: избегай слишком гладких переходов и идеально симметричной "
    "структуры абзацев, чередуй короткие и длинные предложения неравномерно, не начинай "
    "несколько абзацев одинаковыми конструкциями, избегай канцелярских клише "
    "('в современном мире', 'является', 'следует отметить', 'таким образом', "
    "'нельзя переоценить', 'играет важную роль', 'подводя итог'). "
    "Формулировки должны звучать как у обычного студента: где уместно — чуть менее "
    "выверенные обороты, конкретные примеры и детали по теме, живая логика рассуждения, "
    "а не шаблонная академическая гладкость."
)

# Минимальный суммарный объём (в словах) для проверки после генерации - см.
# комментарий у WORD_MIN_TOTAL_WORDS.get(...) в word_build. Соответствует нижней
# границе из length_hint, но отдельной константой, чтобы не парсить текст промпта.
WORD_MIN_TOTAL_WORDS = {
    ("coursework", "long"): 6000,
    ("coursework", "short"): 2000,
    ("referat", "long"): 3000,
    ("referat", "short"): 1200,
    ("report", "long"): 1500,
    ("report", "short"): 600,
    ("essay", "long"): 1500,
    ("essay", "short"): 600,
}

# Минимум слов на один содержательный раздел - используется в довыворота-промпте,
# чтобы указывать модели точную недостачу по каждому разделу, а не только по сумме.
WORD_SECTION_MIN_WORDS = {
    ("coursework", "long"): 600,
    ("coursework", "short"): 200,
    ("referat", "long"): 400,
    ("referat", "short"): 150,
    ("report", "long"): 250,
    ("report", "short"): 100,
    ("essay", "long"): 250,
    ("essay", "short"): 100,
}
# Заголовки, которые не нужно искусственно раздувать при довыворота (оглавление,
# список источников - у них естественно фиксированный, а не текстовый объём).
_WORD_SECTION_SKIP_EXPAND = ("содержание", "список литератур", "список источник", "титульный лист", "задание")

WORD_CATEGORY_TITLES = {
    "physical": {"ru": "👤 Для физлиц", "en": "👤 For individuals", "de": "👤 Für Privatpersonen",
                 "ar": "👤 للأفراد", "zh": "👤 个人用户", "es": "👤 Para particulares", "fr": "👤 Pour particuliers"},
    "entity": {"ru": "🏢 Для юрлиц / ИП", "en": "🏢 For businesses / sole traders", "de": "🏢 Für Unternehmen / Selbstständige",
               "ar": "🏢 للشركات / أصحاب الأعمال", "zh": "🏢 企业/个体户", "es": "🏢 Para empresas / autónomos", "fr": "🏢 Pour entreprises / indépendants"},
    "study": {"ru": "🎓 Для учёбы", "en": "🎓 For studies", "de": "🎓 Für Studium/Schule",
              "ar": "🎓 للدراسة", "zh": "🎓 学习用", "es": "🎓 Para estudios", "fr": "🎓 Pour les études"},
}

WORD_CATEGORIES = {
    "physical": {
        "title": WORD_CATEGORY_TITLES["physical"]["ru"],
        "main": ["dkp", "rent", "proxy", "statement"],
        "more": ["act", "loan", "claim", "consent", "marriage_contract", "gift", "lawsuit", "alimony"],
    },
    "entity": {
        "title": WORD_CATEGORY_TITLES["entity"]["ru"],
        "main": ["offer", "services", "employment", "work_act"],
        "more": ["supply", "agency", "joint_activity", "nonresidential_rent", "cession", "nda",
                 "self_employed", "warranty_letter"],
    },
    "study": {
        "title": WORD_CATEGORY_TITLES["study"]["ru"],
        "main": ["doc", "referat", "report", "essay", "notes", "coursework"],
        "more": [],
    },
}

# Название вида документа на каждом из 7 языков - используется и в кнопках выбора,
# и как fallback-заголовок сгенерированного файла (если модель не даст лучше).
WORD_KIND_LABELS = {
    "doc": {"ru": "📄 Обычный документ", "en": "📄 Generic document", "de": "📄 Allgemeines Dokument",
            "ar": "📄 مستند عام", "zh": "📄 普通文档", "es": "📄 Documento genérico", "fr": "📄 Document générique"},
    "referat": {"ru": "🎓 Реферат", "en": "🎓 Research paper", "de": "🎓 Referat",
                "ar": "🎓 بحث", "zh": "🎓 学术论文", "es": "🎓 Trabajo de investigación", "fr": "🎓 Dossier de recherche"},
    "report": {"ru": "🎤 Доклад", "en": "🎤 Report", "de": "🎤 Vortrag",
               "ar": "🎤 تقرير", "zh": "🎤 报告", "es": "🎤 Exposición", "fr": "🎤 Exposé"},
    "essay": {"ru": "✍️ Эссе", "en": "✍️ Essay", "de": "✍️ Essay",
              "ar": "✍️ مقال", "zh": "✍️ 论文/短文", "es": "✍️ Ensayo", "fr": "✍️ Essai"},
    "notes": {"ru": "📓 Конспект", "en": "📓 Study notes", "de": "📓 Konspekt",
              "ar": "📓 ملخص دراسي", "zh": "📓 课堂笔记", "es": "📓 Apuntes", "fr": "📓 Notes de cours"},
    "coursework": {"ru": "📚 Курсовая работа", "en": "📚 Term paper", "de": "📚 Hausarbeit",
                   "ar": "📚 بحث فصلي", "zh": "📚 学期论文", "es": "📚 Trabajo de curso", "fr": "📚 Dossier de cours"},
    "dkp": {"ru": "📝 Договор купли-продажи", "en": "📝 Sale and purchase agreement", "de": "📝 Kaufvertrag",
            "ar": "📝 عقد بيع", "zh": "📝 买卖合同", "es": "📝 Contrato de compraventa", "fr": "📝 Contrat de vente"},
    "rent": {"ru": "🏠 Договор аренды", "en": "🏠 Rental agreement", "de": "🏠 Mietvertrag",
             "ar": "🏠 عقد إيجار", "zh": "🏠 租赁合同", "es": "🏠 Contrato de alquiler", "fr": "🏠 Contrat de location"},
    "proxy": {"ru": "🧾 Доверенность", "en": "🧾 Power of attorney", "de": "🧾 Vollmacht",
              "ar": "🧾 توكيل", "zh": "🧾 授权委托书", "es": "🧾 Poder notarial", "fr": "🧾 Procuration"},
    "statement": {"ru": "📋 Заявление", "en": "📋 Application", "de": "📋 Antrag",
                  "ar": "📋 طلب", "zh": "📋 申请书", "es": "📋 Solicitud", "fr": "📋 Demande"},
    "act": {"ru": "📦 Акт приёма-передачи", "en": "📦 Handover certificate", "de": "📦 Übergabeprotokoll",
            "ar": "📦 محضر تسليم واستلام", "zh": "📦 交接单", "es": "📦 Acta de entrega", "fr": "📦 Procès-verbal de remise"},
    "loan": {"ru": "💰 Расписка / договор займа", "en": "💰 Promissory note / loan agreement", "de": "💰 Schuldschein / Darlehensvertrag",
             "ar": "💰 سند دين / عقد قرض", "zh": "💰 借条/借款合同", "es": "💰 Pagaré / contrato de préstamo", "fr": "💰 Reconnaissance de dette / contrat de prêt"},
    "claim": {"ru": "⚠️ Претензия", "en": "⚠️ Complaint letter", "de": "⚠️ Mängelrüge",
              "ar": "⚠️ مطالبة", "zh": "⚠️ 索赔函", "es": "⚠️ Reclamación", "fr": "⚠️ Réclamation"},
    "consent": {"ru": "✈️ Согласие на выезд ребёнка", "en": "✈️ Child travel consent", "de": "✈️ Zustimmung zur Ausreise des Kindes",
                "ar": "✈️ موافقة على سفر الطفل", "zh": "✈️ 儿童出境同意书", "es": "✈️ Autorización de viaje para menores", "fr": "✈️ Autorisation de voyage pour mineur"},
    "marriage_contract": {"ru": "💍 Брачный договор", "en": "💍 Prenuptial agreement", "de": "💍 Ehevertrag",
                          "ar": "💍 عقد زواج (اتفاقية)", "zh": "💍 婚前协议", "es": "💍 Capitulaciones matrimoniales", "fr": "💍 Contrat de mariage"},
    "gift": {"ru": "🎁 Договор дарения", "en": "🎁 Gift deed", "de": "🎁 Schenkungsvertrag",
             "ar": "🎁 عقد هبة", "zh": "🎁 赠与合同", "es": "🎁 Contrato de donación", "fr": "🎁 Acte de donation"},
    "lawsuit": {"ru": "⚖️ Исковое заявление", "en": "⚖️ Statement of claim", "de": "⚖️ Klageschrift",
                "ar": "⚖️ لائحة دعوى", "zh": "⚖️ 起诉状", "es": "⚖️ Demanda judicial", "fr": "⚖️ Assignation en justice"},
    "alimony": {"ru": "👶 Соглашение об алиментах", "en": "👶 Child support agreement", "de": "👶 Unterhaltsvereinbarung",
                "ar": "👶 اتفاقية نفقة", "zh": "👶 抚养费协议", "es": "👶 Acuerdo de manutención", "fr": "👶 Convention de pension alimentaire"},
    "offer": {"ru": "💼 Коммерческое предложение", "en": "💼 Commercial proposal", "de": "💼 Geschäftsangebot",
              "ar": "💼 عرض تجاري", "zh": "💼 商业提案", "es": "💼 Propuesta comercial", "fr": "💼 Proposition commerciale"},
    "services": {"ru": "🤝 Договор оказания услуг", "en": "🤝 Service agreement", "de": "🤝 Dienstleistungsvertrag",
                 "ar": "🤝 عقد تقديم خدمات", "zh": "🤝 服务合同", "es": "🤝 Contrato de prestación de servicios", "fr": "🤝 Contrat de prestation de services"},
    "employment": {"ru": "👔 Трудовой договор", "en": "👔 Employment contract", "de": "👔 Arbeitsvertrag",
                   "ar": "👔 عقد عمل", "zh": "👔 劳动合同", "es": "👔 Contrato de trabajo", "fr": "👔 Contrat de travail"},
    "work_act": {"ru": "✅ Акт выполненных работ", "en": "✅ Certificate of completed work", "de": "✅ Abnahmeprotokoll",
                 "ar": "✅ محضر إنجاز الأعمال", "zh": "✅ 工作完成验收单", "es": "✅ Acta de obra terminada", "fr": "✅ Procès-verbal de réception des travaux"},
    "supply": {"ru": "🚚 Договор поставки", "en": "🚚 Supply agreement", "de": "🚚 Liefervertrag",
               "ar": "🚚 عقد توريد", "zh": "🚚 供货合同", "es": "🚚 Contrato de suministro", "fr": "🚚 Contrat de fourniture"},
    "agency": {"ru": "🕴️ Агентский договор", "en": "🕴️ Agency agreement", "de": "🕴️ Agenturvertrag",
               "ar": "🕴️ عقد وكالة", "zh": "🕴️ 代理合同", "es": "🕴️ Contrato de agencia", "fr": "🕴️ Contrat d'agence"},
    "joint_activity": {"ru": "🤝 Договор о совместной деятельности", "en": "🤝 Joint activity agreement", "de": "🤝 Vertrag über gemeinsame Tätigkeit",
                       "ar": "🤝 عقد نشاط مشترك", "zh": "🤝 合作协议", "es": "🤝 Contrato de actividad conjunta", "fr": "🤝 Contrat d'activité commune"},
    "nonresidential_rent": {"ru": "🏢 Аренда нежилого помещения", "en": "🏢 Commercial premises lease", "de": "🏢 Gewerbemietvertrag",
                            "ar": "🏢 عقد إيجار عقار غير سكني", "zh": "🏢 非居住用房租赁合同", "es": "🏢 Alquiler de local comercial", "fr": "🏢 Bail commercial"},
    "cession": {"ru": "🔄 Договор цессии", "en": "🔄 Assignment of claim agreement", "de": "🔄 Abtretungsvertrag",
                "ar": "🔄 عقد حوالة حق", "zh": "🔄 债权转让合同", "es": "🔄 Contrato de cesión de derechos", "fr": "🔄 Contrat de cession de créance"},
    "nda": {"ru": "🔒 Соглашение о неразглашении (NDA)", "en": "🔒 Non-disclosure agreement (NDA)", "de": "🔒 Vertraulichkeitsvereinbarung (NDA)",
            "ar": "🔒 اتفاقية عدم إفصاح", "zh": "🔒 保密协议", "es": "🔒 Acuerdo de confidencialidad (NDA)", "fr": "🔒 Accord de confidentialité (NDA)"},
    "self_employed": {"ru": "🧑‍💻 Договор с самозанятым", "en": "🧑‍💻 Contract with a self-employed contractor", "de": "🧑‍💻 Vertrag mit Selbstständigem",
                      "ar": "🧑‍💻 عقد مع عامل مستقل", "zh": "🧑‍💻 与自由职业者合同", "es": "🧑‍💻 Contrato con autónomo", "fr": "🧑‍💻 Contrat avec un travailleur indépendant"},
    "warranty_letter": {"ru": "✉️ Гарантийное письмо", "en": "✉️ Letter of guarantee", "de": "✉️ Garantieschreiben",
                        "ar": "✉️ خطاب ضمان", "zh": "✉️ 保证函", "es": "✉️ Carta de garantía", "fr": "✉️ Lettre de garantie"},
}
KIND_LABELS = {k: v["ru"] for k, v in WORD_KIND_LABELS.items()}  # обратная совместимость (заголовки по умолчанию и т.п.)
# Обратный словарь label->kind собирается по ВСЕМ языкам сразу, чтобы кнопка, нажатая
# на любом из 7 языков, узнавалась одинаково, независимо от текущего языка пользователя.
LABEL_TO_KIND = {label: k for k, variants in WORD_KIND_LABELS.items() for label in variants.values()}


def word_category_kb(lang="ru"):
    return ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text=WORD_CATEGORY_TITLES["physical"][lang if lang in WORD_CATEGORY_TITLES["physical"] else "ru"])],
        [KeyboardButton(text=WORD_CATEGORY_TITLES["entity"][lang if lang in WORD_CATEGORY_TITLES["entity"] else "ru"])],
        [KeyboardButton(text=WORD_CATEGORY_TITLES["study"][lang if lang in WORD_CATEGORY_TITLES["study"] else "ru"])],
        [KeyboardButton(text=tr("btn_main_menu", lang))],
    ], resize_keyboard=True)


WORD_TITLE_TO_CATEGORY = {label: cat for cat, variants in WORD_CATEGORY_TITLES.items() for label in variants.values()}


def word_kind_kb(category, more=False, lang="ru"):
    cat = WORD_CATEGORIES.get(category, WORD_CATEGORIES["physical"])
    ids = list(cat["main"]) + (list(cat["more"]) if more else [])
    rows = [[KeyboardButton(text=WORD_KIND_LABELS[k].get(lang, WORD_KIND_LABELS[k]["ru"]))] for k in ids]
    if cat["more"] and not more:
        rows.append([KeyboardButton(text=tr("btn_more_docs", lang))])
    rows.append([KeyboardButton(text=tr("btn_back_categories", lang))])
    rows.append([KeyboardButton(text=tr("btn_main_menu", lang))])
    return ReplyKeyboardMarkup(keyboard=rows, resize_keyboard=True)


def word_size_kb(lang="ru"):
    return ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text=tr("btn_shallow", lang))],
        [KeyboardButton(text=tr("btn_deep", lang))],
        [KeyboardButton(text=tr("btn_main_menu", lang))]
    ], resize_keyboard=True)


def word_confirm_kb(lang="ru"):
    return ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text=tr("btn_build_doc", lang))],
        [KeyboardButton(text=tr("btn_add_info", lang))],
        [KeyboardButton(text=tr("btn_change_query", lang))],
        [KeyboardButton(text=tr("btn_main_menu", lang))]
    ], resize_keyboard=True)


EXCEL_CATEGORIES = {
    "physical": {
        "title": WORD_CATEGORY_TITLES["physical"]["ru"],
        "kinds": ["expense_estimate", "family_budget", "startup_model"],
    },
    "entity": {
        "title": WORD_CATEGORY_TITLES["entity"]["ru"],
        "kinds": ["project_budget", "price_list"],
    },
    "study": {
        "title": WORD_CATEGORY_TITLES["study"]["ru"],
        "kinds": ["calc_table"],
    },
}

EXCEL_KIND_LABELS_I18N = {
    "expense_estimate": {"ru": "🧾 Смета расходов", "en": "🧾 Expense estimate", "de": "🧾 Kostenvoranschlag",
                         "ar": "🧾 تقدير المصروفات", "zh": "🧾 支出预算表", "es": "🧾 Presupuesto de gastos", "fr": "🧾 Devis de dépenses"},
    "family_budget": {"ru": "💰 Семейный бюджет", "en": "💰 Family budget", "de": "💰 Familienbudget",
                      "ar": "💰 ميزانية الأسرة", "zh": "💰 家庭预算", "es": "💰 Presupuesto familiar", "fr": "💰 Budget familial"},
    "startup_model": {"ru": "🚀 Финмодель стартапа", "en": "🚀 Startup financial model", "de": "🚀 Startup-Finanzmodell",
                      "ar": "🚀 نموذج مالي لشركة ناشئة", "zh": "🚀 创业财务模型", "es": "🚀 Modelo financiero de startup", "fr": "🚀 Modèle financier de startup"},
    "project_budget": {"ru": "📊 Смета проекта / бизнес-план", "en": "📊 Project budget / business plan", "de": "📊 Projektbudget / Businessplan",
                       "ar": "📊 ميزانية مشروع / خطة عمل", "zh": "📊 项目预算/商业计划", "es": "📊 Presupuesto de proyecto / plan de negocio", "fr": "📊 Budget de projet / plan d'affaires"},
    "price_list": {"ru": "🏷️ Прайс-лист", "en": "🏷️ Price list", "de": "🏷️ Preisliste",
                   "ar": "🏷️ قائمة الأسعار", "zh": "🏷️ 价目表", "es": "🏷️ Lista de precios", "fr": "🏷️ Liste de prix"},
    "calc_table": {"ru": "📐 Расчётная таблица к работе", "en": "📐 Calculation table for coursework", "de": "📐 Berechnungstabelle für die Arbeit",
                   "ar": "📐 جدول حسابي للبحث", "zh": "📐 作业计算表", "es": "📐 Tabla de cálculo para el trabajo", "fr": "📐 Tableau de calcul pour le devoir"},
}
EXCEL_KIND_LABELS = {k: v["ru"] for k, v in EXCEL_KIND_LABELS_I18N.items()}  # обратная совместимость
EXCEL_LABEL_TO_KIND = {label: k for k, variants in EXCEL_KIND_LABELS_I18N.items() for label in variants.values()}

# У финмодели стартапа данные всегда реальные, от пользователя - модель их не придумывает,
# только оформляет реальными формулами Excel (см. решение в этой сессии). Остальные виды
# работают как таблицы в презентациях/курсовых: можно как сгенерировать иллюстративные цифры
# по теме, так и вставить свои реальные данные - и в том, и в другом случае формулы настоящие.
EXCEL_KIND_HINTS_I18N = {
    "expense_estimate": {
        "ai": {"ru": "На что смета? Например: смета на ремонт кухни, смета на свадьбу на 40 человек.",
               "en": "What is the estimate for? E.g. kitchen renovation, a wedding for 40 people.",
               "de": "Wofür ist der Kostenvoranschlag? Z. B. Küchenrenovierung, Hochzeit für 40 Personen.",
               "ar": "لماذا التقدير؟ مثلاً: تجديد المطبخ، حفل زفاف لـ٤٠ شخصاً.",
               "zh": "预算用于什么？例如：厨房装修、40人婚礼。",
               "es": "¿Para qué es el presupuesto? Por ejemplo: reforma de cocina, boda para 40 personas.",
               "fr": "Devis pour quoi ? Ex. : rénovation de cuisine, mariage pour 40 personnes."},
        "user": {"ru": "Пришли список статей расходов с количеством и ценой (например: линолеум - 25 м² - 800 ₽) - я оформлю в таблицу с формулами.",
                 "en": "Send a list of expense items with quantity and price (e.g. linoleum - 25 m² - $80) - I'll format it into a table with formulas.",
                 "de": "Sende eine Liste der Ausgabenposten mit Menge und Preis (z. B. Linoleum - 25 m² - 80 €) - ich erstelle eine Tabelle mit Formeln.",
                 "ar": "أرسل قائمة ببنود المصروفات مع الكمية والسعر (مثال: أرضية فينيل - ٢٥ م² - ٨٠ دولار) - سأنسقها في جدول بمعادلات.",
                 "zh": "请发送费用项目清单（含数量和单价），例如：地板革-25平方米-80美元，我会做成带公式的表格。",
                 "es": "Envía la lista de partidas de gasto con cantidad y precio (ej.: linóleo - 25 m² - 80 €) - lo pondré en una tabla con fórmulas.",
                 "fr": "Envoie la liste des postes de dépense avec quantité et prix (ex. : linoléum - 25 m² - 80 €) - je mettrai ça en tableau avec formules."},
    },
    "family_budget": {
        "ai": {"ru": "Опиши примерный бюджет, который нужно смоделировать (например: бюджет семьи из 3 человек на месяц).",
               "en": "Describe the budget to model (e.g. monthly budget for a family of 3).",
               "de": "Beschreibe das Budget, das modelliert werden soll (z. B. Monatsbudget für eine 3-köpfige Familie).",
               "ar": "صف الميزانية المطلوب نمذجتها (مثال: ميزانية شهرية لأسرة من ٣ أفراد).",
               "zh": "描述需要建模的预算（例如：三口之家的月度预算）。",
               "es": "Describe el presupuesto a modelar (ej.: presupuesto mensual para una familia de 3).",
               "fr": "Décris le budget à modéliser (ex. : budget mensuel pour une famille de 3 personnes)."},
        "user": {"ru": "Пришли свои статьи доходов и расходов с суммами - я оформлю в таблицу с формулами (итого доходы/расходы/остаток).",
                 "en": "Send your income and expense items with amounts - I'll format a table with formulas (total income/expenses/balance).",
                 "de": "Sende deine Einnahmen- und Ausgabenposten mit Beträgen - ich erstelle eine Tabelle mit Formeln (Summe Einnahmen/Ausgaben/Rest).",
                 "ar": "أرسل بنود دخلك ومصروفاتك مع المبالغ - سأنسقها في جدول بمعادلات (إجمالي الدخل/المصروفات/المتبقي).",
                 "zh": "请发送你的收支项目及金额，我会做成带公式的表格（总收入/支出/结余）。",
                 "es": "Envía tus partidas de ingresos y gastos con importes - lo pondré en una tabla con fórmulas (total ingresos/gastos/saldo).",
                 "fr": "Envoie tes postes de revenus et dépenses avec les montants - je ferai un tableau avec formules (total revenus/dépenses/solde)."},
    },
    "project_budget": {
        "ai": {"ru": "На какой проект нужна смета/бизнес-план? Опиши в двух словах.",
               "en": "Which project needs the budget/business plan? Describe briefly.",
               "de": "Für welches Projekt wird der Kostenvoranschlag/Businessplan benötigt? Kurz beschreiben.",
               "ar": "لأي مشروع تحتاج الميزانية/خطة العمل؟ صف بإيجاز.",
               "zh": "哪个项目需要预算/商业计划？请简要描述。",
               "es": "¿Para qué proyecto es el presupuesto/plan de negocio? Descríbelo brevemente.",
               "fr": "Pour quel projet faut-il le budget/plan d'affaires ? Décris brièvement."},
        "user": {"ru": "Пришли статьи доходов и расходов проекта с суммами - оформлю в таблицу с формулами (итого доходы/расходы/прибыль).",
                 "en": "Send the project's income and expense items with amounts - I'll format a table with formulas (total income/expenses/profit).",
                 "de": "Sende die Einnahmen- und Ausgabenposten des Projekts mit Beträgen - ich erstelle eine Tabelle mit Formeln (Summe Einnahmen/Ausgaben/Gewinn).",
                 "ar": "أرسل بنود دخل ومصروفات المشروع مع المبالغ - سأنسقها في جدول بمعادلات (إجمالي الدخل/المصروفات/الربح).",
                 "zh": "请发送项目的收支项目及金额，我会做成带公式的表格（总收入/支出/利润）。",
                 "es": "Envía las partidas de ingresos y gastos del proyecto con importes - lo pondré en una tabla con fórmulas (total ingresos/gastos/beneficio).",
                 "fr": "Envoie les postes de revenus et dépenses du projet avec les montants - je ferai un tableau avec formules (total revenus/dépenses/bénéfice)."},
    },
    "price_list": {
        "ai": {"ru": "Прайс на что? Например: прайс-лист кофейни, прайс на клининговые услуги.",
               "en": "Price list for what? E.g. a coffee shop menu, cleaning services.",
               "de": "Preisliste wofür? Z. B. Café-Preisliste, Reinigungsdienstleistungen.",
               "ar": "قائمة أسعار لماذا؟ مثال: قائمة أسعار مقهى، خدمات تنظيف.",
               "zh": "价目表用于什么？例如：咖啡店价目表、清洁服务。",
               "es": "¿Lista de precios de qué? Ej.: carta de una cafetería, servicios de limpieza.",
               "fr": "Liste de prix pour quoi ? Ex. : carte d'un café, services de nettoyage."},
        "user": {"ru": "Пришли позиции прайса: наименование, единица, количество, цена, скидка (если есть) - оформлю с формулами сумм.",
                 "en": "Send price list items: name, unit, quantity, price, discount (if any) - I'll format it with sum formulas.",
                 "de": "Sende die Positionen: Bezeichnung, Einheit, Menge, Preis, Rabatt (falls vorhanden) - ich erstelle Summenformeln.",
                 "ar": "أرسل بنود القائمة: الاسم، الوحدة، الكمية، السعر، الخصم (إن وجد) - سأضيف معادلات الإجمالي.",
                 "zh": "请发送价目表条目：名称、单位、数量、单价、折扣（如有），我会加上求和公式。",
                 "es": "Envía las partidas: nombre, unidad, cantidad, precio, descuento (si lo hay) - añadiré fórmulas de suma.",
                 "fr": "Envoie les postes : nom, unité, quantité, prix, remise (le cas échéant) - j'ajouterai les formules de somme."},
    },
    "calc_table": {
        "ai": {"ru": "По какой теме нужна расчётная таблица для работы? Например: результаты опроса по теме курсовой.",
               "en": "What topic is the calculation table for? E.g. survey results for a term paper.",
               "de": "Zu welchem Thema wird die Berechnungstabelle benötigt? Z. B. Umfrageergebnisse für die Hausarbeit.",
               "ar": "لأي موضوع يحتاج الجدول الحسابي؟ مثال: نتائج استبيان لبحث دراسي.",
               "zh": "计算表用于什么主题？例如：学期论文的问卷调查结果。",
               "es": "¿Para qué tema es la tabla de cálculo? Ej.: resultados de una encuesta para un trabajo de curso.",
               "fr": "Pour quel sujet faut-il le tableau de calcul ? Ex. : résultats d'une enquête pour un dossier de cours."},
        "user": {"ru": "Пришли показатели и значения - оформлю в таблицу с формулой доли в % и итогом.",
                 "en": "Send indicators and values - I'll format a table with a % share formula and total.",
                 "de": "Sende Kennzahlen und Werte - ich erstelle eine Tabelle mit %-Anteilsformel und Summe.",
                 "ar": "أرسل المؤشرات والقيم - سأنسقها في جدول بمعادلة النسبة المئوية والإجمالي.",
                 "zh": "请发送指标及数值，我会做成带百分比公式和合计的表格。",
                 "es": "Envía los indicadores y valores - lo pondré en una tabla con fórmula de porcentaje y total.",
                 "fr": "Envoie les indicateurs et valeurs - je ferai un tableau avec formule de pourcentage et total."},
    },
}
EXCEL_KIND_HINTS = {k: {"ai": v["ai"]["ru"], "user": v["user"]["ru"]} for k, v in EXCEL_KIND_HINTS_I18N.items()}  # обратная совместимость


def excel_kind_hint(kind, mode_key, lang):
    variants = EXCEL_KIND_HINTS_I18N.get(kind, EXCEL_KIND_HINTS_I18N["expense_estimate"])[mode_key]
    return variants.get(lang, variants["ru"])


def excel_category_kb(lang="ru"):
    return ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text=WORD_CATEGORY_TITLES["physical"].get(lang, WORD_CATEGORY_TITLES["physical"]["ru"]))],
        [KeyboardButton(text=WORD_CATEGORY_TITLES["entity"].get(lang, WORD_CATEGORY_TITLES["entity"]["ru"]))],
        [KeyboardButton(text=WORD_CATEGORY_TITLES["study"].get(lang, WORD_CATEGORY_TITLES["study"]["ru"]))],
        [KeyboardButton(text=tr("btn_main_menu", lang))],
    ], resize_keyboard=True)


def excel_kind_kb(category, lang="ru"):
    cat = EXCEL_CATEGORIES.get(category, EXCEL_CATEGORIES["physical"])
    rows = [[KeyboardButton(text=EXCEL_KIND_LABELS_I18N[k].get(lang, EXCEL_KIND_LABELS_I18N[k]["ru"]))] for k in cat["kinds"]]
    rows.append([KeyboardButton(text=tr("btn_back_categories", lang))])
    rows.append([KeyboardButton(text=tr("btn_main_menu", lang))])
    return ReplyKeyboardMarkup(keyboard=rows, resize_keyboard=True)


def excel_confirm_kb(lang="ru"):
    return ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text=tr("btn_build_table", lang))],
        [KeyboardButton(text=tr("btn_add_info", lang))],
        [KeyboardButton(text=tr("btn_change_query", lang))],
        [KeyboardButton(text=tr("btn_main_menu", lang))]
    ], resize_keyboard=True)


def style_kb(include_keep=False, lang="ru"):
    rows = []
    if include_keep:
        rows.append([KeyboardButton(text=tr("btn_keep_style", lang))])
    theme_order = ["nature", "business", "tech", "school", "history", "science", "sport", "travel",
                   "food", "art", "eco", "minimal", "fashion", "default"]
    for i in range(0, len(theme_order), 2):
        pair = theme_order[i:i + 2]
        rows.append([KeyboardButton(text=THEME_LABELS_I18N[k].get(lang, THEME_LABELS_I18N[k]["ru"])) for k in pair])
    rows.append([KeyboardButton(text=tr("btn_main_menu", lang))])
    return ReplyKeyboardMarkup(keyboard=rows, resize_keyboard=True)


STYLE_BY_LABEL = {label: k for k, variants in THEME_LABELS_I18N.items() for label in variants.values()}


@dp.message(Command("start"))
async def cmd_start(m: Message, state: FSMContext):
    u = get_user(m.from_user.id)
    new_name = m.from_user.first_name or ""
    if u.get("name") != new_name:
        u["name"] = new_name
        save_users()
    await state.clear()
    if u.get("lang_chosen"):
        lang = user_lang(m.from_user.id)
        await m.answer(tr("msg_welcome", lang, name=u["name"] or "🙂"), reply_markup=main_kb(lang))
        return
    await state.set_state(Form.waiting_language)
    intro = " / ".join(TR["msg_choose_lang"][c] for c in ("ru", "en", "ar", "zh"))
    await m.answer(intro, reply_markup=lang_kb())


@dp.message(Form.waiting_language)
async def set_language(m: Message, state: FSMContext):
    code = LANG_LABEL_TO_CODE.get((m.text or "").strip())
    if not code:
        await m.answer(" / ".join(TR["msg_choose_lang"][c] for c in ("ru", "en", "ar", "zh")), reply_markup=lang_kb())
        return
    u = get_user(m.from_user.id)
    u["lang"] = code
    u["lang_chosen"] = True
    save_users()
    await state.clear()
    await m.answer(tr("msg_welcome", code, name=u.get("name") or "🙂"), reply_markup=main_kb(code))


@dp.message(F.text.in_(TR["btn_language"].values()))
async def change_language(m: Message, state: FSMContext):
    await state.set_state(Form.waiting_language)
    await m.answer(tr("msg_choose_lang", user_lang(m.from_user.id)), reply_markup=lang_kb())


@dp.message(Command("cancel"))
async def cmd_cancel(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await state.clear()
    finish_job(m.from_user.id)
    await m.answer(tr("msg_cancelled", lang), reply_markup=main_kb(lang))


@dp.message(Command("menu"))
@dp.message(F.text.in_(ALL_MAIN_MENU_LABELS))
async def to_main_menu(m: Message, state: FSMContext):
    """Кнопка/команда возврата в главное меню - работает на любом шаге любого процесса
    (презентация/документ/таблица), а не только в конце. Регистрируется здесь, до всех
    хендлеров конкретных состояний (Form.waiting_*), которые иначе перехватили бы этот
    текст как обычный пользовательский ввод (тему, данные и т.п.) - aiogram проверяет
    хендлеры в порядке регистрации и берёт первый подошедший."""
    lang = user_lang(m.from_user.id)
    await state.clear()
    finish_job(m.from_user.id)
    await m.answer(tr("msg_main_menu", lang), reply_markup=main_kb(lang))


@dp.message(F.text.in_(ALL_BTN_PRES_LABELS))
async def start_pres(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    if not can_generate(m.from_user.id):
        await m.answer(tr("msg_limit", lang))
        return
    await m.answer(tr("msg_how_build_pres", lang), reply_markup=mode_kb(lang=lang))
    await state.set_state(Form.waiting_mode)


@dp.message(Form.waiting_mode, F.text.in_(ALL_BTN_AI_GENERATE_LABELS))
async def mode_ai(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await state.update_data(mode="ai", user_text="")
    await m.answer(tr("msg_content_lang_prompt", lang), reply_markup=content_lang_kb(lang))
    await state.set_state(Form.waiting_pres_content_lang)


@dp.message(Form.waiting_pres_content_lang)
async def pres_content_lang(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    code = resolve_content_lang((m.text or "").strip(), lang)
    if not code:
        await m.answer(tr("msg_pick_content_lang", lang), reply_markup=content_lang_kb(lang))
        return
    await state.update_data(content_lang=code)
    await m.answer(tr("msg_topic_prompt", lang), reply_markup=cancel_kb(lang))
    await state.set_state(Form.waiting_topic)


@dp.message(Form.waiting_mode, F.text.in_(ALL_BTN_OWN_TEXT_LABELS))
async def mode_user(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await state.update_data(mode="user")
    await m.answer(tr("msg_own_text_prompt", lang), reply_markup=cancel_kb(lang))
    await state.set_state(Form.waiting_user_text)


@dp.message(Form.waiting_mode)
async def waiting_mode_fallback(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await m.answer(tr("msg_didnt_understand", lang), reply_markup=mode_kb(lang=lang))


@dp.message(Form.waiting_topic)
async def process_topic(m: Message, state: FSMContext):
    text = m.text or ""
    lang = user_lang(m.from_user.id)
    if len(text) > 500:
        await m.answer(tr("msg_too_long_500", lang))
        return
    name, _ = pick_theme(text)
    await state.update_data(topic=text, extra="", extra_used=0, theme_name=name)
    style_label = THEME_LABELS_I18N.get(name, {}).get(lang, THEME_LABELS.get(name, name))
    await m.answer(
        tr("msg_style_fits_topic", lang, style=style_label),
        reply_markup=style_kb(include_keep=True, lang=lang)
    )
    await state.set_state(Form.waiting_theme)


@dp.message(Form.waiting_user_text)
async def process_user_text(m: Message, state: FSMContext):
    text = m.text or ""
    lang = user_lang(m.from_user.id)
    if len(text) < 40:
        await m.answer(tr("msg_too_little_text", lang))
        return
    if len(text) > 4000:
        await m.answer(tr("msg_too_long_4000", lang))
        return
    name, _ = pick_theme(text)
    topic = text[:80].replace("\n", " ")
    await state.update_data(user_text=text, topic=topic, extra="", extra_used=0, theme_name=name)
    style_label = THEME_LABELS_I18N.get(name, {}).get(lang, THEME_LABELS.get(name, name))
    await m.answer(
        tr("msg_style_fits_text", lang, style=style_label),
        reply_markup=style_kb(include_keep=True, lang=lang)
    )
    await state.set_state(Form.waiting_theme)


@dp.message(Form.waiting_theme)
async def process_theme(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    label = m.text or ""
    if label not in TR["btn_keep_style"].values():
        name = STYLE_BY_LABEL.get(label)
        if not name:
            await m.answer(tr("msg_pick_style", lang), reply_markup=style_kb(include_keep=True, lang=lang))
            return
        await state.update_data(theme_name=name)
    await m.answer(tr("msg_how_many_slides", lang), reply_markup=slides_kb(lang))
    await state.set_state(Form.waiting_slides)


@dp.message(Form.waiting_slides)
async def process_slides(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    slides = 8
    t = m.text or ""
    if "16" in t:
        slides = 16
    elif "12" in t:
        slides = 12
    data = await state.get_data()
    angle = random.choice(ANGLES)
    await state.update_data(slides=slides, angle=angle)
    await m.answer(tr("msg_building_sample", lang), reply_markup=cancel_kb(lang))
    lang_instr = grok_lang_instruction(content_gen_lang(data, lang) or "ru")

    if data.get("mode") == "user":
        prompt = f"""Ты — редактор презентаций высокого уровня. Пользователь прислал свой текст.
Исправь ошибки, сохрани смысл и факты, но выстрой драматургию: у презентации должна быть чёткая идея,
а не пересказ по порядку.
Текст:
{data.get('user_text')}
Слайдов: {slides}
Угол подачи: {angle}
Дай короткий план текстом: цепляющее название (не общее, а конкретное и живое) и 3 пункта,
каждый — не абстрактная категория, а конкретный тезис с сутью, который будет на слайде.
Без JSON, без вводных фраз, без нумерации "Слайд 1".{lang_instr}"""
    else:
        prompt = f"""Ты — редактор презентаций высокого уровня, который умеет находить неочевидный
и интересный угол в любой теме, избегая шаблонных заголовков вроде "Введение" или "Что это такое".
Тема: {data.get('topic')}
Слайдов: {slides}
Доп: {data.get('extra')}
Угол подачи: {angle}
Дай короткий уникальный план текстом: цепляющее конкретное название и 3 пункта — каждый должен быть
содержательным тезисом (что именно будет рассказано), а не общей категорией.
Без JSON, без вводных фраз, без нумерации "Слайд 1".{lang_instr}"""

    sample = await ask_grok(prompt)
    if grok_failed(sample):
        await m.answer(
            tr("msg_grok_error", lang),
            reply_markup=slides_kb(lang)
        )
        return
    theme_name = data.get("theme_name", "default")
    style_label = THEME_LABELS_I18N.get(theme_name, {}).get(lang, THEME_LABELS.get(theme_name, theme_name))
    await state.update_data(sample=sample)
    await send_draft(
        m,
        tr("msg_draft_ready_pres", lang, sample=sample, style=style_label),
        title=data.get("topic", tr("msg_start_pres_again", lang)),
        reply_markup=confirm_kb(lang)
    )
    await state.set_state(Form.waiting_confirm)


@dp.message(Form.waiting_confirm, F.text.in_(ALL_BTN_CHANGE_TOPIC_LABELS))
async def change_topic(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    if data.get("mode") == "user":
        await m.answer(tr("msg_new_text", lang), reply_markup=cancel_kb(lang))
        await state.set_state(Form.waiting_user_text)
    else:
        await m.answer(tr("msg_new_topic", lang), reply_markup=cancel_kb(lang))
        await state.set_state(Form.waiting_topic)


@dp.message(Form.waiting_confirm, F.text.in_(ALL_BTN_CHANGE_STYLE_LABELS))
async def change_style(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await m.answer(tr("msg_which_style", lang), reply_markup=style_kb(lang=lang))
    await state.set_state(Form.waiting_style)


@dp.message(Form.waiting_style)
async def process_style(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    label = m.text or ""
    name = STYLE_BY_LABEL.get(label)
    if not name:
        await m.answer(tr("msg_pick_style", lang), reply_markup=style_kb(lang=lang))
        return
    data = await state.get_data()
    await state.update_data(theme_name=name)
    await m.answer(
        tr("msg_style_changed", lang, style=label, sample=data.get('sample')),
        reply_markup=confirm_kb(lang)
    )
    await state.set_state(Form.waiting_confirm)


@dp.message(Form.waiting_confirm, F.text.in_(ALL_BTN_ADD_INFO_LABELS))
async def add_extra(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    if (await state.get_data()).get("extra_used", 0) >= 3:
        await m.answer(tr("msg_extra_limit", lang), reply_markup=confirm_kb(lang))
        return
    await m.answer(tr("msg_extra_prompt", lang), reply_markup=cancel_kb(lang))
    await state.set_state(Form.waiting_extra)


@dp.message(Form.waiting_extra)
async def process_extra(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    text = m.text or ""
    if len(text) > 800:
        await m.answer(tr("msg_too_long_800", lang))
        return
    extra = ((data.get("extra") or "") + "\n" + text).strip()
    angle = random.choice(ANGLES)
    await state.update_data(extra=extra, extra_used=data.get("extra_used", 0) + 1, angle=angle)
    await m.answer(tr("msg_updating_draft", lang))
    lang_instr = grok_lang_instruction(content_gen_lang(data, lang) or "ru")
    if data.get("mode") == "user":
        prompt = f"""Исправь и обнови структуру.
Исходный текст:
{data.get('user_text')}
Дополнительно:
{extra}
Слайдов: {data.get('slides')}
Угол: {angle}
Короткий план: название и 3 пункта. Без JSON.{lang_instr}"""
    else:
        prompt = f"""Тема: {data.get('topic')}
Доп: {extra}
Угол: {angle}
Новый короткий план, 3 пункта. Без JSON.{lang_instr}"""
    sample = await ask_grok(prompt)
    if grok_failed(sample):
        await m.answer(
            tr("msg_draft_update_failed", lang),
            reply_markup=confirm_kb(lang)
        )
        await state.set_state(Form.waiting_confirm)
        return
    await state.update_data(sample=sample)
    await send_draft(m, tr("msg_draft_updated", lang, sample=sample), title=data.get("topic", tr("msg_start_pres_again", lang)), reply_markup=confirm_kb(lang))
    await state.set_state(Form.waiting_confirm)


@dp.message(Form.waiting_confirm, F.text.in_(ALL_BTN_FULL_VERSION_LABELS))
async def ask_photo_source(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await state.update_data(user_photos=[])
    await m.answer(
        tr("msg_photo_or_own", lang),
        reply_markup=photo_source_kb(lang)
    )
    await state.set_state(Form.waiting_pres_photo_choice)


@dp.message(Form.waiting_pres_photo_choice, F.text.in_(ALL_BTN_AI_PHOTOS_LABELS))
async def photo_source_ai(m: Message, state: FSMContext):
    await _build_presentation(m, state)


@dp.message(Form.waiting_pres_photo_choice, F.text.in_(ALL_BTN_OWN_PHOTOS_LABELS))
async def photo_source_own(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    slides = data.get("slides", 8)
    await m.answer(
        tr("msg_send_photos_one_by_one", lang, slides=slides),
        reply_markup=photos_done_kb(lang)
    )
    await state.set_state(Form.waiting_pres_photos)


@dp.message(Form.waiting_pres_photo_choice)
async def waiting_pres_photo_choice_fallback(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await m.answer(tr("msg_didnt_understand", lang), reply_markup=photo_source_kb(lang))


@dp.message(Form.waiting_pres_photos, F.photo)
async def collect_pres_photo(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    uid = m.from_user.id
    data = await state.get_data()
    # Лимит фото = число выбранных слайдов (одно фото на слайд) - больше просто некуда
    # использовать, лишние присланные фото молча игнорируются с вежливым сообщением.
    max_photos = data.get("slides", 20)
    hint_count = len(data.get("user_photos") or [])
    if hint_count >= max_photos:
        await m.answer(tr("msg_enough_photos", lang))
        return
    path = f"/tmp/{uid}_userphoto_{hint_count}_{random.randint(1000, 9999)}.jpg"
    try:
        file = await bot.get_file(m.photo[-1].file_id)
        await bot.download_file(file.file_path, path)
    except Exception as e:
        print("Photo download error:", e)
        await m.answer(tr("msg_photo_download_fail", lang))
        return
    # Скачивание выше идёт БЕЗ блокировки (можно параллельно, это самая долгая часть);
    # а вот изменение общего списка user_photos - под блокировкой на пользователя, иначе
    # несколько фото одного альбома, скачавшихся почти одновременно, затрут друг друга
    # в state (см. _get_photo_lock).
    async with _get_photo_lock(uid):
        data = await state.get_data()
        photos = data.get("user_photos") or []
        max_photos = data.get("slides", 20)
        if len(photos) >= max_photos:
            await m.answer(tr("msg_enough_photos", lang))
            return
        photos.append(path)
        await state.update_data(user_photos=photos)
        n = len(photos)
    await m.answer(tr("msg_photo_received", lang, n=n))


@dp.message(Form.waiting_pres_photos, F.text.in_(ALL_BTN_PHOTOS_DONE_LABELS))
async def finish_pres_photos(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    if not (data.get("user_photos") or []):
        await m.answer(tr("msg_photo_wait", lang), reply_markup=photos_done_kb(lang))
        return
    await _build_presentation(m, state)


@dp.message(Form.waiting_pres_photos)
async def pres_photos_fallback(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await m.answer(tr("msg_send_photos", lang), reply_markup=photos_done_kb(lang))


async def _build_presentation(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    uid = m.from_user.id
    u = get_user(uid)
    ok, reason = start_job(uid)
    if not ok:
        await m.answer(reason)
        return
    await m.answer(tr("msg_building_pres_photos", lang))
    try:

        theme_name = data.get("theme_name") or pick_theme(data.get("topic", ""))[0]
        colors = THEMES.get(theme_name, THEMES["default"])
        angle = data.get("angle") or random.choice(ANGLES)
        cgl = content_gen_lang(data, lang)

        if data.get("mode") == "user":
            common_rules = """
    Правила качества (обязательны для каждого слайда):
    - Заголовок слайда: 3-6 слов, конкретный и живой, называет суть тезиса, а не общую категорию
      (плохо: "Введение", "Общая информация"; хорошо: конкретный факт или вывод).
    - Текст: ровно 2 абзаца, каждый 2-4 предложения. Первый абзац - конкретный факт, пример или наблюдение,
      второй - что это значит или почему это важно. Никаких вводных фраз типа "в этом слайде" или "стоит отметить".
    - Между слайдами не повторяй одну и ту же структуру фразы - разнообразь подачу (факт, вопрос, сравнение, история).
    - Никаких общих фраз без содержания ("это важная тема", "мир меняется") - только конкретика: цифры, имена,
      примеры, детали, если они есть в материале или логично следуют из темы.
    - Фото: описание живого кадра (реальная сцена, человек, объект, место), не стоковый шаблон и не абстракция.
    - График (ключ "chart", НЕОБЯЗАТЕЛЬНЫЙ): добавляй его ТОЛЬКО если в тексте пользователя
      реально ЕСТЬ конкретные цифры, статистика, доли или динамика по времени, которые можно
      честно визуализировать - НИКОГДА не придумывай цифры от себя ради красивой картинки.
      Если в присланном тексте нет настоящих числовых данных - не добавляй chart вообще, у слайда
      остаётся только фото. Не больше 1-2 слайдов с графиком на всю презентацию, и только там, где
      график реально помогает понять данные быстрее, чем абзац текста.
      Формат: "chart": {{"chart_type": "bar" | "line" | "pie", "title": "короткий заголовок графика",
      "categories": ["...", "..."], "series": [{{"name": "...", "values": [12, 34, ...]}}]}} -
      цифры бери СТРОГО из текста пользователя, ничего не домысливая и не округляя "для красоты".
      У слайда с графиком image_prompt всё равно укажи (на случай, если график не соберётся), но
      использоваться будет либо график, либо фото, не оба сразу."""
            raw = await ask_grok(f"""Собери уникальную презентацию из текста пользователя.
    Исправь ошибки, сохрани смысл и все факты из текста.
    Текст:
    {data.get('user_text')}
    Доп:
    {data.get('extra')}
    Слайдов: {data.get('slides')}
    Угол: {angle}
    Стиль: {theme_name}
    {common_rules}
    Только JSON (ключ "chart" добавляй лишь на 1-2 слайдах и только если в тексте пользователя
    есть настоящие числа для него, иначе не пиши его вовсе):
    {{"title":"...","slides":[{{"title":"...","content":"абзац1\n\nабзац2","image_prompt":"unique cinematic scene","chart":null}}]}}{grok_json_lang_instruction(cgl or "ru")}""")
        else:
            common_rules = """
    Правила качества (обязательны для каждого слайда):
    - Заголовок слайда: 3-6 слов, конкретный и живой, называет суть тезиса, а не общую категорию
      (плохо: "Введение", "Общая информация"; хорошо: конкретный факт или вывод).
    - Текст: ровно 2 абзаца, каждый 2-4 предложения. Первый абзац - конкретный факт, пример или наблюдение,
      второй - что это значит или почему это важно. Никаких вводных фраз типа "в этом слайде" или "стоит отметить".
    - Между слайдами не повторяй одну и ту же структуру фразы - разнообразь подачу (факт, вопрос, сравнение, история).
    - Никаких общих фраз без содержания ("это важная тема", "мир меняется") - только конкретика: цифры, имена,
      примеры, детали, логично следующие из темы.
    - Фото: описание живого кадра (реальная сцена, человек, объект, место), не стоковый шаблон и не абстракция.
    - Графики (диаграммы) в этом режиме НЕ используем: тема задана вами без исходного текста, а значит любые
      цифры для графика пришлось бы выдумать - это выглядело бы как настоящие данные, но ими не являлось бы.
      Ключ "chart" всегда оставляй null."""
            raw = await ask_grok(f"""Собери уникальную презентацию уровня лучшего журнала на эту тему,
    найди неочевидный и интересный угол, избегай банальностей.
    Тема: {data.get('topic')}
    Слайдов: {data.get('slides')}
    Доп: {data.get('extra')}
    Угол: {angle}
    Стиль: {theme_name}
    {common_rules}
    Только JSON (ключ "chart" всегда null в этом режиме):
    {{"title":"...","slides":[{{"title":"...","content":"абзац1\n\nабзац2","image_prompt":"unique cinematic scene","chart":null}}]}}{grok_json_lang_instruction(cgl or "ru")}""")

        try:
            content = extract_json(raw)
            if not isinstance(content.get("slides"), list) or not content["slides"]:
                raise ValueError("В ответе модели нет слайдов")
        except Exception as e:
            print("Presentation JSON parse error:", e)
            await m.answer(tr("msg_no_text", lang, btn=tr("msg_start_pres_again", lang)), reply_markup=main_kb(lang))
            await state.clear()
            return

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slides_data = content.get("slides", [])
        n = len(slides_data)

        # 7 раскладок (3 базовых + 3 зеркальных фото слева/справа/сверху/снизу/крупно/мелко,
        # плюс "карточки" - две текстовые плашки рядом без фото, для слайдов с двумя явными
        # смысловыми блоками в content) - выбираются случайно без повтора одной и той же
        # раскладки два слайда подряд, чтобы презентация не выглядела как один и тот же шаблон,
        # повторённый N раз, и чтобы разные презентации не были визуально неотличимы друг от друга.
        LAYOUT_COUNT = 7
        layout_sequence = []
        prev_layout = None
        for i in range(n):
            pool = [l for l in range(LAYOUT_COUNT) if l != prev_layout]
            content_blocks_i = [x.strip() for x in (slides_data[i].get("content") or "").split("\n\n") if x.strip()]
            if len(content_blocks_i) < 2:
                # "Карточкам" нужно ровно 2 смысловых блока - если модель не разбила
                # content на два абзаца, одна из карточек останется пустой и будет выглядеть
                # криво, поэтому для таких слайдов этот вариант просто не предлагаем.
                pool = [l for l in pool if l != 6]
            next_layout = random.choice(pool)
            layout_sequence.append(next_layout)
            prev_layout = next_layout

        # Свои фото пользователя (если он их прислал) идут в очередь: первое - на обложку,
        # остальные - по слайдам по порядку; на то, что не хватило, генерируем через ИИ как раньше.
        user_photos = list(data.get("user_photos") or [])
        # Общий уровень яркости по ВСЕМ присланным фото - ориентир для enhance_user_photo(),
        # чтобы фото с разных камер/освещения в одной презентации смотрелись как единый комплект.
        user_photos_luminance = average_luminance(user_photos) if user_photos else None

        cover_img = None
        cover_panel_img = None
        cover_src = f"/tmp/{uid}_cover.png"
        cover_own = user_photos.pop(0) if user_photos else None
        user_photo_originals = []
        if cover_own:
            user_photo_originals.append(cover_own)
            cover_own = prepare_user_photo(cover_own, colors, user_photos_luminance)
        cover_ok = bool(cover_own) or await generate_image(
            f"{content.get('title')}, wide cinematic opening shot, {colors['photo']}", cover_src
        )
        # Два стиля обложки, выбираются случайно - чтобы первый (самый запоминающийся) слайд
        # тоже не был всегда одинаковым в каждой презентации:
        # "band" - фото на весь кадр, плашка с заголовком снизу (классический вариант);
        # "split" - фото на половину кадра (слева или справа), заголовок на цветной панели рядом.
        cover_style = random.choice(["band", "split"]) if cover_ok else "band"
        cover_split_side = random.choice(["left", "right"])
        if cover_ok:
            if cover_style == "split":
                cover_panel = f"/tmp/{uid}_cover_panel.png"
                cover(cover_own or cover_src, cover_panel, 1280, 1500)
                cover_panel_img = cover_panel
            else:
                cover_wide = f"/tmp/{uid}_cover_w.png"
                cover(cover_own or cover_src, cover_wide, 1920, 1080)
                cover_img = cover_wide

        # Если модель дала валидный chart для слайда - фото для него не генерируем вообще
        # (график займёт то же место в раскладке, экономим время и токены на генерацию картинки).
        MAX_CHARTS = 2
        charts = []
        chart_count = 0
        for s in slides_data:
            c = s.get("chart")
            valid = (
                isinstance(c, dict) and c.get("categories") and c.get("series")
                and chart_count < MAX_CHARTS
            )
            if valid:
                chart_count += 1
                charts.append(c)
            else:
                charts.append(None)

        images = []
        raw_sources = []
        for i, s in enumerate(slides_data):
            if charts[i] or layout_sequence[i] == 6:
                # Раскладке "карточки" фото не нужно вовсе - не тратим на неё ни платную
                # генерацию через ИИ, ни личное фото пользователя (пусть достанется слайду,
                # который реально его покажет).
                images.append(None)
                continue
            own = user_photos.pop(0) if user_photos else None
            if own:
                user_photo_originals.append(own)
                src, ok = prepare_user_photo(own, colors, user_photos_luminance), True
            else:
                src = f"/tmp/{uid}_{i}_{random.randint(1000, 9999)}.png"
                prompt = f"{s.get('image_prompt') or s.get('title')}, {colors['photo']}, unique composition"
                ok = await generate_image(prompt, src)
            if ok:
                raw_sources.append(src)
                wide, tall = f"/tmp/{uid}_{i}_{random.randint(1000, 9999)}_w.png", f"/tmp/{uid}_{i}_{random.randint(1000, 9999)}_t.png"
                cover(src, wide, 1920, 1080)
                cover(src, tall, 1260, 1500)
                images.append((wide, tall))
            else:
                images.append(None)

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_background(slide, colors)
        if cover_panel_img:
            photo_x = 0 if cover_split_side == "left" else 6.933
            panel_x = 6.933 if cover_split_side == "left" else 0
            slide.shapes.add_picture(cover_panel_img, Inches(photo_x), Inches(0), width=Inches(6.4), height=Inches(7.5))
            rect(slide, panel_x, 0, 6.4, 7.5, colors["bg"])
            txt(slide, panel_x + 0.5, 2.9, 5.4, 2.0, content.get("title", "Презентация"), 32, colors["ink"], True, font_name=colors["heading_font"])
            rect(slide, panel_x + 0.5, 4.9, 0.85, 0.05, colors["line"])
            txt(slide, panel_x + 0.5, 6.6, 5.4, 0.4, "01  /  введение", 13, colors["mute"])
        else:
            if cover_img:
                slide.shapes.add_picture(cover_img, Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
                rect(slide, 0, 4.7, 13.333, 2.8, colors["bg"])
            txt(slide, 0.7, 5.0, 12, 1.5, content.get("title", "Презентация"), 40, colors["ink"], True, font_name=colors["heading_font"])
            txt(slide, 0.7, 6.6, 12, 0.4, "01  /  введение", 13, colors["mute"])

        for idx, s in enumerate(slides_data):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide_background(slide, colors)
            layout = layout_sequence[idx]
            img = images[idx] if idx < len(images) else None
            chart_data = charts[idx] if idx < len(charts) else None
            if layout == 0:
                if chart_data:
                    add_chart(slide, 0.5, 0.6, 5.4, 6.3, chart_data, colors)
                elif img:
                    slide.shapes.add_picture(img[1], Inches(0), Inches(0), width=Inches(6.4), height=Inches(7.5))
                txt(slide, 7.05, 1.5, 5.5, 1.6, s.get("title", ""), 30, colors["ink"], True, font_name=colors["heading_font"])
                rect(slide, 7.05, 3.25, 0.85, 0.05, colors["line"])
                box = slide.shapes.add_textbox(Inches(7.05), Inches(3.5), Inches(5.5), Inches(3.2))
                tf = box.text_frame
                tf.word_wrap = True
                blocks = [x.strip() for x in (s.get("content") or "").split("\n") if x.strip()][:2]
                for i, b in enumerate(blocks or [""]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = b
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(*colors["mid"])
                    p.font.name = "Calibri"
                    p.space_after = Pt(16)
                txt(slide, 7.05, 6.95, 5.5, 0.3, f"{idx + 2:02}  /  {n + 1:02}", 12, colors["mute"])
            elif layout == 1:
                if chart_data:
                    add_chart(slide, 0.7, 0.35, 11.9, 4.0, chart_data, colors)
                elif img:
                    slide.shapes.add_picture(img[0], Inches(0), Inches(0), width=Inches(13.333), height=Inches(4.55))
                txt(slide, 0.7, 4.85, 12, 1.0, s.get("title", ""), 28, colors["ink"], True, font_name=colors["heading_font"])
                box = slide.shapes.add_textbox(Inches(0.7), Inches(5.85), Inches(12), Inches(1.2))
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = " ".join((s.get("content") or "").split())
                p.font.size = Pt(15)
                p.font.color.rgb = RGBColor(*colors["mid"])
                p.font.name = "Calibri"
            elif layout == 2:
                # Заголовок уже НЕ должен доходить до x=8.7 (где начинается фото/график
                # справа) - раньше ширина плашки (8.2") залезала на 0.2" в зону картинки,
                # и длинные заголовки визуально обрезались, т.к. фото рисуется поверх текста.
                txt(slide, 0.7, 1.3, 7.6, 2.2, s.get("title", ""), 36, colors["ink"], True, font_name=colors["heading_font"])
                rect(slide, 0.7, 3.6, 1.1, 0.06, colors["line"])
                box = slide.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(7.4), Inches(2.6))
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = " ".join((s.get("content") or "").split())
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(*colors["mid"])
                p.font.name = "Calibri"
                if chart_data:
                    add_chart(slide, 8.5, 1.3, 4.6, 4.7, chart_data, colors)
                elif img:
                    slide.shapes.add_picture(img[1], Inches(8.7), Inches(1.3), width=Inches(3.9), height=Inches(4.7))
                txt(slide, 0.7, 6.95, 5.5, 0.3, f"{idx + 2:02}  /  {n + 1:02}", 12, colors["mute"])
            elif layout == 3:
                # Зеркало layout 0: фото - правая половина кадра, текст - слева.
                if chart_data:
                    add_chart(slide, 0.7, 0.6, 5.4, 6.3, chart_data, colors)
                elif img:
                    slide.shapes.add_picture(img[1], Inches(6.933), Inches(0), width=Inches(6.4), height=Inches(7.5))
                txt(slide, 0.7, 1.5, 5.5, 1.6, s.get("title", ""), 30, colors["ink"], True, font_name=colors["heading_font"])
                rect(slide, 0.7, 3.25, 0.85, 0.05, colors["line"])
                box = slide.shapes.add_textbox(Inches(0.7), Inches(3.5), Inches(5.5), Inches(3.2))
                tf = box.text_frame
                tf.word_wrap = True
                blocks = [x.strip() for x in (s.get("content") or "").split("\n") if x.strip()][:2]
                for i, b in enumerate(blocks or [""]):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = b
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(*colors["mid"])
                    p.font.name = "Calibri"
                    p.space_after = Pt(16)
                txt(slide, 0.7, 6.95, 5.5, 0.3, f"{idx + 2:02}  /  {n + 1:02}", 12, colors["mute"])
            elif layout == 4:
                # Зеркало layout 1: фото - нижняя полоса кадра, текст - сверху.
                txt(slide, 0.7, 0.4, 12, 1.0, s.get("title", ""), 28, colors["ink"], True, font_name=colors["heading_font"])
                box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12), Inches(1.2))
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = " ".join((s.get("content") or "").split())
                p.font.size = Pt(15)
                p.font.color.rgb = RGBColor(*colors["mid"])
                p.font.name = "Calibri"
                if chart_data:
                    add_chart(slide, 0.7, 3.15, 11.9, 4.0, chart_data, colors)
                elif img:
                    slide.shapes.add_picture(img[0], Inches(0), Inches(2.95), width=Inches(13.333), height=Inches(4.55))
            elif layout == 5:
                # Зеркало layout 2: фото - слева мелко, текст - справа.
                txt(slide, 5.2, 1.3, 7.4, 2.2, s.get("title", ""), 36, colors["ink"], True, font_name=colors["heading_font"])
                rect(slide, 5.2, 3.6, 1.1, 0.06, colors["line"])
                box = slide.shapes.add_textbox(Inches(5.2), Inches(3.9), Inches(7.4), Inches(2.6))
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = " ".join((s.get("content") or "").split())
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(*colors["mid"])
                p.font.name = "Calibri"
                if chart_data:
                    add_chart(slide, 0.7, 1.3, 4.6, 4.7, chart_data, colors)
                elif img:
                    slide.shapes.add_picture(img[1], Inches(0.7), Inches(1.3), width=Inches(3.9), height=Inches(4.7))
                txt(slide, 5.2, 6.95, 7.4, 0.3, f"{idx + 2:02}  /  {n + 1:02}", 12, colors["mute"])
            else:
                # Раскладка "карточки" (layout 6) - без фото: заголовок сверху, ниже -
                # две текстовые плашки рядом с мягкой заливкой под тему. Используется только
                # для слайдов, где content уже разбит на 2 абзаца (см. фильтр при выборе
                # раскладки выше) - каждый абзац идёт в свою карточку.
                txt(slide, 0.7, 0.5, 11.9, 1.0, s.get("title", ""), 32, colors["ink"], True, font_name=colors["heading_font"])
                rect(slide, 0.7, 1.55, 1.1, 0.06, colors["line"])
                blocks = [x.strip() for x in (s.get("content") or "").split("\n\n") if x.strip()][:2]
                card_w, gap, card_top, card_h = 5.85, 0.3, 2.0, 4.65
                card_fill = _pptx_tint(colors["line"], 0.9)
                for i, block in enumerate(blocks):
                    card_l = 0.7 + i * (card_w + gap)
                    rect(slide, card_l, card_top, card_w, card_h, card_fill)
                    box = slide.shapes.add_textbox(Inches(card_l + 0.35), Inches(card_top + 0.3),
                                                    Inches(card_w - 0.7), Inches(card_h - 0.6))
                    tf = box.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = block
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(*colors["mid"])
                    p.font.name = "Calibri"
                txt(slide, 0.7, 6.95, 5.5, 0.3, f"{idx + 2:02}  /  {n + 1:02}", 12, colors["mute"])

        pptx_path = f"/tmp/pres_{uid}.pptx"
        prs.save(pptx_path)
        pdf_path = f"/tmp/pres_{uid}.pdf"
        pdf = canvas.Canvas(pdf_path, pagesize=A4)
        w, h = A4
        fn, fb = register_pdf_fonts()

        # Титульная "обложка" копии — как в самой презентации.
        pdf.setFont(fb, 22)
        for line in wrap_lines(content.get("title", "Презентация"), fb, 22, w - 80):
            pdf.drawString(40, h - 90, line)
            break  # заголовок короткий по промпту, одной строки почти всегда достаточно
        pdf.setFont(fn, 11)
        pdf.drawString(40, h - 115, "Текстовая копия презентации")
        pdf.showPage()

        y = h - 60
        for i, s in enumerate(slides_data, 1):
            if y < 100:
                pdf.showPage()
                y = h - 60
            pdf.setFont(fb, 14)
            for line in wrap_lines(f"{i}. {s.get('title', '')}", fb, 14, w - 80):
                if y < 60:
                    pdf.showPage()
                    y = h - 60
                    pdf.setFont(fb, 14)
                pdf.drawString(40, y, line)
                y -= 18
            y -= 6
            pdf.setFont(fn, 11)
            paragraphs = [x.strip() for x in (s.get("content") or "").split("\n") if x.strip()]
            for para in paragraphs:
                for line in wrap_lines(para, fn, 11, w - 80):
                    if y < 60:
                        pdf.showPage()
                        y = h - 60
                        pdf.setFont(fn, 11)
                    pdf.drawString(40, y, line)
                    y -= 15
                y -= 8
            y -= 14
        pdf.save()

        pres_fname = safe_filename(content.get("title"), fallback="Презентация")
        await m.answer_document(FSInputFile(pptx_path, filename=f"{pres_fname}.pptx"), caption=tr("msg_pptx_caption", lang))
        await m.answer_document(FSInputFile(pdf_path, filename=f"{pres_fname}.pdf"), caption=tr("msg_pptx_pdf_caption", lang))
        u["generations"] += 1
        u["history"].append(f"{datetime.now().strftime('%d.%m %H:%M')} — {content.get('title')}")
        note_success(uid)
        for p in (cover_src, cover_own, cover_img, cover_panel_img, pptx_path, pdf_path, *raw_sources, *user_photo_originals, *[f for pair in images if pair for f in pair]):
            try:
                if p and os.path.exists(p):
                    os.remove(p)
            except Exception:
                pass
        await m.answer(
            tr("msg_pptx_ready", lang) + "\n\n" + await signature_line(lang),
            reply_markup=main_kb(lang)
        )
        await state.clear()
    finally:
        finish_job(uid)


@dp.message(Form.waiting_confirm)
async def waiting_confirm_fallback(m: Message, state: FSMContext):
    """Ловит любой текст, не совпавший с кнопками выше, чтобы диалог никогда
    не зависал без ответа."""
    lang = user_lang(m.from_user.id)
    await m.answer(
        tr("msg_didnt_understand", lang),
        reply_markup=confirm_kb(lang)
    )


def _run(p, text, size=12, bold=False, name="Times New Roman", align=None):
    if align is not None:
        p.alignment = align
    r = p.add_run(text or "")
    r.bold = bold
    r.font.size = DocxPt(size)
    r.font.name = name
    r.font.color.rgb = DocxRGB(0, 0, 0)
    return r


def _p(doc, text, size=12, bold=False, align=WD_ALIGN_PARAGRAPH.LEFT, before=0, after=6, first=None, line=1.15):
    p = doc.add_paragraph()
    p.alignment = align
    p.paragraph_format.space_before = DocxPt(before)
    p.paragraph_format.space_after = DocxPt(after)
    p.paragraph_format.line_spacing = line
    if first is not None:
        p.paragraph_format.first_line_indent = Cm(first)
    _run(p, text, size, bold)
    return p


def _tabbed_line(doc, parts, tabs_cm=(9, 13), size=12, before=0, after=0):
    """Строка с элементами по табуляции - для строк вида 'подпись ___ И.О. Фамилия'.
    parts: список (текст, bold)."""
    p = doc.add_paragraph()
    p.paragraph_format.space_before = DocxPt(before)
    p.paragraph_format.space_after = DocxPt(after)
    for cm in tabs_cm:
        p.paragraph_format.tab_stops.add_tab_stop(Cm(cm))
    for i, (text, bold) in enumerate(parts):
        if i > 0:
            p.add_run("\t")
        _run(p, text, size, bold)
    return p


def _student_title_page(doc, kind_label, title, meta):
    """Титульный лист по образцу техникума/ссуза: жирные лейблы ВЫПОЛНИЛ/ПРОВЕРИЛ/ОЦЕНКА
    слева, линия подписи с ФИО справа от неё, мелкие подписи (подпись)/(ФИО) под линией."""
    _p(doc, meta.get("org") or "Министерство образования", 13, align=WD_ALIGN_PARAGRAPH.CENTER, after=0)
    _p(doc, meta.get("school") or "[Название учебного заведения]", 13, align=WD_ALIGN_PARAGRAPH.CENTER)
    for _ in range(5):
        doc.add_paragraph()
    _p(doc, title, 16, True, WD_ALIGN_PARAGRAPH.CENTER, after=10)
    _p(doc, kind_label, 14, True, WD_ALIGN_PARAGRAPH.CENTER, after=10)
    if meta.get("discipline"):
        _p(doc, f"по предмету/дисциплине: {meta.get('discipline')}", 12, True, WD_ALIGN_PARAGRAPH.CENTER)
    for _ in range(5):
        doc.add_paragraph()
    _p(doc, "ВЫПОЛНИЛ", 12, True, after=0)
    _tabbed_line(
        doc,
        [(f"студент {meta.get('group') or '[группа]'} группы", False), ("__________", False), (meta.get("author") or "И.О. Фамилия", True)],
        after=0,
    )
    _tabbed_line(doc, [("", False), ("(подпись)", False), ("(ФИО)", False)], size=9, after=8)
    _p(doc, "«___»____________ 20___ г.", 12, before=0, after=16)
    _p(doc, "ПРОВЕРИЛ", 12, True, after=0)
    _tabbed_line(
        doc,
        [(meta.get("teacher_position") or "", False), ("__________", False), (meta.get("teacher") or "И.О. Фамилия", True)],
        after=0,
    )
    _tabbed_line(doc, [("", False), ("(подпись)", False), ("(ФИО)", False)], size=9, after=8)
    doc.add_paragraph()
    _p(doc, "ОЦЕНКА", 12, True, after=0)
    _tabbed_line(doc, [("", False), ("__________________________", False)], after=16)
    _p(doc, "«___»____________ 20___ г.", 12, after=30)
    _p(doc, f"{meta.get('city') or '[Город]'} {meta.get('year') or '2026'}", 12, align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_page_break()


def _set_cell_borders(cell):
    """Тонкая чёрная рамка у ячейки таблицы - python-docx не даёт это напрямую, только через XML."""
    tcPr = cell._tc.get_or_add_tcPr()
    borders = OxmlElement("w:tcBorders")
    for edge in ("top", "left", "bottom", "right"):
        el = OxmlElement(f"w:{edge}")
        el.set(qn("w:val"), "single")
        el.set(qn("w:sz"), "4")
        el.set(qn("w:color"), "000000")
        borders.append(el)
    tcPr.append(borders)


def _add_table(doc, table_data, size=11):
    """table_data: {"headers": [...], "rows": [[...], ...]}. Рисует таблицу с рамками
    и жирной шапкой - для перечней товаров/имущества в договорах и актах, статистики
    и практических данных в курсовых/рефератах."""
    headers = table_data.get("headers") or []
    rows = table_data.get("rows") or []
    if not headers and not rows:
        return
    ncols = len(headers) if headers else max((len(r) for r in rows), default=0)
    if ncols == 0:
        return
    t = doc.add_table(rows=0, cols=ncols)
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    if headers:
        row = t.add_row()
        for i, htext in enumerate(headers):
            cell = row.cells[i]
            cell.text = ""
            _run(cell.paragraphs[0], str(htext), size, True)
            _set_cell_borders(cell)
    for r in rows:
        row = t.add_row()
        for i in range(ncols):
            cell = row.cells[i]
            cell.text = ""
            val = str(r[i]) if i < len(r) else ""
            _run(cell.paragraphs[0], val, size, False)
            _set_cell_borders(cell)
    doc.add_paragraph()  # отступ после таблицы


def _parse_markdown_table(lines):
    """Модель иногда пишет таблицу markdown-синтаксисом прямо в тексте раздела
    (| ячейка | ячейка |) вместо структурированного ключа "table" - на такой случай
    разбираем это защитным парсером и рендерим как настоящую таблицу Word.
    Возвращает {"headers": [...], "rows": [[...], ...]} или None, если это не таблица."""
    rows = []
    for ln in lines:
        ln = ln.strip()
        if not (ln.startswith("|") and ln.endswith("|") and ln.count("|") >= 3):
            return None
        rows.append([c.strip() for c in ln.strip("|").split("|")])
    if len(rows) < 2:
        return None
    sep = rows[1]
    if sep and all(set(c) <= set("-: ") for c in sep):
        rows.pop(1)
    if not rows:
        return None
    return {"headers": rows[0], "rows": rows[1:]}


def _body(doc, sections, indent=True, head_center=False, size=12, line=1.15, first=1.25, head_size=13):
    for block in sections or []:
        head = (block.get("title") or "").strip()
        body = block.get("content") or ""
        if head:
            _p(
                doc, head, head_size, True,
                align=WD_ALIGN_PARAGRAPH.CENTER if head_center else WD_ALIGN_PARAGRAPH.LEFT,
                before=10, after=6
            )
        raw_paras = [x.strip() for x in body.split("\n") if x.strip()]
        i = 0
        while i < len(raw_paras):
            if raw_paras[i].startswith("|") and raw_paras[i].endswith("|"):
                j = i
                while j < len(raw_paras) and raw_paras[j].startswith("|") and raw_paras[j].endswith("|"):
                    j += 1
                md_table = _parse_markdown_table(raw_paras[i:j])
                if md_table:
                    _add_table(doc, md_table, size=size - 1 if size > 9 else size)
                    i = j
                    continue
            _p(doc, raw_paras[i], size, first=first if indent else None, after=6, line=line)
            i += 1
        table_data = block.get("table")
        if isinstance(table_data, dict):
            _add_table(doc, table_data, size=size - 1 if size > 9 else size)


# Семейства новых юридических документов: чтобы не дублировать почти одинаковый
# код форматирования 22 раза, новые виды сгруппированы по "семьям" (двусторонний
# договор / акт / заявление в суд-организацию / односторонний документ), и
# build_word() рендерит их по общему шаблону, подставляя только заголовок и роли сторон.
CONTRACT_FAMILY = {
    "marriage_contract": {"caption": "БРАЧНЫЙ ДОГОВОР", "left": "СУПРУГ", "right": "СУПРУГА"},
    "gift": {"caption": "ДОГОВОР ДАРЕНИЯ", "left": "ДАРИТЕЛЬ", "right": "ОДАРЯЕМЫЙ"},
    "supply": {"caption": "ДОГОВОР ПОСТАВКИ", "left": "ПОСТАВЩИК", "right": "ПОКУПАТЕЛЬ"},
    "agency": {"caption": "АГЕНТСКИЙ ДОГОВОР", "left": "ПРИНЦИПАЛ", "right": "АГЕНТ"},
    "joint_activity": {"caption": "ДОГОВОР О СОВМЕСТНОЙ ДЕЯТЕЛЬНОСТИ", "left": "СТОРОНА 1", "right": "СТОРОНА 2"},
    "nonresidential_rent": {"caption": "ДОГОВОР АРЕНДЫ НЕЖИЛОГО ПОМЕЩЕНИЯ", "left": "АРЕНДОДАТЕЛЬ", "right": "АРЕНДАТОР"},
    "cession": {"caption": "ДОГОВОР ЦЕССИИ (УСТУПКИ ПРАВА ТРЕБОВАНИЯ)", "left": "ЦЕДЕНТ", "right": "ЦЕССИОНАРИЙ"},
    "nda": {"caption": "СОГЛАШЕНИЕ О НЕРАЗГЛАШЕНИИ (NDA)", "left": "СТОРОНА 1", "right": "СТОРОНА 2"},
    "employment": {"caption": "ТРУДОВОЙ ДОГОВОР", "left": "РАБОТОДАТЕЛЬ", "right": "РАБОТНИК"},
    "self_employed": {"caption": "ДОГОВОР С САМОЗАНЯТЫМ", "left": "ЗАКАЗЧИК", "right": "ИСПОЛНИТЕЛЬ"},
    "services": {"caption": "ДОГОВОР ОКАЗАНИЯ УСЛУГ", "left": "ЗАКАЗЧИК", "right": "ИСПОЛНИТЕЛЬ"},
}

ACT_FAMILY = {
    "work_act": {"caption": "АКТ ВЫПОЛНЕННЫХ РАБОТ (ОКАЗАННЫХ УСЛУГ)"},
}

STATEMENT_FAMILY = {
    "lawsuit": {
        "caption": "ИСКОВОЕ ЗАЯВЛЕНИЕ",
        "closing": "На основании изложенного и в соответствии с законодательством РФ прошу суд удовлетворить исковые требования.",
    },
}

DECLARATION_FAMILY = {
    "alimony": {
        "caption": "СОГЛАШЕНИЕ ОБ УПЛАТЕ АЛИМЕНТОВ",
        "intro": "Я, {author}, обязуюсь уплачивать алименты на условиях, указанных ниже:",
    },
    "warranty_letter": {
        "caption": "ГАРАНТИЙНОЕ ПИСЬМО",
        "intro": "Настоящим письмом {author} гарантирует следующее:",
    },
}

# Единая схема мета-полей под каждый тип документа. Используется и в промпте для ИИ
# (как образец JSON, который нужно заполнить), и напрямую как список плейсхолдеров
# для режима "Скачать шаблон" (без обращения к ИИ).
META_SCHEMAS = {
    "referat": '{"org":"Министерство образования","school":"[учебное заведение]","discipline":"[предмет/дисциплина/МДК]","group":"[номер группы]","author":"[ФИО студента]","teacher":"[должность и/или ФИО преподавателя]","city":"[город]","year":"2026"}',
    "coursework": '{"org":"Министерство образования","school":"[учебное заведение]","specialty":"[специальность]","author":"[ФИО студента]","teacher":"[ФИО научного руководителя]","city":"[город]","year":"2026"}',
    "report": '{"school":"[учебное заведение]","discipline":"[предмет/дисциплина]","group":"[номер группы]","author":"[ФИО]","teacher":"[должность и/или ФИО преподавателя]"}',
    "essay": '{"school":"[учебное заведение]","discipline":"[предмет/дисциплина]","group":"[номер группы]","author":"[ФИО]","teacher":"[должность и/или ФИО преподавателя]"}',
    "notes": '{}',
    "dkp": '{"city":"[город]","date":"«___» __________ 20___ г."}',
    "rent": '{"city":"[город]","date":"«___» __________ 20___ г."}',
    "offer": '{"number":"[номер КП]"}',
    "act": '{"city":"[город]","date":"«___» __________ 20___ г.","basis":"[договор №/основание]","from_party":"[передающая сторона, ФИО/наименование]","to_party":"[принимающая сторона, ФИО/наименование]"}',
    "statement": '{"to":"[должность руководителя]","to_name":"[ФИО руководителя]","from":"[должность, ФИО заявителя]","date":"[дата]","author":"[ФИО]"}',
    "proxy": '{"city":"[город]","date":"«___» __________ 20___ г.","author":"[ФИО доверителя]"}',
    "loan": '{"city":"[город]","date":"«___» __________ 20___ г.","from_party":"[заимодавец, ФИО/паспорт]","to_party":"[заёмщик, ФИО/паспорт]"}',
    "claim": '{"to":"[адресат претензии]","from":"[заявитель, ФИО/контакты]","date":"[дата]","basis":"[договор №/основание]"}',
    "consent": '{"city":"[город]","date":"«___» __________ 20___ г.","author":"[ФИО родителя]"}',
    "marriage_contract": '{"city":"[город]","date":"«___» __________ 20___ г."}',
    "gift": '{"city":"[город]","date":"«___» __________ 20___ г."}',
    "lawsuit": '{"to":"[наименование суда]","from":"[ФИО истца, адрес, контакты]","date":"[дата]"}',
    "alimony": '{"city":"[город]","date":"«___» __________ 20___ г.","author":"[ФИО плательщика]"}',
    "services": '{"city":"[город]","date":"«___» __________ 20___ г."}',
    "employment": '{"city":"[город]","date":"«___» __________ 20___ г."}',
    "work_act": '{"city":"[город]","date":"«___» __________ 20___ г.","basis":"[договор №/основание]","from_party":"[исполнитель]","to_party":"[заказчик]"}',
    "supply": '{"city":"[город]","date":"«___» __________ 20___ г."}',
    "agency": '{"city":"[город]","date":"«___» __________ 20___ г."}',
    "joint_activity": '{"city":"[город]","date":"«___» __________ 20___ г."}',
    "nonresidential_rent": '{"city":"[город]","date":"«___» __________ 20___ г."}',
    "cession": '{"city":"[город]","date":"«___» __________ 20___ г."}',
    "nda": '{"city":"[город]","date":"«___» __________ 20___ г."}',
    "self_employed": '{"city":"[город]","date":"«___» __________ 20___ г."}',
    "warranty_letter": '{"city":"[город]","date":"«___» __________ 20___ г.","author":"[наименование/ФИО гаранта]"}',
    "doc": '{"org":"","sign":""}',
}

# Описание вида документа для промпта ИИ (что именно должно быть в разделах).
WORD_KIND_DESC = {
    "doc": "обычный документ",
    "referat": "реферат для университета. Титульный лист, содержание, введение (актуальность темы, цель и задачи), 2–3 главы, заключение с выводами (не пересказ содержания), список литературы не менее 5 источников. Текст живой, связный, как для сдачи преподавателю, не набор абзацев.",
    "report": "доклад для школы/университета. Титульный лист (упрощённый), краткое вступление, основная часть по теме с конкретными фактами, короткий вывод. Компактно и по существу, рассчитан на устное зачитывание.",
    "essay": "эссе. Без титульного листа или с упрощённым. Личная позиция автора, аргументированное рассуждение с примерами, живой связный текст, не сухой пересказ.",
    "notes": "конспект. Без титульного листа. Сжатое изложение материала по пунктам и подпунктам, ключевые определения выделены, без художественных отступлений — удобно для повторения перед экзаменом.",
    "coursework": "курсовая работа. Титульный лист, содержание, введение (актуальность темы, объект и предмет, цель и задачи исследования), 2-3 главы с подразделами (например 1.1, 1.2 и 2.1, 2.2) — теоретическая и практическая часть, заключение с самостоятельными выводами по каждой задаче (не пересказ содержания глав), список литературы не менее 5 источников. Это черновик-каркас под доработку с научным руководителем, а не финальная работа.",
    "dkp": "договор купли-продажи: предмет договора, цена и порядок расчётов, права и обязанности сторон, порядок передачи товара, ответственность сторон, срок действия и заключительные положения",
    "rent": "договор аренды: предмет аренды с точным описанием, срок аренды, размер и порядок внесения арендной платы, права и обязанности сторон, порядок передачи и возврата имущества, ответственность сторон",
    "offer": "коммерческое предложение уровня 2026 года: заголовок с конкретной выгодой или болью клиента (не 'КП от компании'), короткое описание сути в 1-2 абзаца, о компании, что именно входит в предложение, стоимость и условия оплаты, сроки и этапы, почему стоит выбрать именно нас (доказательства, кейсы), контакты и призыв к действию с дедлайном",
    "act": "акт приёма-передачи: дата и место составления, номер и основание (реквизиты договора), полные данные передающей и принимающей стороны, подробный перечень передаваемого имущества/товара/документов с количеством и состоянием, отметка об отсутствии претензий, место для подписей обеих сторон",
    "statement": "заявление",
    "proxy": "доверенность: кто выдаёт (доверитель), кому (представитель), точный перечень полномочий, срок действия",
    "loan": "расписка / договор займа: заимодавец и заёмщик (ФИО, паспортные данные), сумма займа цифрами и прописью, срок возврата, проценты (если есть) или указание на беспроцентный заём, порядок возврата, ответственность за просрочку (неустойка/пени)",
    "claim": "досудебная претензия: реквизиты адресата и заявителя, описание нарушения со ссылкой на договор/закон, конкретное требование (сумма, срок исполнения), срок для добровольного удовлетворения, предупреждение об обращении в суд",
    "consent": "согласие на выезд ребёнка за границу: ФИО и паспортные данные родителя (доверителя), ФИО, дата рождения и данные свидетельства о рождении/паспорта ребёнка, ФИО сопровождающего (если есть), страна/страны выезда, срок действия согласия",
    "marriage_contract": "брачный договор: ФИО супругов, реквизиты свидетельства о браке, режим собственности на имущество (совместное/раздельное/долевое), перечень конкретного имущества и порядок его раздела в случае развода",
    "gift": "договор дарения: ФИО и паспортные данные дарителя и одаряемого, точное описание предмета дарения, отсутствие встречного предоставления, момент перехода права собственности",
    "lawsuit": "исковое заявление в суд: наименование суда, данные истца и ответчика, цена иска (если применимо), обстоятельства дела по порядку, правовое обоснование со ссылками на закон, исковые требования, перечень приложений",
    "alimony": "соглашение об уплате алиментов: ФИО плательщика и получателя, ФИО и дата рождения ребёнка, размер алиментов (доля дохода или твёрдая сумма), периодичность и способ выплаты, индексация",
    "services": "договор возмездного оказания услуг: заказчик и исполнитель, точное описание услуги, стоимость и порядок оплаты, сроки оказания, порядок сдачи-приёмки, ответственность сторон",
    "employment": "трудовой договор: работодатель и работник, должность и трудовая функция, дата начала работы, режим рабочего времени, оклад и порядок выплаты зарплаты, права и обязанности сторон, испытательный срок (если есть)",
    "work_act": "акт выполненных работ (оказанных услуг): реквизиты основного договора, заказчик и исполнитель, перечень выполненных работ/услуг с объёмом, стоимостью и сроками, отметка о соответствии качества, отсутствие претензий",
    "supply": "договор поставки: поставщик и покупатель, наименование, ассортимент и количество товара, цена, сроки и порядок поставки, порядок приёмки товара по количеству и качеству, ответственность за просрочку",
    "agency": "агентский договор: принципал и агент, точный перечень юридических и фактических действий, которые агент совершает от имени и за счёт принципала, размер и порядок выплаты агентского вознаграждения, срок действия",
    "joint_activity": "договор о совместной деятельности (простого товарищества): участники, общая цель, вклад каждого участника (деньги, имущество, навыки), порядок ведения общих дел, распределение прибыли и убытков",
    "nonresidential_rent": "договор аренды нежилого помещения: арендодатель и арендатор, точный адрес и площадь помещения, назначение использования, срок аренды, размер и порядок внесения арендной платы, порядок передачи и возврата помещения",
    "cession": "договор цессии (уступки права требования): цедент и цессионарий, реквизиты первоначального обязательства/договора, объём и содержание уступаемого права требования, цена уступки, момент перехода права",
    "nda": "соглашение о неразглашении (NDA): стороны, точное определение конфиденциальной информации, цель передачи информации, обязательства по неразглашению и ограничению доступа, срок действия обязательств, ответственность за нарушение",
    "self_employed": "договор с самозанятым (на выполнение работ/оказание услуг): заказчик и исполнитель-самозанятый (ФИО, ИНН, номер чека НПД), предмет договора, стоимость, сроки, порядок сдачи-приёмки и оплаты",
    "warranty_letter": "гарантийное письмо: кому адресовано, кто гарантирует, чёткая суть гарантии (оплата, выполнение обязательства, устранение недостатков), срок исполнения гарантии",
}


def template_sections_for(kind):
    """Плейсхолдер-структура документа для режима 'Скачать шаблон' (без ИИ)."""
    if kind in ("dkp", "rent") or kind in CONTRACT_FAMILY:
        return [
            {"title": "Предмет договора", "content": "[точное описание предмета договора]"},
            {"title": "Цена и порядок расчётов", "content": "[сумма, порядок и сроки оплаты]"},
            {"title": "Права и обязанности сторон", "content": "[права и обязанности каждой стороны]"},
            {"title": "Ответственность сторон", "content": "[ответственность за нарушение условий]"},
            {"title": "Заключительные положения", "content": "[срок действия, порядок разрешения споров, реквизиты сторон]"},
        ]
    if kind == "act" or kind in ACT_FAMILY:
        return [{"title": "", "content": "[перечень передаваемого/выполненного — наименование, количество, состояние, стоимость]"}]
    if kind in ("statement", "claim") or kind in STATEMENT_FAMILY:
        return [{"title": "", "content": "Прошу [указать суть просьбы, требования или обстоятельства дела]."}]
    if kind in ("proxy", "consent") or kind in DECLARATION_FAMILY:
        return [{"title": "", "content": "[точный перечень полномочий или условий]"}]
    if kind == "loan":
        return [{"title": "", "content": "Сумма займа: [сумма цифрами и прописью]. Срок возврата: [дата]. Проценты: [указать или «беспроцентный»]."}]
    if kind == "offer":
        return [
            {"title": "Суть предложения", "content": "[что именно предлагается]"},
            {"title": "Стоимость и условия", "content": "[цена, порядок оплаты, сроки]"},
        ]
    return [{"title": "", "content": "[текст документа]"}]


def build_word(path, title, sections, kind="doc", meta=None):
    meta = meta or {}
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Cm(2)
    sec.bottom_margin = Cm(2)
    sec.left_margin = Cm(2.5)
    sec.right_margin = Cm(2)
    title = title or "Документ"

    if kind == "referat":
        # Поля, шрифт, интервал и отступ ниже соответствуют официальным требованиям
        # вузов к оформлению учебных работ (Times New Roman 14pt, интервал 1.2,
        # поля 20мм со всех сторон, абзацный отступ 1см).
        sec.left_margin = Cm(2)
        sec.right_margin = Cm(2)
        _student_title_page(doc, "РЕФЕРАТ", title, meta)
        _body(doc, sections, indent=True, head_center=True, size=14, line=1.2, first=1.0, head_size=14)

    elif kind == "coursework":
        sec.left_margin = Cm(2)
        sec.right_margin = Cm(2)
        _p(doc, meta.get("org") or "Министерство образования", 14, align=WD_ALIGN_PARAGRAPH.CENTER, after=0)
        _p(doc, meta.get("school") or "[Название учебного заведения]", 14, align=WD_ALIGN_PARAGRAPH.CENTER)
        if meta.get("specialty"):
            _p(doc, f"Специальность: {meta.get('specialty')}", 12, align=WD_ALIGN_PARAGRAPH.CENTER)
        for _ in range(4):
            doc.add_paragraph()
        _p(doc, "КУРСОВАЯ РАБОТА", 20, True, WD_ALIGN_PARAGRAPH.CENTER, after=8)
        _p(doc, title, 16, True, WD_ALIGN_PARAGRAPH.CENTER)
        for _ in range(6):
            doc.add_paragraph()
        _p(doc, f"Выполнил: {meta.get('author') or '[ФИО студента]'}", 14, align=WD_ALIGN_PARAGRAPH.RIGHT, after=0)
        _p(doc, f"Научный руководитель: {meta.get('teacher') or '[ФИО преподавателя]'}", 14, align=WD_ALIGN_PARAGRAPH.RIGHT)
        for _ in range(6):
            doc.add_paragraph()
        _p(doc, f"{meta.get('city') or '[Город]'} {meta.get('year') or '2026'}", 14, align=WD_ALIGN_PARAGRAPH.CENTER)
        doc.add_page_break()
        _body(doc, sections, indent=True, head_center=True, size=14, line=1.2, first=1.0, head_size=14)

    elif kind == "report":
        sec.left_margin = Cm(2)
        sec.right_margin = Cm(2)
        _student_title_page(doc, "ДОКЛАД", title, meta)
        _body(doc, sections, indent=True, head_center=False)

    elif kind == "essay":
        sec.left_margin = Cm(2)
        sec.right_margin = Cm(2)
        _student_title_page(doc, "ЭССЕ", title, meta)
        _body(doc, sections, indent=True, head_center=False)

    elif kind == "notes":
        _p(doc, "КОНСПЕКТ", 12, True, WD_ALIGN_PARAGRAPH.CENTER, after=0)
        _p(doc, title, 18, True, WD_ALIGN_PARAGRAPH.CENTER, before=4, after=14)
        _body(doc, sections, indent=False, head_center=False)

    elif kind in ("dkp", "rent"):
        cap = "ДОГОВОР КУПЛИ-ПРОДАЖИ" if kind == "dkp" else "ДОГОВОР АРЕНДЫ"
        _p(doc, cap, 16, True, WD_ALIGN_PARAGRAPH.CENTER, after=2)
        _p(doc, title if title not in (cap, "Документ") else "№ ______", 12, align=WD_ALIGN_PARAGRAPH.CENTER, after=10)
        p = doc.add_paragraph()
        p.paragraph_format.tab_stops.add_tab_stop(Cm(16), WD_TAB_ALIGNMENT.RIGHT)
        _run(p, f"{meta.get('city') or 'г. _______________'}\t{meta.get('date') or '«___» __________ 20___ г.'}", 12)
        _body(doc, sections, indent=True, head_center=False)
        left = "ПРОДАВЕЦ" if kind == "dkp" else "АРЕНДОДАТЕЛЬ"
        right = "ПОКУПАТЕЛЬ" if kind == "dkp" else "АРЕНДАТОР"
        _p(doc, f"{left}                    {right}", 12, True, before=18)
        _p(doc, "__________ / [ФИО] /          __________ / [ФИО] /", 12)

    elif kind == "offer":
        _p(doc, "КОММЕРЧЕСКОЕ ПРЕДЛОЖЕНИЕ", 11, True, after=0)
        _p(doc, meta.get("number") or "", 10, after=12)
        _p(doc, title, 20, True, after=14)
        for i, block in enumerate(sections or [], 1):
            head = (block.get("title") or "").strip()
            body = block.get("content") or ""
            if head:
                _p(doc, f"{i}. {head}", 14, True, before=12, after=6)
            for para in [x.strip() for x in body.split("\n") if x.strip()]:
                _p(doc, para, 12, after=6, line=1.2)

    elif kind == "act":
        _p(doc, "АКТ ПРИЁМА-ПЕРЕДАЧИ", 16, True, WD_ALIGN_PARAGRAPH.CENTER, after=2)
        _p(doc, title if title not in ("АКТ ПРИЁМА-ПЕРЕДАЧИ", "Документ") else "№ ______", 12,
           align=WD_ALIGN_PARAGRAPH.CENTER, after=10)
        p = doc.add_paragraph()
        p.paragraph_format.tab_stops.add_tab_stop(Cm(16), WD_TAB_ALIGNMENT.RIGHT)
        _run(p, f"{meta.get('city') or 'г. _______________'}\t{meta.get('date') or '«___» __________ 20___ г.'}", 12)
        if meta.get("basis"):
            _p(doc, f"Основание: {meta.get('basis')}", 12, after=10)
        _p(
            doc,
            f"{meta.get('from_party') or '[передающая сторона, ФИО/наименование]'}, именуемый(ая) в дальнейшем "
            "«Передающая сторона», с одной стороны, и "
            f"{meta.get('to_party') or '[принимающая сторона, ФИО/наименование]'}, именуемый(ая) в дальнейшем "
            "«Принимающая сторона», с другой стороны, составили настоящий акт о нижеследующем:",
            12, after=10, first=1.25
        )
        _body(doc, sections, indent=True, head_center=False)
        _p(doc, "Претензий по количеству, качеству и комплектности стороны друг к другу не имеют.", 12, before=8, after=18)
        _p(doc, "ПЕРЕДАЛ                    ПРИНЯЛ", 12, True, before=6)
        _p(doc, "__________ / [ФИО] /          __________ / [ФИО] /", 12)

    elif kind == "statement":
        _p(doc, meta.get("to") or "Директору [организация]", 12, align=WD_ALIGN_PARAGRAPH.RIGHT, after=0)
        _p(doc, meta.get("to_name") or "[ФИО руководителя]", 12, align=WD_ALIGN_PARAGRAPH.RIGHT, after=0)
        _p(doc, f"от {meta.get('from') or '[должность, ФИО]'}", 12, align=WD_ALIGN_PARAGRAPH.RIGHT, after=18)
        _p(doc, "ЗАЯВЛЕНИЕ", 16, True, WD_ALIGN_PARAGRAPH.CENTER, after=14)
        _body(doc, sections, indent=True, head_center=False)
        _p(doc, f"{meta.get('date') or '[дата]'}                    __________ / {meta.get('author') or '[ФИО]'} /", 12, before=24)

    elif kind == "proxy":
        _p(doc, "ДОВЕРЕННОСТЬ", 16, True, WD_ALIGN_PARAGRAPH.CENTER, after=10)
        p = doc.add_paragraph()
        p.paragraph_format.tab_stops.add_tab_stop(Cm(16), WD_TAB_ALIGNMENT.RIGHT)
        _run(p, f"{meta.get('city') or 'г. _______________'}\t{meta.get('date') or '«___» __________ 20___ г.'}", 12)
        _body(doc, sections, indent=True, head_center=False)
        _p(doc, f"Подпись доверителя: __________ / {meta.get('author') or '[ФИО]'} /", 12, before=20)

    elif kind == "loan":
        _p(doc, "РАСПИСКА В ПОЛУЧЕНИИ ДЕНЕЖНЫХ СРЕДСТВ", 16, True, WD_ALIGN_PARAGRAPH.CENTER, after=2)
        _p(doc, "(договор займа)", 11, align=WD_ALIGN_PARAGRAPH.CENTER, after=10)
        p = doc.add_paragraph()
        p.paragraph_format.tab_stops.add_tab_stop(Cm(16), WD_TAB_ALIGNMENT.RIGHT)
        _run(p, f"{meta.get('city') or 'г. _______________'}\t{meta.get('date') or '«___» __________ 20___ г.'}", 12)
        _p(
            doc,
            f"{meta.get('from_party') or '[Заимодавец, ФИО, паспортные данные]'}, именуемый(ая) в дальнейшем "
            "«Заимодавец», передал(а) в долг, а "
            f"{meta.get('to_party') or '[Заёмщик, ФИО, паспортные данные]'}, именуемый(ая) в дальнейшем "
            "«Заёмщик», получил(а) в долг денежные средства на нижеследующих условиях:",
            12, after=10, first=1.25
        )
        _body(doc, sections, indent=True, head_center=False)
        _p(doc, "ЗАИМОДАВЕЦ                    ЗАЁМЩИК", 12, True, before=18)
        _p(doc, "__________ / [ФИО] /          __________ / [ФИО] /", 12)

    elif kind == "claim":
        _p(doc, meta.get("to") or "[адресат претензии]", 12, align=WD_ALIGN_PARAGRAPH.RIGHT, after=0)
        _p(doc, f"от {meta.get('from') or '[ФИО, контакты заявителя]'}", 12, align=WD_ALIGN_PARAGRAPH.RIGHT, after=18)
        _p(doc, "ПРЕТЕНЗИЯ", 16, True, WD_ALIGN_PARAGRAPH.CENTER, after=4)
        if meta.get("basis"):
            _p(doc, f"Основание: {meta.get('basis')}", 11, align=WD_ALIGN_PARAGRAPH.CENTER, after=12)
        _body(doc, sections, indent=True, head_center=False)
        _p(
            doc,
            "В случае неудовлетворения настоящей претензии в указанный срок я буду вынужден(а) "
            "обратиться в суд за защитой своих прав со взысканием всех сопутствующих расходов.",
            12, before=10, after=18
        )
        _p(doc, f"{meta.get('date') or '[дата]'}                    __________ / {meta.get('from') or '[ФИО]'} /", 12)

    elif kind == "consent":
        _p(doc, "СОГЛАСИЕ", 16, True, WD_ALIGN_PARAGRAPH.CENTER, after=2)
        _p(doc, "на выезд несовершеннолетнего ребёнка за границу", 12, align=WD_ALIGN_PARAGRAPH.CENTER, after=10)
        p = doc.add_paragraph()
        p.paragraph_format.tab_stops.add_tab_stop(Cm(16), WD_TAB_ALIGNMENT.RIGHT)
        _run(p, f"{meta.get('city') or 'г. _______________'}\t{meta.get('date') or '«___» __________ 20___ г.'}", 12)
        _p(
            doc,
            f"Я, {meta.get('author') or '[ФИО родителя, паспортные данные]'}, даю согласие на выезд "
            "моего несовершеннолетнего ребёнка за пределы Российской Федерации на условиях, указанных ниже:",
            12, after=10, first=1.25
        )
        _body(doc, sections, indent=True, head_center=False)
        _p(doc, f"Подпись: __________ / {meta.get('author') or '[ФИО]'} /", 12, before=20)

    elif kind in CONTRACT_FAMILY:
        fam = CONTRACT_FAMILY[kind]
        _p(doc, fam["caption"], 16, True, WD_ALIGN_PARAGRAPH.CENTER, after=2)
        _p(doc, title if title not in (fam["caption"], "Документ") else "№ ______", 12,
           align=WD_ALIGN_PARAGRAPH.CENTER, after=10)
        p = doc.add_paragraph()
        p.paragraph_format.tab_stops.add_tab_stop(Cm(16), WD_TAB_ALIGNMENT.RIGHT)
        _run(p, f"{meta.get('city') or 'г. _______________'}\t{meta.get('date') or '«___» __________ 20___ г.'}", 12)
        _body(doc, sections, indent=True, head_center=False)
        _p(doc, f"{fam['left']}                    {fam['right']}", 12, True, before=18)
        _p(doc, "__________ / [ФИО/наименование] /          __________ / [ФИО/наименование] /", 12)

    elif kind in ACT_FAMILY:
        fam = ACT_FAMILY[kind]
        _p(doc, fam["caption"], 16, True, WD_ALIGN_PARAGRAPH.CENTER, after=2)
        _p(doc, title if title not in (fam["caption"], "Документ") else "№ ______", 12,
           align=WD_ALIGN_PARAGRAPH.CENTER, after=10)
        p = doc.add_paragraph()
        p.paragraph_format.tab_stops.add_tab_stop(Cm(16), WD_TAB_ALIGNMENT.RIGHT)
        _run(p, f"{meta.get('city') or 'г. _______________'}\t{meta.get('date') or '«___» __________ 20___ г.'}", 12)
        if meta.get("basis"):
            _p(doc, f"Основание: {meta.get('basis')}", 12, after=10)
        _p(
            doc,
            f"{meta.get('from_party') or '[исполнитель, ФИО/наименование]'}, именуемый(ая) в дальнейшем "
            "«Исполнитель», с одной стороны, и "
            f"{meta.get('to_party') or '[заказчик, ФИО/наименование]'}, именуемый(ая) в дальнейшем "
            "«Заказчик», с другой стороны, составили настоящий акт о нижеследующем:",
            12, after=10, first=1.25
        )
        _body(doc, sections, indent=True, head_center=False)
        _p(doc, "Претензий по объёму, качеству и срокам стороны друг к другу не имеют.", 12, before=8, after=18)
        _p(doc, "ИСПОЛНИТЕЛЬ                    ЗАКАЗЧИК", 12, True, before=6)
        _p(doc, "__________ / [ФИО] /          __________ / [ФИО] /", 12)

    elif kind in STATEMENT_FAMILY:
        fam = STATEMENT_FAMILY[kind]
        _p(doc, meta.get("to") or "[наименование суда/адресата]", 12, align=WD_ALIGN_PARAGRAPH.RIGHT, after=0)
        _p(doc, f"от {meta.get('from') or '[ФИО, адрес, контакты]'}", 12, align=WD_ALIGN_PARAGRAPH.RIGHT, after=18)
        _p(doc, fam["caption"], 16, True, WD_ALIGN_PARAGRAPH.CENTER, after=4)
        _body(doc, sections, indent=True, head_center=False)
        _p(doc, fam["closing"], 12, before=10, after=18)
        _p(doc, f"{meta.get('date') or '[дата]'}                    __________ / {meta.get('from') or '[ФИО]'} /", 12)

    elif kind in DECLARATION_FAMILY:
        fam = DECLARATION_FAMILY[kind]
        _p(doc, fam["caption"], 16, True, WD_ALIGN_PARAGRAPH.CENTER, after=10)
        p = doc.add_paragraph()
        p.paragraph_format.tab_stops.add_tab_stop(Cm(16), WD_TAB_ALIGNMENT.RIGHT)
        _run(p, f"{meta.get('city') or 'г. _______________'}\t{meta.get('date') or '«___» __________ 20___ г.'}", 12)
        _p(doc, fam["intro"].format(author=meta.get("author") or "[ФИО]"), 12, after=10, first=1.25)
        _body(doc, sections, indent=True, head_center=False)
        _p(doc, f"Подпись: __________ / {meta.get('author') or '[ФИО]'} /", 12, before=20)

    else:
        _p(doc, meta.get("org") or "", 11, True, after=0)
        _p(doc, title, 16, True, WD_ALIGN_PARAGRAPH.CENTER, before=8, after=12)
        _body(doc, sections, indent=True, head_center=False)
        if meta.get("sign"):
            _p(doc, "С уважением,", 12, before=16, after=0)
            _p(doc, meta.get("sign"), 12)

    doc.save(path)


# ==================== EXCEL ====================

def _xl_hex(color_tuple):
    return "%02X%02X%02X" % color_tuple


def _xl_tint(color_tuple, factor=0.9):
    """Осветляет цвет темы к белому (factor 0..1, чем больше - тем светлее).
    Используется для полос "зебры" в таблицах - лёгкий цветной акцент, а не серый по умолчанию."""
    r, g, b = color_tuple
    return (int(r + (255 - r) * factor), int(g + (255 - g) * factor), int(b + (255 - b) * factor))


def _xl_num(v, default=0.0):
    # Модель/пользователь должны прислать чистое число, но это не гарантировано -
    # иногда проскакивает "500 000", "500 000 ₽" или "12 месяцев" вместо 500000/12.
    # Вместо падения с ValueError на float()/int() достаём первое число из строки,
    # иначе одна нечисловая ячейка ломает формулу умножения/суммы на всём листе
    # (Excel даёт #VALUE! на любую арифметику с текстом).
    if v is None:
        return default
    if isinstance(v, (int, float)):
        return float(v)
    s = str(v).replace("\xa0", " ").replace(" ", "").replace(",", ".")
    m = re.search(r"-?\d+(\.\d+)?", s)
    return float(m.group()) if m else default


# Приблизительный курс к рублю на момент написания (конец августа 2026) - используется
# только для финмодели стартапа, если пользователь называет суммы не в рублях, а всё
# остальное в таблице считается в ₽. Курс плавает, поэтому в итоговый файл всегда
# добавляется примечание с датой и курсом, чтобы пользователь мог свериться и при
# необходимости попросить пересчитать по актуальному курсу.
FX_TO_RUB = {"RUB": 1.0, "РУБ": 1.0, "₽": 1.0, "USD": 85.0, "$": 85.0, "EUR": 95.0, "€": 95.0}
FX_NOTE_DATE = "конец августа 2026"


def _xl_fx_rate(currency):
    return FX_TO_RUB.get((currency or "RUB").strip().upper(), 1.0)


def _xl_style_sheet(ws, colors, n_cols, col_widths=None):
    """Базовое оформление листа под тему презентаций/документов: акцентная шапка,
    тонкие границы у таблицы, читаемая ширина колонок, альбомная печать по ширине
    (чтобы при печати/экспорте в PDF таблица не резалась на колонки по границе листа).
    col_widths - необязательный список ширин по колонкам (1-индексация по смыслу, но
    список 0-индексирован), подбирается под реальное содержимое конкретной таблицы -
    иначе везде была одна и та же ширина 22/16, из-за которой длинные названия обрезались,
    а короткие цифровые колонки занимали лишнее место."""
    ws.sheet_view.showGridLines = False
    for i in range(1, n_cols + 1):
        if col_widths and i - 1 < len(col_widths):
            width = col_widths[i - 1]
        else:
            width = 22 if i == 1 else 16
        ws.column_dimensions[get_column_letter(i)].width = width
    ws.page_setup.orientation = "landscape"
    ws.sheet_properties.pageSetUpPr.fitToPage = True
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 0
    ws.print_options.horizontalCentered = False
    ws.sheet_properties.tabColor = _xl_hex(colors["line"])


def _xl_name_col_width(rows, key="name", min_w=18, max_w=44, pad=4):
    """Подбирает ширину текстовой колонки под самое длинное реальное значение в данных,
    а не держит её фиксированной - иначе длинные названия статей обрезались визуально."""
    longest = max((len(str(r.get(key, ""))) for r in rows), default=min_w)
    return max(min_w, min(max_w, longest + pad))


def _xl_finalize_table(ws, colors, header_row, first_data_row, last_data_row, n_cols):
    """Общие штрихи для готовой таблицы: заморозка шапки при прокрутке, автофильтр
    по шапке (быстро сортировать/фильтровать позиции без формул), автовысота шапки."""
    ws.freeze_panes = f"A{first_data_row}"
    last_col_letter = get_column_letter(n_cols)
    ws.auto_filter.ref = f"A{header_row}:{last_col_letter}{last_data_row}"
    ws.row_dimensions[header_row].height = 26


def _xl_header_row(ws, row_idx, headers, colors):
    fill = PatternFill(start_color=_xl_hex(colors["line"]), end_color=_xl_hex(colors["line"]), fill_type="solid")
    font = XlFont(bold=True, color=_xl_hex((255, 255, 255) if sum(colors["line"]) < 400 else (20, 20, 20)), size=11)
    thin = Side(style="thin", color=_xl_hex(colors["mute"]))
    for i, h in enumerate(headers, 1):
        c = ws.cell(row=row_idx, column=i, value=h)
        c.fill = fill
        c.font = font
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        c.border = Border(left=thin, right=thin, top=thin, bottom=thin)


def _xl_body_row(ws, row_idx, values, colors, number_cols=None, money_cols=None, percent_cols=None, zebra=False):
    thin = Side(style="thin", color=_xl_hex(colors["mute"]))
    zebra_fill = None
    if zebra:
        zebra_fill = PatternFill(start_color=_xl_hex(_xl_tint(colors["line"], 0.9)),
                                  end_color=_xl_hex(_xl_tint(colors["line"], 0.9)), fill_type="solid")
    for i, v in enumerate(values, 1):
        c = ws.cell(row=row_idx, column=i, value=v)
        c.border = Border(left=thin, right=thin, top=thin, bottom=thin)
        c.alignment = Alignment(horizontal="left" if i == 1 else "center", vertical="center")
        if zebra_fill:
            c.fill = zebra_fill
        # Формулы (вычисляемые колонки вроде "Сумма") остаются защищёнными при
        # включённой защите листа, а введённые вручную/моделью значения - редактируемые.
        c.protection = Protection(locked=isinstance(v, str) and v.startswith("="))
        if money_cols and i in money_cols:
            c.number_format = '#,##0.00 ₽'
        elif percent_cols and i in percent_cols:
            c.number_format = '0.0%'
        elif number_cols and i in number_cols:
            c.number_format = '#,##0.##'
    ws.row_dimensions[row_idx].height = 20


def _xl_total_row(ws, row_idx, label, values, colors, n_cols, money_cols=None, percent_cols=None):
    fill = PatternFill(start_color=_xl_hex(colors["bg"] if sum(colors["bg"]) < 400 else (235, 235, 235)),
                        end_color=_xl_hex(colors["bg"] if sum(colors["bg"]) < 400 else (235, 235, 235)), fill_type="solid")
    thin = Side(style="thin", color=_xl_hex(colors["mute"]))
    ws.cell(row=row_idx, column=1, value=label).font = XlFont(bold=True)
    for i in range(1, n_cols + 1):
        c = ws.cell(row=row_idx, column=i)
        c.border = Border(left=thin, right=thin, top=Side(style="double", color=_xl_hex(colors["mute"])), bottom=thin)
        c.font = XlFont(bold=True)
        if i > 1 and i - 2 < len(values) and values[i - 2] is not None:
            c.value = values[i - 2]
            if money_cols and i in money_cols:
                c.number_format = '#,##0.00 ₽'
            elif percent_cols and i in percent_cols:
                c.number_format = '0.0%'


def _xl_title(ws, title, colors, n_cols):
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=n_cols)
    c = ws.cell(row=1, column=1, value=title)
    c.font = XlFont(bold=True, size=14, color=_xl_hex(colors["ink"] if sum(colors["bg"]) > 400 else (30, 30, 30)))
    ws.row_dimensions[1].height = 28


def _xl_subtitle(ws, text, colors, n_cols, row=2):
    """Раньше строка 2 всегда пустовала (просто визуальный отступ перед шапкой) - теперь
    несёт дату создания файла, чтобы было видно, насколько таблица свежая."""
    ws.merge_cells(start_row=row, start_column=1, end_row=row, end_column=n_cols)
    c = ws.cell(row=row, column=1, value=text)
    c.font = XlFont(italic=True, size=9, color=_xl_hex(colors["mute"]))


def _xl_category_list(rows):
    """Список уникальных категорий в порядке первого появления, или None, если ни у одной
    строки не заполнено поле category (тогда доп. колонка и блок итогов по категориям не нужны)."""
    seen = []
    for row in rows:
        cat = (row.get("category") or "").strip()
        if cat and cat not in seen:
            seen.append(cat)
    return seen or None


def _xl_category_summary(ws, colors, start_row, cats, money_col_letter, cat_col_letter,
                          first_item_row, last_item_row, sums_by_cat, n_cols):
    """Небольшой блок 'Итоги по категориям' под основной таблицей - реальная формула SUMIF
    по вспомогательной колонке с категорией, а не текст. cats - список категорий по порядку
    появления, sums_by_cat - {категория: посчитанная кодом сумма} для кеша значений формул.
    Подпись категории кладём в колонку "Статья"/"Наименование" (она достаточно широкая под
    длинные названия), а сумму - прямо под колонкой "Сумма" основной таблицы, чтобы визуально
    совпадало с ней по вертикали."""
    cache = {}
    money_col_idx = column_index_from_string(money_col_letter)
    r = start_row
    ws.cell(row=r, column=1, value="Итоги по категориям").font = XlFont(bold=True, size=11)
    r += 1
    for cat in cats:
        label_c = ws.cell(row=r, column=2, value=cat)
        label_c.font = XlFont(italic=True)
        formula = (f'=SUMIF({cat_col_letter}{first_item_row}:{cat_col_letter}{last_item_row},"{cat}",'
                   f'{money_col_letter}{first_item_row}:{money_col_letter}{last_item_row})')
        val_c = ws.cell(row=r, column=money_col_idx, value=formula)
        val_c.font = XlFont(italic=True, bold=True)
        val_c.number_format = '#,##0.00 ₽'
        cache[f"{money_col_letter}{r}"] = sums_by_cat.get(cat, 0.0)
        r += 1
    return r, cache


def _xl_add_type_dropdown(ws, col_letter, first_row, last_row, options):
    """Выпадающий список вместо свободного текста в колонке 'Тип' - чтобы при редактировании
    пользователь не сломал формулу SUMIF опечаткой ("доходы" вместо "доход")."""
    if last_row < first_row:
        return
    dv = DataValidation(type="list", formula1='"' + ",".join(options) + '"', allow_blank=True)
    ws.add_data_validation(dv)
    dv.add(f"{col_letter}{first_row}:{col_letter}{last_row}")


def _xl_color_scale(ws, col_letter, first_row, last_row):
    """Цветовая шкала на денежной/числовой колонке - крупные суммы визуально заметнее
    без необходимости вчитываться в цифры."""
    if last_row < first_row:
        return
    rule = ColorScaleRule(
        start_type="min", start_color="FFFFE9CC",
        mid_type="percentile", mid_value=50, mid_color="FFFFB25C",
        end_type="max", end_color="FFCC5200",
    )
    ws.conditional_formatting.add(f"{col_letter}{first_row}:{col_letter}{last_row}", rule)


def _xl_add_bar_chart(ws, n_cols, cat_col_idx, val_col_idx, header_row, first_row, last_row,
                       anchor_row=None, y_title="Сумма, ₽"):
    """Столбчатая диаграмма по данным, которые уже лежат в таблице (никаких лишних вычислений -
    просто ссылка на диапазон). header_row нужен, чтобы подписью серии стал заголовок колонки."""
    if last_row < first_row:
        return
    chart = BarChart()
    chart.type = "col"
    chart.style = 10
    chart.y_axis.title = y_title
    chart.x_axis.title = None
    chart.legend = None
    data_ref = Reference(ws, min_col=val_col_idx, min_row=header_row, max_row=last_row)
    cats_ref = Reference(ws, min_col=cat_col_idx, min_row=first_row, max_row=last_row)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)
    chart.height = 9
    chart.width = 17
    anchor = f"{get_column_letter(n_cols + 2)}{anchor_row or header_row}"
    ws.add_chart(chart, anchor)


def _xl_add_pie_chart(ws, n_cols, cat_col_idx, val_col_idx, header_row, first_row, last_row, anchor_row=None):
    """Круговая диаграмма долей - для показателей/значений, где важна доля от целого,
    а не сравнение абсолютных величин. Подписи - только процент (без дублирования названия
    и значения на каждом секторе), полное название категории и так видно в легенде."""
    if last_row < first_row:
        return
    chart = PieChart()
    chart.style = 10
    data_ref = Reference(ws, min_col=val_col_idx, min_row=header_row, max_row=last_row)
    cats_ref = Reference(ws, min_col=cat_col_idx, min_row=first_row, max_row=last_row)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)
    chart.height = 9
    chart.width = 19
    dl = DataLabelList()
    dl.showCatName = False
    dl.showLegendKey = False
    dl.showSerName = False
    dl.showVal = False
    dl.showPercent = True
    dl.showBubbleSize = False
    chart.dataLabels = dl
    anchor = f"{get_column_letter(n_cols + 2)}{anchor_row or header_row}"
    ws.add_chart(chart, anchor)


def _xl_protect_sheet(ws):
    """Защищает формулы/шапку/итоги от случайной правки (нельзя затереть формулу опечаткой),
    но оставляет редактируемыми ячейки с реальными данными - их пользователь мог бы захотеть
    подправить руками. Без пароля - это мягкая защита от опечаток, а не секьюрити: снимается
    в один клик через Рецензирование → Снять защиту листа, если понадобится."""
    ws.protection.sheet = True
    ws.protection.formatCells = False
    ws.protection.formatColumns = False
    ws.protection.formatRows = False
    ws.protection.sort = False
    ws.protection.autoFilter = False
    ws.protection.selectLockedCells = False
    ws.protection.selectUnlockedCells = False


def _xl_inject_cached_values(path, cache):
    """openpyxl пишет в ячейку с формулой только сам текст формулы ("=SUM(...)"),
    без результата - настоящий Excel/Numbers/Google Таблицы пересчитывают всё сами
    при открытии, но быстрый просмотр файла (например Quick Look на iPhone) формулы
    не считает и показывает пустую ячейку или 0. Чтобы число было видно в любом
    просмотрщике, а не только в "настоящем" Excel, здесь вручную дописывается
    закешированный результат прямо в XML уже сохранённого файла - формула остаётся
    рабочей и редактируемой, просто рядом с ней лежит готовое число на первый показ.
    cache: {"E8": 4500, "B13": 59900.0, ...} - адрес ячейки -> посчитанное код."""
    import zipfile
    import shutil

    sheet_path = "xl/worksheets/sheet1.xml"
    tmp_path = path + ".tmp"
    with zipfile.ZipFile(path, "r") as zin:
        data = {n: zin.read(n) for n in zin.namelist()}
    xml = data[sheet_path].decode("utf-8")
    for ref, val in cache.items():
        if val is None:
            continue
        # openpyxl уже пишет пустой <v></v> (или самозакрывающийся <v/>) сразу после
        # </f> - подставляем число внутрь него, а не добавляем новый тег.
        pattern = re.compile(r'(<c r="%s"[^>]*>(?:(?!</c>).)*?</f>)<v\s*/?>(?:</v>)?(</c>)' % re.escape(ref), re.DOTALL)
        xml, n = pattern.subn(lambda m: f"{m.group(1)}<v>{val}</v>{m.group(2)}", xml, count=1)
        if not n:
            print(f"_xl_inject_cached_values: не нашёл ячейку {ref} с формулой (не критично)")
    data[sheet_path] = xml.encode("utf-8")
    with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for n, content in data.items():
            zout.writestr(n, content)
    shutil.move(tmp_path, path)


# Фиксированные схемы колонок под каждый вид - модель (или сам пользователь) поставляет
# только содержимое строк (rows), формулы и оформление всегда собираются кодом, а не ИИ,
# чтобы суммы/проценты/итоги были настоящими формулами Excel, а не текстом.
EXCEL_ROW_SCHEMA = {
    "expense_estimate": {"fields": ["name", "qty", "price", "category"], "example": '{"name":"...", "qty": 1, "price": 1000, "category": "..."}'},
    "family_budget": {"fields": ["name", "type", "amount", "category"], "example": '{"name":"...", "type": "доход"|"расход", "amount": 1000, "category": "..."}'},
    "project_budget": {"fields": ["name", "type", "amount", "category"], "example": '{"name":"...", "type": "доход"|"расход", "amount": 1000, "category": "..."}'},
    "price_list": {"fields": ["name", "unit", "qty", "price", "discount", "category"], "example": '{"name":"...", "unit":"шт", "qty": 1, "price": 1000, "discount": 0, "category": "..."}'},
    "calc_table": {"fields": ["name", "value"], "example": '{"name":"...", "value": 100}'},
}

# Минимум строк для режима "сгенерировать самому" (ai) - раньше промпт просил всего
# 5-8 строк, из-за чего таблицы выглядели пустыми/скудными. Для family_budget/project_budget
# минимум относится к каждой категории (доход/расход) по отдельности, а не к общему числу строк,
# иначе модель могла просто выдать 10 расходов и 0 доходов.
EXCEL_MIN_ROWS_AI = {
    "expense_estimate": 12,
    "family_budget": 12,
    "project_budget": 12,
    "price_list": 14,
    "calc_table": 10,
}


def build_excel_items(path, title, kind, colors, rows, subtitle=None):
    """Собирает смету/бюджет/прайс/расчётную таблицу с реальными формулами Excel.
    rows - список dict по схеме EXCEL_ROW_SCHEMA[kind]. Если хотя бы у одной строки заполнено
    поле category - добавляется колонка "Категория" и блок промежуточных итогов по ней (реальные
    формулы SUMIF, не текст). В конце - диаграмма по данным таблицы и мягкая защита формул."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Таблица"
    cache = {}  # адрес ячейки -> посчитанное код значение, для _xl_inject_cached_values
    cats = _xl_category_list(rows) if kind != "calc_table" else None
    has_cat = cats is not None

    if kind == "expense_estimate":
        headers = ["№", "Статья расходов", "Кол-во", "Цена, ₽", "Сумма, ₽"] + (["Категория"] if has_cat else [])
        n_cols = len(headers)
        col_widths = [6, _xl_name_col_width(rows), 10, 14, 15] + ([18] if has_cat else [])
        _xl_style_sheet(ws, colors, n_cols, col_widths=col_widths)
        _xl_title(ws, title, colors, n_cols)
        if subtitle:
            _xl_subtitle(ws, subtitle, colors, n_cols)
        _xl_header_row(ws, 3, headers, colors)
        r = 4
        total = 0.0
        sums_by_cat = {}
        for i, row in enumerate(rows, 1):
            qty, price = _xl_num(row.get("qty"), 1), _xl_num(row.get("price"))
            line_sum = qty * price
            total += line_sum
            cat = (row.get("category") or "").strip()
            values = [i, row.get("name", ""), qty, price, f"=C{r}*D{r}"]
            if has_cat:
                values.append(cat or "Прочее")
                sums_by_cat[cat or "Прочее"] = sums_by_cat.get(cat or "Прочее", 0.0) + line_sum
            _xl_body_row(ws, r, values, colors, number_cols={3}, money_cols={4, 5}, zebra=(i % 2 == 0))
            cache[f"E{r}"] = line_sum
            r += 1
        last_item_row = r - 1
        total_values = [None, None, None, f"=SUM(E4:E{last_item_row})"] + ([None] if has_cat else [])
        _xl_total_row(ws, r, "ИТОГО", total_values, colors, n_cols, money_cols={5})
        cache[f"E{r}"] = total
        _xl_finalize_table(ws, colors, 3, 4, last_item_row, n_cols)
        _xl_color_scale(ws, "E", 4, last_item_row)
        _xl_add_bar_chart(ws, n_cols, cat_col_idx=2, val_col_idx=5, header_row=3, first_row=4, last_row=last_item_row)
        if has_cat:
            summary_row, summary_cache = _xl_category_summary(
                ws, colors, r + 2, cats, "E", get_column_letter(n_cols), 4, last_item_row, sums_by_cat, n_cols)
            cache.update(summary_cache)

    elif kind in ("family_budget", "project_budget"):
        headers = ["№", "Статья", "Тип", "Сумма, ₽"] + (["Категория"] if has_cat else [])
        n_cols = len(headers)
        col_widths = [6, _xl_name_col_width(rows), 12, 15] + ([18] if has_cat else [])
        _xl_style_sheet(ws, colors, n_cols, col_widths=col_widths)
        _xl_title(ws, title, colors, n_cols)
        if subtitle:
            _xl_subtitle(ws, subtitle, colors, n_cols)
        _xl_header_row(ws, 3, headers, colors)
        income_label, expense_label = "Доход", "Расход"
        r = 4
        total_inc, total_exp = 0.0, 0.0
        sums_by_cat = {}
        for i, row in enumerate(rows, 1):
            # Нормализуем "Тип" в самом коде, а не полагаемся на то, что модель/пользователь
            # напишет ровно "доход"/"расход" - иначе SUMIF ниже просто не найдёт совпадение
            # (он сравнивает строки целиком, а не по смыслу) и итог молча окажется нулевым.
            raw_type = (row.get("type") or "").strip().lower()
            norm_type = expense_label if "расход" in raw_type else income_label
            amount = _xl_num(row.get("amount"))
            if norm_type == income_label:
                total_inc += amount
            else:
                total_exp += amount
            cat = (row.get("category") or "").strip()
            values = [i, row.get("name", ""), norm_type, amount]
            if has_cat:
                values.append(cat or "Прочее")
                key = cat or "Прочее"
                # Кешируем сумму ТАК ЖЕ, как её посчитает настоящая формула SUMIF в Excel
                # (она суммирует "Сумма, ₽" как есть, без учёта знака дохода/расхода) -
                # иначе быстрый просмотр показывал бы одно число, а пересчёт в Excel - другое.
                sums_by_cat[key] = sums_by_cat.get(key, 0.0) + amount
            _xl_body_row(ws, r, values, colors, money_cols={4}, zebra=(i % 2 == 0))
            r += 1
        last = r - 1
        _xl_finalize_table(ws, colors, 3, 4, last, n_cols)
        _xl_color_scale(ws, "D", 4, last)
        _xl_add_type_dropdown(ws, "C", 4, last, [income_label, expense_label])
        r_inc, r_exp = r, r + 1
        ws.cell(row=r_inc, column=1, value="").border = Border()
        inc_values = [None, None, f'=SUMIF(C4:C{last},"{income_label}",D4:D{last})'] + ([None] if has_cat else [])
        exp_values = [None, None, f'=SUMIF(C4:C{last},"{expense_label}",D4:D{last})'] + ([None] if has_cat else [])
        _xl_total_row(ws, r_inc, "Итого доходы", inc_values, colors, n_cols, money_cols={4})
        cache[f"D{r_inc}"] = total_inc
        _xl_total_row(ws, r_exp, "Итого расходы", exp_values, colors, n_cols, money_cols={4})
        cache[f"D{r_exp}"] = total_exp
        r_res = r_exp + 1
        res_label = "Остаток" if kind == "family_budget" else "Прибыль"
        res_values = [None, None, f"=D{r_inc}-D{r_exp}"] + ([None] if has_cat else [])
        _xl_total_row(ws, r_res, res_label, res_values, colors, n_cols, money_cols={4})
        result_val = total_inc - total_exp
        cache[f"D{r_res}"] = result_val
        # Цветовой акцент итога - зелёный, если остаток/прибыль положительные, красный,
        # если ушли в минус: так сразу видно результат, не читая формулу.
        res_cell = ws.cell(row=r_res, column=4)
        res_cell.font = XlFont(bold=True, color=_xl_hex((30, 140, 70) if result_val >= 0 else (190, 40, 40)))
        chart = BarChart()
        chart.type = "col"
        chart.style = 10
        chart.legend = None
        chart.y_axis.title = "₽"
        data_ref = Reference(ws, min_col=4, min_row=r_inc, max_row=r_exp)
        cats_ref = Reference(ws, min_col=1, min_row=r_inc, max_row=r_exp)
        chart.add_data(data_ref, titles_from_data=False)
        chart.set_categories(cats_ref)
        chart.title = "Доходы и расходы"
        chart.height = 9
        chart.width = 17
        ws.add_chart(chart, f"{get_column_letter(n_cols + 2)}3")
        if has_cat:
            summary_row, summary_cache = _xl_category_summary(
                ws, colors, r_res + 2, cats, "D", get_column_letter(n_cols), 4, last, sums_by_cat, n_cols)
            cache.update(summary_cache)

    elif kind == "price_list":
        headers = ["№", "Наименование", "Ед.", "Кол-во", "Цена, ₽", "Скидка, %", "Сумма, ₽"] + (["Категория"] if has_cat else [])
        n_cols = len(headers)
        col_widths = [6, _xl_name_col_width(rows), 8, 10, 14, 12, 15] + ([18] if has_cat else [])
        _xl_style_sheet(ws, colors, n_cols, col_widths=col_widths)
        _xl_title(ws, title, colors, n_cols)
        if subtitle:
            _xl_subtitle(ws, subtitle, colors, n_cols)
        _xl_header_row(ws, 3, headers, colors)
        r = 4
        total = 0.0
        sums_by_cat = {}
        for i, row in enumerate(rows, 1):
            disc = _xl_num(row.get("discount"))
            qty, price = _xl_num(row.get("qty"), 1), _xl_num(row.get("price"))
            disc_frac = disc / 100 if disc else 0
            line_sum = qty * price * (1 - disc_frac)
            total += line_sum
            cat = (row.get("category") or "").strip()
            values = [i, row.get("name", ""), row.get("unit", "шт"), qty, price, disc_frac, f"=D{r}*E{r}*(1-F{r})"]
            if has_cat:
                values.append(cat or "Прочее")
                sums_by_cat[cat or "Прочее"] = sums_by_cat.get(cat or "Прочее", 0.0) + line_sum
            _xl_body_row(ws, r, values, colors, number_cols={4}, money_cols={5, 7}, percent_cols={6}, zebra=(i % 2 == 0))
            cache[f"G{r}"] = line_sum
            r += 1
        last_item_row = r - 1
        total_values = [None, None, None, None, None, f"=SUM(G4:G{last_item_row})"] + ([None] if has_cat else [])
        _xl_total_row(ws, r, "ИТОГО", total_values, colors, n_cols, money_cols={7})
        cache[f"G{r}"] = total
        _xl_finalize_table(ws, colors, 3, 4, last_item_row, n_cols)
        _xl_color_scale(ws, "G", 4, last_item_row)
        _xl_add_bar_chart(ws, n_cols, cat_col_idx=2, val_col_idx=7, header_row=3, first_row=4, last_row=last_item_row)
        if has_cat:
            summary_row, summary_cache = _xl_category_summary(
                ws, colors, r + 2, cats, "G", get_column_letter(n_cols), 4, last_item_row, sums_by_cat, n_cols)
            cache.update(summary_cache)

    elif kind == "calc_table":
        headers = ["№", "Показатель", "Значение", "Доля, %"]
        n_cols = len(headers)
        _xl_style_sheet(ws, colors, n_cols, col_widths=[6, _xl_name_col_width(rows), 13, 12])
        _xl_title(ws, title, colors, n_cols)
        if subtitle:
            _xl_subtitle(ws, subtitle, colors, n_cols)
        _xl_header_row(ws, 3, headers, colors)
        r = 4
        first_r = r
        values = []
        for i, row in enumerate(rows, 1):
            v = _xl_num(row.get("value"))
            values.append(v)
            _xl_body_row(ws, r, [i, row.get("name", ""), v, None], colors, number_cols={3}, percent_cols={4}, zebra=(i % 2 == 0))
            r += 1
        last_r = r - 1
        _xl_finalize_table(ws, colors, 3, first_r, last_r, n_cols)
        _xl_add_pie_chart(ws, n_cols, cat_col_idx=2, val_col_idx=3, header_row=3, first_row=first_r, last_row=last_r)
        total_value = sum(values)
        # IFERROR - если все значения окажутся нулевыми (или их сумма равна нулю), формула
        # деления вернёт #DIV/0! в каждой строке; выводим 0% вместо ошибки на весь лист.
        for rr, v in zip(range(first_r, r), values):
            ws.cell(row=rr, column=4, value=f"=IFERROR(C{rr}/SUM($C${first_r}:$C${last_r}),0)")
            ws.cell(row=rr, column=4).number_format = '0.0%'
            ws.cell(row=rr, column=4).protection = Protection(locked=True)
            cache[f"D{rr}"] = (v / total_value) if total_value else 0.0
        total_pct_formula = f"=IFERROR(SUM(C{first_r}:C{last_r})/SUM(C{first_r}:C{last_r}),0)"
        _xl_total_row(ws, r, "ИТОГО", [None, f"=SUM(C{first_r}:C{last_r})", total_pct_formula], colors, n_cols, percent_cols={4})
        cache[f"C{r}"] = total_value
        cache[f"D{r}"] = 1.0 if total_value else 0.0

    _xl_protect_sheet(ws)
    wb.save(path)
    if cache:
        _xl_inject_cached_values(path, cache)


def build_excel_startup(path, title, colors, data):
    """Финмодель стартапа - данные (инвестиции, расходы, цена, продажи, срок) реальные,
    от пользователя, никогда не придумываются моделью. Всё, что видно в файле кроме
    самих исходных чисел - настоящие формулы Excel (выручка/прибыль/накопительный итог)."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Финмодель"

    horizon = max(1, min(int(_xl_num(data.get("horizon_months"), 12)), 36))
    investment = _xl_num(data.get("investment"))
    price = _xl_num(data.get("price"))
    monthly_sales = _xl_num(data.get("monthly_sales"))
    expenses = [{"name": (e or {}).get("name", ""), "amount": _xl_num((e or {}).get("amount"))}
                for e in (data.get("expenses") or [])]

    n_cols = 5
    _xl_style_sheet(ws, colors, n_cols)
    _xl_title(ws, title, colors, n_cols)
    cache = {}  # адрес ячейки -> посчитанное код значение, для _xl_inject_cached_values

    ws.cell(row=2, column=1, value=f"Стартовые вложения: {investment:,.0f} ₽ · Сформировано: {datetime.now().strftime('%d.%m.%Y')}".replace(",", " ")).font = XlFont(italic=True, size=10)

    # блок постоянных ежемесячных расходов - отдельной таблицей, чтобы сумма расходов
    # в помесячном расчёте тоже была формулой (SUM), а не готовым числом
    r = 4
    ws.cell(row=r, column=1, value="Ежемесячные расходы").font = XlFont(bold=True)
    r += 1
    exp_first = r
    total_exp_value = sum((e.get("amount") or 0) for e in expenses)
    for i, e in enumerate(expenses, 1):
        _xl_body_row(ws, r, [None, e.get("name", ""), None, None, e.get("amount", 0)], colors, money_cols={5}, zebra=(i % 2 == 0))
        r += 1
    exp_last = r - 1 if expenses else exp_first
    if not expenses:
        _xl_body_row(ws, r, [None, "—", None, None, 0], colors, money_cols={5})
        exp_last = r
        r += 1
    total_exp_row = r
    _xl_total_row(ws, total_exp_row, "Итого расходов в месяц", [None, None, None, f"=SUM(E{exp_first}:E{exp_last})"], colors, n_cols, money_cols={5})
    cache[f"E{total_exp_row}"] = total_exp_value
    r = total_exp_row + 1

    fx_note = data.get("_fx_note")
    if fx_note:
        ws.cell(row=r, column=1, value=f"💱 {fx_note}").font = XlFont(italic=True, size=9, color=_xl_hex(colors["mute"]))
        r += 1

    r += 1

    ws.cell(row=r, column=1, value=f"Цена за единицу: {price:,.0f} ₽ · Продажи в месяц: {monthly_sales:,.0f} шт.".replace(",", " ")).font = XlFont(italic=True, size=10)
    r += 2

    headers = ["Месяц", "Выручка, ₽", "Расходы, ₽", "Прибыль, ₽", "Накопительно, ₽"]
    _xl_header_row(ws, r, headers, colors)
    table_start = r + 1
    r = table_start

    # cum_values считаем заранее (не только для "срока окупаемости", но и чтобы
    # сразу знать, какое число закешировать в каждую ячейку E - формула та же математика)
    rev_value = price * monthly_sales
    profit_value = rev_value - total_exp_value
    cum_values = []
    running = -investment
    for month in range(1, horizon + 1):
        running += profit_value
        cum_values.append(running)

    for month in range(1, horizon + 1):
        rev_formula = f"={price}*{monthly_sales}"
        exp_formula = f"=E{total_exp_row}"
        profit_formula = f"=B{r}-C{r}"
        if month == 1:
            cum_formula = f"=D{r}-{investment}"
        else:
            cum_formula = f"=D{r}+E{r - 1}"
        _xl_body_row(ws, r, [f"Месяц {month}", rev_formula, exp_formula, profit_formula, cum_formula], colors, money_cols={2, 3, 4, 5}, zebra=(month % 2 == 0))
        cache[f"B{r}"] = rev_value
        cache[f"C{r}"] = total_exp_value
        cache[f"D{r}"] = profit_value
        cache[f"E{r}"] = cum_values[month - 1]
        r += 1
    table_end = r - 1
    _xl_finalize_table(ws, colors, table_start - 1, table_start, table_end, n_cols)

    # срок окупаемости - первый месяц, где накопительный итог формулой стал неотрицательным;
    # считается по-настоящему из реальных чисел пользователя (это математика, а не выдумка),
    # но записывается как обычная сводная ячейка, а не как сложная формула поиска,
    # чтобы не зависеть от array-формул, которые по-разному ведут себя в разных версиях Excel.
    payback = next((i + 1 for i, v in enumerate(cum_values) if v >= 0), None)

    r += 1
    ws.cell(row=r, column=1, value="Срок окупаемости").font = XlFont(bold=True)
    ws.cell(row=r, column=2, value=(f"{payback} мес." if payback else f"не достигается за {horizon} мес."))
    r += 1
    ws.cell(row=r, column=1, value="Точка безубыточности (продаж/мес)").font = XlFont(bold=True)
    be_cell = ws.cell(row=r, column=2, value=f"=ROUNDUP(E{total_exp_row}/{price},0)" if price else "—")
    if price:
        be_cell.number_format = '#,##0 "шт."'
        import math
        cache[f"B{r}"] = math.ceil(total_exp_value / price)

    for i in range(1, n_cols + 1):
        ws.column_dimensions[get_column_letter(i)].width = 20
    ws.column_dimensions["A"].width = 40

    # Линейный график прибыли и накопительного итога по месяцам - сразу видно момент
    # выхода в плюс (пересечение накопительной линии с нулём), а не только число в ячейке.
    chart = LineChart()
    chart.style = 12
    chart.y_axis.title = "₽"
    chart.x_axis.title = "Месяц"
    data_ref = Reference(ws, min_col=4, max_col=5, min_row=table_start - 1, max_row=table_end)
    cats_ref = Reference(ws, min_col=1, min_row=table_start, max_row=table_end)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)
    chart.height = 9
    chart.width = 19
    ws.add_chart(chart, f"{get_column_letter(n_cols + 2)}{table_start - 1}")

    _xl_protect_sheet(ws)
    wb.save(path)
    _xl_inject_cached_values(path, cache)


EXCEL_KIND_DESC = {
    "expense_estimate": "смета расходов",
    "family_budget": "семейный бюджет",
    "project_budget": "смета проекта / бизнес-план",
    "price_list": "прайс-лист",
    "calc_table": "расчётная таблица к учебной работе",
    "startup_model": "финансовая модель стартапа",
}


def excel_mode_kb(lang="ru"):
    return ReplyKeyboardMarkup(keyboard=[
        [KeyboardButton(text=tr("btn_ai_generate", lang))],
        [KeyboardButton(text=tr("btn_own_data", lang))],
        [KeyboardButton(text=tr("btn_main_menu", lang))],
    ], resize_keyboard=True)


@dp.message(F.text.in_(ALL_BTN_EXCEL_LABELS))
async def start_excel(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    if not can_generate(m.from_user.id):
        await m.answer(tr("msg_limit", lang))
        return
    await m.answer(tr("msg_for_whom_table", lang), reply_markup=excel_category_kb(lang))
    await state.set_state(Form.waiting_excel_category)


@dp.message(Form.waiting_excel_category)
async def excel_category(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    cat = WORD_TITLE_TO_CATEGORY.get((m.text or "").strip())
    if not cat:
        await m.answer(tr("msg_pick_category", lang), reply_markup=excel_category_kb(lang))
        return
    await state.update_data(excel_category=cat)
    await m.answer(tr("msg_which_table", lang), reply_markup=excel_kind_kb(cat, lang))
    await state.set_state(Form.waiting_excel_kind)


@dp.message(Form.waiting_excel_kind, F.text.in_(ALL_BACK_CATEGORIES_LABELS))
async def excel_kind_back(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await m.answer(tr("msg_for_whom_table", lang), reply_markup=excel_category_kb(lang))
    await state.set_state(Form.waiting_excel_category)


@dp.message(Form.waiting_excel_kind)
async def excel_kind(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    kind = EXCEL_LABEL_TO_KIND.get((m.text or "").strip())
    if not kind:
        data = await state.get_data()
        cat = data.get("excel_category", "physical")
        await m.answer(tr("msg_pick_table", lang), reply_markup=excel_kind_kb(cat, lang))
        return
    await state.update_data(excel_kind=kind, extra="")
    if kind == "startup_model":
        # Финмодель - данные всегда реальные от пользователя, режим "сгенерировать" тут не предлагаем.
        await m.answer(
            tr("msg_startup_data_prompt", lang),
            reply_markup=cancel_kb(lang)
        )
        await state.set_state(Form.waiting_excel_startup_data)
        return
    await m.answer(
        tr("msg_how_build_table", lang),
        reply_markup=excel_mode_kb(lang)
    )
    await state.set_state(Form.waiting_excel_mode)


@dp.message(Form.waiting_excel_mode, F.text.in_(ALL_BTN_AI_GENERATE_LABELS))
async def excel_mode_ai(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await state.update_data(excel_mode="ai")
    await m.answer(tr("msg_content_lang_prompt", lang), reply_markup=content_lang_kb(lang))
    await state.set_state(Form.waiting_excel_content_lang)


@dp.message(Form.waiting_excel_content_lang)
async def excel_content_lang(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    code = resolve_content_lang((m.text or "").strip(), lang)
    if not code:
        await m.answer(tr("msg_pick_content_lang", lang), reply_markup=content_lang_kb(lang))
        return
    await state.update_data(content_lang=code)
    data = await state.get_data()
    kind = data.get("excel_kind", "expense_estimate")
    await m.answer(excel_kind_hint(kind, "ai", lang), reply_markup=cancel_kb(lang))
    await state.set_state(Form.waiting_excel_topic)


@dp.message(Form.waiting_excel_mode, F.text.in_(ALL_BTN_OWN_DATA_LABELS))
async def excel_mode_user(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    kind = data.get("excel_kind", "expense_estimate")
    await state.update_data(excel_mode="user")
    await m.answer(excel_kind_hint(kind, "user", lang), reply_markup=cancel_kb(lang))
    await state.set_state(Form.waiting_excel_data)


@dp.message(Form.waiting_excel_mode)
async def waiting_excel_mode_fallback(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await m.answer(tr("msg_didnt_understand", lang), reply_markup=excel_mode_kb(lang))


@dp.message(Form.waiting_excel_topic)
async def excel_topic(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await state.update_data(excel_topic=m.text or "")
    await m.answer(tr("msg_build_table_q", lang), reply_markup=excel_confirm_kb(lang))
    await state.set_state(Form.waiting_excel_confirm)


@dp.message(Form.waiting_excel_data)
async def excel_data(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await state.update_data(excel_topic=m.text or "")
    await m.answer(tr("msg_build_table_q", lang), reply_markup=excel_confirm_kb(lang))
    await state.set_state(Form.waiting_excel_confirm)


@dp.message(Form.waiting_excel_startup_data)
async def excel_startup_data(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await state.update_data(excel_topic=m.text or "")
    await m.answer(tr("msg_build_model_q", lang), reply_markup=excel_confirm_kb(lang))
    await state.set_state(Form.waiting_excel_confirm)


@dp.message(Form.waiting_excel_confirm, F.text.in_(ALL_BTN_CHANGE_QUERY_LABELS))
async def excel_change_request(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    kind = data.get("excel_kind", "expense_estimate")
    if kind == "startup_model":
        await m.answer(tr("msg_resend_data", lang), reply_markup=cancel_kb(lang))
        await state.set_state(Form.waiting_excel_startup_data)
    else:
        mode = data.get("excel_mode", "ai")
        hint_key = "ai" if mode == "ai" else "user"
        await m.answer(excel_kind_hint(kind, hint_key, lang), reply_markup=cancel_kb(lang))
        await state.set_state(Form.waiting_excel_topic if mode == "ai" else Form.waiting_excel_data)


@dp.message(Form.waiting_excel_confirm, F.text.in_(ALL_BTN_ADD_INFO_LABELS))
async def excel_add_extra(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await m.answer(tr("msg_what_add", lang), reply_markup=cancel_kb(lang))
    await state.set_state(Form.waiting_excel_extra)


@dp.message(Form.waiting_excel_extra)
async def excel_extra(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    extra = (data.get("extra", "") + "\n" + (m.text or "")).strip()
    await state.update_data(extra=extra)
    await m.answer(tr("msg_build_table_q", lang), reply_markup=excel_confirm_kb(lang))
    await state.set_state(Form.waiting_excel_confirm)


@dp.message(Form.waiting_excel_confirm, F.text.in_(ALL_BTN_BUILD_TABLE_LABELS))
async def excel_build(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    uid = m.from_user.id
    u = get_user(uid)
    ok, reason = start_job(uid)
    if not ok:
        await m.answer(reason)
        return
    await m.answer(tr("msg_building_excel", lang))
    kind = data.get("excel_kind", "expense_estimate")
    topic = data.get("excel_topic", "")
    extra = data.get("extra", "")
    theme_name, colors = pick_theme(topic)
    lang_instr = grok_json_lang_instruction(content_gen_lang(data, lang, mode_key="excel_mode") or "ru")

    try:
        if kind == "startup_model":
            # Только извлечение реальных чисел, которые пользователь уже прислал -
            # модель не имеет права ничего досочинять или менять числа.
            raw = await ask_grok(f"""Извлеки структурированные данные для финансовой модели стартапа
из сообщения пользователя. Числа НИЧЕГО не выдумывай и не досчитывай за пользователя - если какой-то
суммы нет в тексте, ставь 0 (кроме horizon_months - если не указано, ставь 12).
В список expenses включи не только статьи, явно названные пользователем, но и типичные для такого
проекта категории расходов (аренда, зарплаты, маркетинг/реклама, оборудование/техника, хостинг/ИТ,
налоги и обязательные платежи, прочее - выбери те, что подходят по смыслу проекта), даже если сумма
по ним не упомянута - в таком случае ставь amount=0. Это делает список ПОЛНЕЕ, а не выдуманнее: суммы
остаются честными, добавляются только названия недостающих статей. Не меньше 6 строк в expenses.
Для КАЖДОЙ суммы (стартовые вложения, каждая статья расходов, цена за единицу) отдельно определи,
в какой валюте её назвал пользователь - по явному указанию ("долларов", "$", "евро", "€") или по
контексту. Если валюта явно не указана - ставь "RUB". НИКОГДА не конвертируй суммы сам и не пересчитывай
их в рубли - оставляй исходное число ровно как в сообщении пользователя, конвертация делается отдельно
кодом по курсу, а не тобой.
Сообщение пользователя:
{topic}
Доп. уточнения: {extra}
Только JSON:
{{"title":"короткое название проекта по смыслу сообщения","investment":0,"investment_currency":"RUB",
"expenses":[{{"name":"...","amount":0,"currency":"RUB"}}],"price":0,"price_currency":"RUB",
"monthly_sales":0,"horizon_months":12}}{lang_instr}""")
            parsed = extract_json(raw)

            # Конвертация в рубли делается здесь кодом по фиксированному курсу (FX_TO_RUB),
            # а не доверяется модели - иначе есть риск, что она либо забудет сконвертировать
            # (как было раньше: доллары просто вставлялись как есть вместо рублей), либо
            # применит собственный устаревший/придуманный курс. Собираем список реальных
            # конвертаций, чтобы честно показать пользователю, что и по какому курсу пересчитано.
            fx_notes = []

            def _conv(amount, currency, label):
                rate = _xl_fx_rate(currency)
                amount = _xl_num(amount)
                if rate != 1.0 and amount:
                    converted = amount * rate
                    fx_notes.append(f"{label}: {amount:,.0f} {currency} → {converted:,.0f} ₽".replace(",", " "))
                    return converted
                return amount

            parsed["investment"] = _conv(parsed.get("investment"), parsed.get("investment_currency"), "Стартовые вложения")
            parsed["price"] = _conv(parsed.get("price"), parsed.get("price_currency"), "Цена за единицу")
            for e in (parsed.get("expenses") or []):
                e["amount"] = _conv(e.get("amount"), e.get("currency"), e.get("name") or "Расход")
            if fx_notes:
                parsed["_fx_note"] = (
                    f"Курс на {FX_NOTE_DATE} (примерный, уточни актуальный при необходимости): "
                    + "; ".join(fx_notes)
                )

            xlsx_path = f"/tmp/xl_{uid}.xlsx"
            build_excel_startup(xlsx_path, parsed.get("title") or "Финансовая модель", colors, parsed)
            title_for_name = parsed.get("title") or "Финмодель"
        else:
            schema = EXCEL_ROW_SCHEMA[kind]
            min_rows = EXCEL_MIN_ROWS_AI.get(kind, 10)
            ai_mode = data.get("excel_mode") != "user"
            has_category_field = "category" in schema["fields"]
            category_line = ""
            if has_category_field:
                category_line = (
                    " Если для темы естественно разбить позиции на смысловые категории (обычно 3-6 штук, "
                    "например 'Материалы'/'Работа'/'Техника' для сметы или 'Жильё'/'Еда'/'Транспорт' для "
                    "бюджета) - заполни поле category одинаковым коротким названием категории у всех строк, "
                    "которые к ней относятся. Если разбивка на категории неуместна для темы - оставь "
                    "category пустой строкой у всех строк, ничего не выдумывая для галочки."
                )
            if not ai_mode:
                task_line = (
                    "Ниже данные прислал сам пользователь - структурируй их в JSON СТРОГО как есть, "
                    "ничего не досочиняя и не меняя числа. Извлеки КАЖДУЮ отдельную позицию, которую "
                    "упомянул пользователь, включая мелкие и не сгруппированные - не объединяй и не "
                    "сокращай список для краткости, даже если позиций получится много. Если каких-то "
                    "полей не хватает - оставь разумные значения по умолчанию (0 или пустая строка), "
                    "но не выдумывай новые позиции, которых не было в сообщении."
                    + category_line
                )
            else:
                task_line = (
                    f"Пользователь не прислал реальные данные - подбери сам подробные и разнообразные "
                    f"иллюстративные данные по теме (как в примерах для учебных работ). Верни МИНИМУМ "
                    f"{min_rows} строк (в family_budget/project_budget - минимум {min_rows} доходных И "
                    f"минимум {min_rows} расходных строк раздельно, если применимо к теме), реальные "
                    f"по масштабу цифры. Не дублируй формулировки и не используй общие заглушки вроде "
                    f"\"Позиция 1\" - каждая строка должна называть конкретную, отличную от других статью."
                    + category_line
                )
            raw = await ask_grok(f"""Собери содержимое для Excel-таблицы: {EXCEL_KIND_DESC.get(kind)}.
{task_line}
Тема/данные: {topic}
Доп: {extra}
Каждая строка - объект с полями {schema['fields']}, например {schema['example']}.
Только JSON:
{{"title":"короткое название таблицы по теме","rows":[{schema['example']}]}}{lang_instr}""")
            parsed = extract_json(raw)
            rows = parsed.get("rows") or []

            # Если модель всё равно поскупилась на строки в режиме "сгенерировать самому" -
            # переспрашиваем ещё раз с явным указанием, сколько строк получилось и что нужно больше.
            # В режиме "мои данные" ничего не дозапрашиваем - короткая таблица там может быть
            # честным отражением того, что реально прислал пользователь.
            if ai_mode and len(rows) < min_rows:
                raw2 = await ask_grok(f"""Собери содержимое для Excel-таблицы: {EXCEL_KIND_DESC.get(kind)}.
{task_line}
В прошлый раз получилось только {len(rows)} строк - это слишком мало, нужно строго не меньше {min_rows}
строк, разверни тему подробнее, добавь больше конкретных статей/позиций.
Тема/данные: {topic}
Доп: {extra}
Каждая строка - объект с полями {schema['fields']}, например {schema['example']}.
Только JSON:
{{"title":"короткое название таблицы по теме","rows":[{schema['example']}]}}{lang_instr}""")
                parsed2 = extract_json(raw2)
                rows2 = parsed2.get("rows") or []
                if len(rows2) > len(rows):
                    parsed, rows = parsed2, rows2

            if not rows:
                raise ValueError("Модель не вернула строки таблицы")
            xlsx_path = f"/tmp/xl_{uid}.xlsx"
            xl_subtitle = f"Сформировано: {datetime.now().strftime('%d.%m.%Y')}"
            if topic and data.get("excel_mode") != "user":
                topic_short = topic if len(topic) <= 60 else topic[:57] + "..."
                xl_subtitle += f" · Тема: {topic_short}"
            build_excel_items(xlsx_path, parsed.get("title") or EXCEL_KIND_DESC.get(kind, "Таблица"), kind, colors, rows, subtitle=xl_subtitle)
            title_for_name = parsed.get("title") or EXCEL_KIND_DESC.get(kind)

        fname = safe_filename(title_for_name, fallback=EXCEL_KIND_DESC.get(kind, "Таблица"))
        await m.answer_document(FSInputFile(xlsx_path, filename=f"{fname}.xlsx"), caption=tr("msg_excel_caption", lang))
        u["generations"] += 1
        u["history"].append(f"{datetime.now().strftime('%d.%m %H:%M')} — {title_for_name}")
        note_success(uid)
        try:
            os.remove(xlsx_path)
        except Exception:
            pass
        note = tr("msg_table_ready", lang)
        if kind == "startup_model":
            note += tr("msg_check_source_numbers", lang)
            if parsed.get("_fx_note"):
                note += f"\n\n💱 {parsed['_fx_note']}"
        note += "\n\n" + await signature_line(lang)
        await m.answer(note, reply_markup=main_kb(lang))
        await state.clear()
    except Exception as e:
        print("Excel build error:", e)
        await m.answer(tr("msg_build_error_table", lang), reply_markup=main_kb(lang))
        await state.clear()
    finally:
        finish_job(uid)


@dp.message(Form.waiting_excel_confirm)
async def excel_confirm_fallback(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await m.answer(tr("msg_choose_action", lang), reply_markup=excel_confirm_kb(lang))


@dp.message(F.text.in_(ALL_BTN_WORD_LABELS))
async def start_word(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    if not can_generate(m.from_user.id):
        await m.answer(tr("msg_limit", lang))
        return
    await m.answer(tr("msg_for_whom_doc", lang), reply_markup=word_category_kb(lang))
    await state.set_state(Form.waiting_word_category)


@dp.message(Form.waiting_word_category)
async def word_category(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    cat = WORD_TITLE_TO_CATEGORY.get((m.text or "").strip())
    if not cat:
        await m.answer(tr("msg_pick_category", lang), reply_markup=word_category_kb(lang))
        return
    await state.update_data(word_category=cat)
    await m.answer(tr("msg_which_doc", lang), reply_markup=word_kind_kb(cat, lang=lang))
    await state.set_state(Form.waiting_word_kind)


@dp.message(Form.waiting_word_kind, F.text.in_(ALL_MORE_DOCS_LABELS))
async def word_kind_more(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    cat = data.get("word_category", "physical")
    await m.answer(tr("msg_more_docs", lang), reply_markup=word_kind_kb(cat, more=True, lang=lang))


@dp.message(Form.waiting_word_kind, F.text.in_(ALL_BACK_CATEGORIES_LABELS))
async def word_kind_back(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await m.answer(tr("msg_for_whom_doc", lang), reply_markup=word_category_kb(lang))
    await state.set_state(Form.waiting_word_category)


@dp.message(Form.waiting_word_kind)
async def word_kind(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    kind = LABEL_TO_KIND.get((m.text or "").strip())
    if not kind:
        data = await state.get_data()
        cat = data.get("word_category", "physical")
        await m.answer(tr("msg_pick_doc", lang), reply_markup=word_kind_kb(cat, lang=lang))
        return
    await state.update_data(word_kind=kind, extra="", extra_used=0)
    show_template = kind not in STUDY_KINDS
    hint = tr("msg_how_build_doc", lang)
    if show_template:
        hint += tr("msg_template_hint_suffix", lang)
    await m.answer(hint, reply_markup=mode_kb(show_template, lang=lang))
    await state.set_state(Form.waiting_word_mode)



WORD_KIND_HINTS_I18N = {
    "doc": {
        "ai": {"ru": "Напиши тему документа. Например: доклад про историю китов.",
               "en": "Write the document's topic. For example: a report on the history of whales.",
               "de": "Schreibe das Thema des Dokuments. Zum Beispiel: ein Vortrag über die Geschichte der Wale.",
               "ar": "اكتب موضوع المستند. مثال: تقرير عن تاريخ الحيتان.",
               "zh": "请输入文档主题。例如：关于鲸鱼历史的报告。",
               "es": "Escribe el tema del documento. Por ejemplo: un informe sobre la historia de las ballenas.",
               "fr": "Écris le sujet du document. Par exemple : un exposé sur l'histoire des baleines."},
        "user": {"ru": "Пришли текст документа — я поправлю грамотность и оформлю по стандарту.",
                 "en": "Send the document's text — I'll fix the grammar and format it to standard.",
                 "de": "Sende den Text des Dokuments — ich korrigiere die Grammatik und formatiere ihn nach Standard.",
                 "ar": "أرسل نص المستند — سأصحح الأخطاء وأنسقه وفق المعايير.",
                 "zh": "请发送文档文本——我会修正语法并按标准排版。",
                 "es": "Envía el texto del documento — corregiré la gramática y le daré formato estándar.",
                 "fr": "Envoie le texte du document — je corrigerai la grammaire et je le mettrai en forme selon le standard."}},
    "referat": {
        "ai": {"ru": "Напиши тему реферата. Также укажи, если знаешь: учебное заведение, свои ФИО, ФИО преподавателя, город и год — иначе оставлю поля пустыми для заполнения.",
               "en": "Write the topic of the research paper. Also specify, if known: institution, your full name, instructor's full name, city and year — otherwise I'll leave those fields blank to fill in.",
               "de": "Schreibe das Thema des Referats. Gib außerdem an, falls bekannt: Bildungseinrichtung, deinen vollen Namen, den Namen des Dozenten, Stadt und Jahr — sonst lasse ich diese Felder leer zum Ausfüllen.",
               "ar": "اكتب موضوع البحث. حدد أيضاً إن كنت تعرف: المؤسسة التعليمية، اسمك الكامل، اسم الأستاذ، المدينة والسنة — وإلا سأترك هذه الحقول فارغة للتعبئة.",
               "zh": "请输入学术论文主题。如果知道的话，也请注明：学校、你的姓名、指导教师姓名、城市和年份——否则这些字段将留空待填。",
               "es": "Escribe el tema del trabajo de investigación. Indica también, si lo sabes: institución, tu nombre completo, nombre del profesor, ciudad y año — si no, dejaré esos campos en blanco para completar.",
               "fr": "Écris le sujet du dossier de recherche. Indique aussi, si tu les connais : établissement, ton nom complet, nom de l'enseignant, ville et année — sinon je laisserai ces champs vides à compléter."},
        "user": {"ru": "Пришли текст реферата (или тезисы). Также укажи учебное заведение, свои ФИО, ФИО преподавателя, город и год, если нужно их вставить.",
                 "en": "Send the text of the research paper (or an outline). Also specify institution, your full name, instructor's full name, city and year if you need them inserted.",
                 "de": "Sende den Text des Referats (oder Thesen). Gib außerdem Bildungseinrichtung, deinen vollen Namen, den Namen des Dozenten, Stadt und Jahr an, falls diese eingefügt werden sollen.",
                 "ar": "أرسل نص البحث (أو النقاط الرئيسية). حدد أيضاً المؤسسة التعليمية واسمك الكامل واسم الأستاذ والمدينة والسنة إذا كنت تريد إدراجها.",
                 "zh": "请发送论文文本（或提纲）。如需插入，请同时注明学校、你的姓名、指导教师姓名、城市和年份。",
                 "es": "Envía el texto del trabajo (o las ideas principales). Indica también institución, tu nombre completo, nombre del profesor, ciudad y año si necesitas incluirlos.",
                 "fr": "Envoie le texte du dossier (ou les grandes lignes). Indique aussi l'établissement, ton nom complet, le nom de l'enseignant, la ville et l'année si tu veux les insérer."}},
    "report": {
        "ai": {"ru": "Напиши тему доклада. Также укажи, если знаешь: учебное заведение, свои ФИО, класс/курс — иначе оставлю поля пустыми для заполнения.",
               "en": "Write the topic of the report. Also specify, if known: institution, your full name, grade/year — otherwise I'll leave those fields blank to fill in.",
               "de": "Schreibe das Thema des Vortrags. Gib außerdem an, falls bekannt: Bildungseinrichtung, deinen vollen Namen, Klasse/Kurs — sonst lasse ich diese Felder leer zum Ausfüllen.",
               "ar": "اكتب موضوع التقرير. حدد أيضاً إن كنت تعرف: المؤسسة التعليمية، اسمك الكامل، الصف/السنة الدراسية — وإلا سأترك هذه الحقول فارغة للتعبئة.",
               "zh": "请输入报告主题。如果知道的话，也请注明：学校、你的姓名、年级——否则这些字段将留空待填。",
               "es": "Escribe el tema de la exposición. Indica también, si lo sabes: institución, tu nombre completo, curso/grado — si no, dejaré esos campos en blanco para completar.",
               "fr": "Écris le sujet de l'exposé. Indique aussi, si tu les connais : établissement, ton nom complet, classe/année — sinon je laisserai ces champs vides à compléter."},
        "user": {"ru": "Пришли текст доклада (или тезисы). Также укажи учебное заведение, свои ФИО и класс/курс, если нужно их вставить.",
                 "en": "Send the text of the report (or an outline). Also specify institution, your full name and grade/year if you need them inserted.",
                 "de": "Sende den Text des Vortrags (oder Thesen). Gib außerdem Bildungseinrichtung, deinen vollen Namen und Klasse/Kurs an, falls diese eingefügt werden sollen.",
                 "ar": "أرسل نص التقرير (أو النقاط الرئيسية). حدد أيضاً المؤسسة التعليمية واسمك الكامل والصف/السنة الدراسية إذا كنت تريد إدراجها.",
                 "zh": "请发送报告文本（或提纲）。如需插入，请同时注明学校、你的姓名和年级。",
                 "es": "Envía el texto de la exposición (o las ideas principales). Indica también institución, tu nombre completo y curso/grado si necesitas incluirlos.",
                 "fr": "Envoie le texte de l'exposé (ou les grandes lignes). Indique aussi l'établissement, ton nom complet et ta classe/année si tu veux les insérer."}},
    "essay": {
        "ai": {"ru": "Напиши тему эссе и свою позицию по ней, если она уже есть. Укажи, если знаешь: учебное заведение, свои ФИО.",
               "en": "Write the essay topic and your stance on it, if you already have one. Specify, if known: institution, your full name.",
               "de": "Schreibe das Thema des Essays und deine Position dazu, falls du bereits eine hast. Gib an, falls bekannt: Bildungseinrichtung, deinen vollen Namen.",
               "ar": "اكتب موضوع المقال وموقفك منه إن كان لديك بالفعل. حدد إن كنت تعرف: المؤسسة التعليمية، اسمك الكامل.",
               "zh": "请输入论文主题，以及你对该主题的观点（如果已有）。如知道的话，也请注明学校和你的姓名。",
               "es": "Escribe el tema del ensayo y tu postura al respecto, si ya la tienes. Indica, si lo sabes: institución, tu nombre completo.",
               "fr": "Écris le sujet de l'essai et ta position sur celui-ci, si tu en as déjà une. Indique, si tu les connais : établissement, ton nom complet."},
        "user": {"ru": "Пришли текст эссе (или тезисы, свои мысли по теме). Укажи учебное заведение и свои ФИО, если нужно их вставить.",
                 "en": "Send the essay text (or an outline, your thoughts on the topic). Specify institution and your full name if you need them inserted.",
                 "de": "Sende den Text des Essays (oder Thesen, deine Gedanken zum Thema). Gib Bildungseinrichtung und deinen vollen Namen an, falls diese eingefügt werden sollen.",
                 "ar": "أرسل نص المقال (أو النقاط الرئيسية، أفكارك حول الموضوع). حدد المؤسسة التعليمية واسمك الكامل إذا كنت تريد إدراجهما.",
                 "zh": "请发送论文文本（或提纲、你对主题的想法）。如需插入，请注明学校和你的姓名。",
                 "es": "Envía el texto del ensayo (o las ideas principales, tus reflexiones sobre el tema). Indica institución y tu nombre completo si necesitas incluirlos.",
                 "fr": "Envoie le texte de l'essai (ou les grandes lignes, tes réflexions sur le sujet). Indique l'établissement et ton nom complet si tu veux les insérer."}},
    "notes": {
        "ai": {"ru": "Напиши тему или раздел, по которому нужен конспект (например: конспект лекции по клеточной биологии).",
               "en": "Write the topic or section you need notes on (for example: notes on a cell biology lecture).",
               "de": "Schreibe das Thema oder den Abschnitt, für den du Notizen brauchst (zum Beispiel: Mitschrift einer Vorlesung zur Zellbiologie).",
               "ar": "اكتب الموضوع أو القسم الذي تحتاج ملخصاً له (مثال: ملخص محاضرة عن بيولوجيا الخلية).",
               "zh": "请输入需要做笔记的主题或章节（例如：细胞生物学讲座笔记）。",
               "es": "Escribe el tema o la sección de la que necesitas apuntes (por ejemplo: apuntes de una clase de biología celular).",
               "fr": "Écris le sujet ou la section pour lequel tu as besoin de notes (par exemple : notes d'un cours sur la biologie cellulaire)."},
        "user": {"ru": "Пришли материал, по которому нужно сделать конспект — я структурирую его в сжатом виде по пунктам.",
                 "en": "Send the material to make notes from — I'll structure it concisely into points.",
                 "de": "Sende das Material, aus dem Notizen erstellt werden sollen — ich strukturiere es kompakt in Stichpunkten.",
                 "ar": "أرسل المادة التي تريد تلخيصها — سأنظمها بإيجاز في نقاط.",
                 "zh": "请发送需要整理笔记的材料——我会将其精炼为要点结构。",
                 "es": "Envía el material del que hay que hacer apuntes — lo estructuraré de forma concisa en puntos.",
                 "fr": "Envoie le contenu à partir duquel faire des notes — je le structurerai de façon concise en points."}},
    "coursework": {
        "ai": {"ru": "Напиши тему курсовой работы. Укажи, если знаешь: учебное заведение, специальность, свои ФИО, ФИО научного руководителя, город и год.",
               "en": "Write the topic of the coursework. Specify, if known: institution, major, your full name, supervisor's full name, city and year.",
               "de": "Schreibe das Thema der Kursarbeit. Gib an, falls bekannt: Bildungseinrichtung, Studienfach, deinen vollen Namen, den Namen des Betreuers, Stadt und Jahr.",
               "ar": "اكتب موضوع البحث الفصلي. حدد إن كنت تعرف: المؤسسة التعليمية، التخصص، اسمك الكامل، اسم المشرف، المدينة والسنة.",
               "zh": "请输入课程论文主题。如知道的话，也请注明学校、专业、你的姓名、导师姓名、城市和年份。",
               "es": "Escribe el tema del trabajo de curso. Indica, si lo sabes: institución, carrera, tu nombre completo, nombre del director, ciudad y año.",
               "fr": "Écris le sujet du mémoire. Indique, si tu les connais : établissement, filière, ton nom complet, nom du directeur, ville et année."},
        "user": {"ru": "Пришли текст или тезисы курсовой работы. Укажи учебное заведение, специальность, свои ФИО, ФИО научного руководителя, город и год, если нужно их вставить.",
                 "en": "Send the text or outline of the coursework. Specify institution, major, your full name, supervisor's full name, city and year if you need them inserted.",
                 "de": "Sende den Text oder die Thesen der Kursarbeit. Gib Bildungseinrichtung, Studienfach, deinen vollen Namen, den Namen des Betreuers, Stadt und Jahr an, falls diese eingefügt werden sollen.",
                 "ar": "أرسل نص أو نقاط البحث الفصلي. حدد المؤسسة التعليمية والتخصص واسمك الكامل واسم المشرف والمدينة والسنة إذا كنت تريد إدراجها.",
                 "zh": "请发送课程论文文本或提纲。如需插入，请注明学校、专业、你的姓名、导师姓名、城市和年份。",
                 "es": "Envía el texto o las ideas principales del trabajo de curso. Indica institución, carrera, tu nombre completo, nombre del director, ciudad y año si necesitas incluirlos.",
                 "fr": "Envoie le texte ou les grandes lignes du mémoire. Indique l'établissement, la filière, ton nom complet, le nom du directeur, la ville et l'année si tu veux les insérer."}},
    "dkp": {
        "ai": {"ru": "Напиши, что покупается/продаётся (например: продажа автомобиля). Укажи, если знаешь: ФИО и паспортные данные продавца и покупателя, точное описание вещи, цену, порядок оплаты, дату и место.",
               "en": "Write what's being bought/sold (for example: sale of a car). Specify, if known: full names and ID details of seller and buyer, exact description of the item, price, payment terms, date and place.",
               "de": "Schreibe, was gekauft/verkauft wird (zum Beispiel: Verkauf eines Autos). Gib an, falls bekannt: Namen und Ausweisdaten von Verkäufer und Käufer, genaue Beschreibung der Sache, Preis, Zahlungsbedingungen, Datum und Ort.",
               "ar": "اكتب ما يُشترى/يُباع (مثال: بيع سيارة). حدد إن كنت تعرف: الاسم الكامل وبيانات الهوية للبائع والمشتري، وصفاً دقيقاً للشيء، السعر، شروط الدفع، التاريخ والمكان.",
               "zh": "请说明买卖的是什么（例如：出售汽车）。如知道的话，请注明卖方和买方的姓名及证件信息、物品的准确描述、价格、付款方式、日期和地点。",
               "es": "Escribe qué se compra/vende (por ejemplo: venta de un coche). Indica, si lo sabes: nombre completo y datos de identificación del vendedor y comprador, descripción exacta del bien, precio, forma de pago, fecha y lugar.",
               "fr": "Écris ce qui est acheté/vendu (par exemple : vente d'une voiture). Indique, si tu les connais : noms complets et pièces d'identité du vendeur et de l'acheteur, description exacte du bien, prix, modalités de paiement, date et lieu."},
        "user": {"ru": "Пришли данные для договора купли-продажи: ФИО и паспортные данные продавца и покупателя, описание предмета продажи, цену, порядок оплаты, дату и место составления.",
                 "en": "Send the details for the sale agreement: full names and ID details of seller and buyer, description of the item, price, payment terms, date and place of drafting.",
                 "de": "Sende die Daten für den Kaufvertrag: Namen und Ausweisdaten von Verkäufer und Käufer, Beschreibung des Kaufgegenstands, Preis, Zahlungsbedingungen, Datum und Ort der Erstellung.",
                 "ar": "أرسل بيانات عقد البيع: الاسم الكامل وبيانات الهوية للبائع والمشتري، وصف الشيء المباع، السعر، شروط الدفع، تاريخ ومكان التحرير.",
                 "zh": "请发送买卖合同所需信息：卖方和买方的姓名及证件信息、销售物品描述、价格、付款方式、签订日期和地点。",
                 "es": "Envía los datos para el contrato de compraventa: nombre completo y datos de identificación del vendedor y comprador, descripción del bien vendido, precio, forma de pago, fecha y lugar de redacción.",
                 "fr": "Envoie les données du contrat de vente : noms complets et pièces d'identité du vendeur et de l'acheteur, description du bien vendu, prix, modalités de paiement, date et lieu de rédaction."}},
    "rent": {
        "ai": {"ru": "Напиши, что сдаётся в аренду (например: аренда квартиры). Укажи, если знаешь: данные сторон, точный адрес/описание объекта, срок аренды, сумму и порядок оплаты.",
               "en": "Write what's being rented out (for example: apartment rental). Specify, if known: parties' details, exact address/description of the property, lease term, amount and payment terms.",
               "de": "Schreibe, was vermietet wird (zum Beispiel: Vermietung einer Wohnung). Gib an, falls bekannt: Daten der Parteien, genaue Adresse/Beschreibung des Objekts, Mietdauer, Betrag und Zahlungsbedingungen.",
               "ar": "اكتب ما الذي يُؤجَّر (مثال: تأجير شقة). حدد إن كنت تعرف: بيانات الطرفين، العنوان الدقيق/وصف العقار، مدة الإيجار، المبلغ وشروط الدفع.",
               "zh": "请说明出租的是什么（例如：出租公寓）。如知道的话，请注明双方信息、准确地址/物业描述、租期、金额和付款方式。",
               "es": "Escribe qué se alquila (por ejemplo: alquiler de un piso). Indica, si lo sabes: datos de las partes, dirección exacta/descripción del inmueble, plazo del alquiler, importe y forma de pago.",
               "fr": "Écris ce qui est loué (par exemple : location d'un appartement). Indique, si tu les connais : coordonnées des parties, adresse exacte/description du bien, durée du bail, montant et modalités de paiement."},
        "user": {"ru": "Пришли данные для договора аренды: ФИО сторон, точный адрес или описание объекта аренды, срок, сумму и порядок оплаты, кто оплачивает коммунальные услуги/ремонт.",
                 "en": "Send the details for the lease agreement: full names of the parties, exact address or description of the rented property, term, amount and payment terms, who pays for utilities/repairs.",
                 "de": "Sende die Daten für den Mietvertrag: Namen der Parteien, genaue Adresse oder Beschreibung des Mietobjekts, Mietdauer, Betrag und Zahlungsbedingungen, wer für Nebenkosten/Reparaturen zahlt.",
                 "ar": "أرسل بيانات عقد الإيجار: الاسم الكامل للطرفين، العنوان الدقيق أو وصف العقار المؤجر، المدة، المبلغ وشروط الدفع، من يدفع الخدمات/الصيانة.",
                 "zh": "请发送租赁合同所需信息：双方姓名、租赁物业的准确地址或描述、租期、金额和付款方式、由谁支付水电费/维修费。",
                 "es": "Envía los datos para el contrato de alquiler: nombre completo de las partes, dirección exacta o descripción del inmueble, plazo, importe y forma de pago, quién paga los suministros/reparaciones.",
                 "fr": "Envoie les données du bail : noms complets des parties, adresse exacte ou description du bien loué, durée, montant et modalités de paiement, qui paie les charges/réparations."}},
    "offer": {
        "ai": {"ru": "Напиши, что за услугу/товар предлагаете и кому. Укажи, если знаешь: название компании, суть предложения, цену и условия, срок действия КП.",
               "en": "Write what service/product you're offering and to whom. Specify, if known: company name, essence of the offer, price and terms, validity period of the proposal.",
               "de": "Schreibe, welche Dienstleistung/welches Produkt du anbietest und wem. Gib an, falls bekannt: Firmenname, Kern des Angebots, Preis und Bedingungen, Gültigkeitsdauer des Angebots.",
               "ar": "اكتب ما هي الخدمة/المنتج الذي تقدمه ولمن. حدد إن كنت تعرف: اسم الشركة، جوهر العرض، السعر والشروط، مدة صلاحية العرض.",
               "zh": "请说明你提供的服务/产品是什么，面向谁。如知道的话，请注明公司名称、提案要点、价格和条件、报价有效期。",
               "es": "Escribe qué servicio/producto ofreces y a quién. Indica, si lo sabes: nombre de la empresa, esencia de la propuesta, precio y condiciones, plazo de validez de la oferta.",
               "fr": "Écris quel service/produit tu proposes et à qui. Indique, si tu les connais : nom de l'entreprise, essence de l'offre, prix et conditions, durée de validité de la proposition."},
        "user": {"ru": "Пришли данные для коммерческого предложения: название компании, суть предложения, цену и условия оплаты, сроки, контакты.",
                 "en": "Send the details for the commercial proposal: company name, essence of the offer, price and payment terms, deadlines, contacts.",
                 "de": "Sende die Daten für das Angebot: Firmenname, Kern des Angebots, Preis und Zahlungsbedingungen, Fristen, Kontaktdaten.",
                 "ar": "أرسل بيانات العرض التجاري: اسم الشركة، جوهر العرض، السعر وشروط الدفع، المواعيد، بيانات التواصل.",
                 "zh": "请发送商业提案所需信息：公司名称、提案要点、价格和付款条件、期限、联系方式。",
                 "es": "Envía los datos para la propuesta comercial: nombre de la empresa, esencia de la oferta, precio y condiciones de pago, plazos, contactos.",
                 "fr": "Envoie les données de l'offre commerciale : nom de l'entreprise, essence de l'offre, prix et modalités de paiement, délais, contacts."}},
    "act": {
        "ai": {"ru": "Напиши, что передаётся и по какому договору (например: акт передачи оборудования по договору №12). Укажи, если знаешь: кто передаёт, кто принимает, перечень предметов.",
               "en": "Write what's being handed over and under which agreement (for example: equipment handover act under contract No. 12). Specify, if known: who's handing over, who's receiving, list of items.",
               "de": "Schreibe, was übergeben wird und aufgrund welchen Vertrags (zum Beispiel: Übergabeprotokoll für Geräte laut Vertrag Nr. 12). Gib an, falls bekannt: wer übergibt, wer empfängt, Liste der Gegenstände.",
               "ar": "اكتب ما الذي يُسلَّم وبموجب أي عقد (مثال: محضر تسليم معدات بموجب العقد رقم 12). حدد إن كنت تعرف: من يسلّم، من يستلم، قائمة الأصناف.",
               "zh": "请说明交接的内容及依据的合同（例如：根据12号合同的设备交接单）。如知道的话，请注明交付方、接收方、物品清单。",
               "es": "Escribe qué se entrega y en virtud de qué contrato (por ejemplo: acta de entrega de equipo según el contrato n.º 12). Indica, si lo sabes: quién entrega, quién recibe, listado de objetos.",
               "fr": "Écris ce qui est remis et en vertu de quel contrat (par exemple : procès-verbal de remise de matériel selon le contrat n° 12). Indique, si tu les connais : qui remet, qui reçoit, liste des objets."},
        "user": {"ru": "Пришли данные для акта приёма-передачи: номер и дату основного договора, кто передаёт и кто принимает (ФИО/организация), перечень передаваемого с количеством и состоянием.",
                 "en": "Send the details for the handover act: number and date of the main agreement, who's handing over and who's receiving (name/organization), list of items with quantity and condition.",
                 "de": "Sende die Daten für das Übergabeprotokoll: Nummer und Datum des Hauptvertrags, wer übergibt und wer empfängt (Name/Organisation), Liste der übergebenen Gegenstände mit Menge und Zustand.",
                 "ar": "أرسل بيانات محضر التسليم والاستلام: رقم وتاريخ العقد الأساسي، من يسلّم ومن يستلم (الاسم/المنظمة)، قائمة الأصناف المسلَّمة مع الكمية والحالة.",
                 "zh": "请发送交接单所需信息：主合同编号和日期、交付方和接收方（姓名/单位）、交接物品清单（含数量和状态）。",
                 "es": "Envía los datos para el acta de entrega-recepción: número y fecha del contrato principal, quién entrega y quién recibe (nombre/organización), listado de lo entregado con cantidad y estado.",
                 "fr": "Envoie les données du procès-verbal de remise : numéro et date du contrat principal, qui remet et qui reçoit (nom/organisation), liste des objets remis avec quantité et état."}},
    "statement": {
        "ai": {"ru": "Напиши суть заявления и кому оно адресовано.",
               "en": "Write the essence of the statement and to whom it's addressed.",
               "de": "Schreibe das Anliegen des Antrags und an wen er gerichtet ist.",
               "ar": "اكتب جوهر الطلب ولمن هو موجّه.",
               "zh": "请说明申请的内容以及收件人。",
               "es": "Escribe el motivo de la solicitud y a quién va dirigida.",
               "fr": "Écris l'objet de la demande et à qui elle est adressée."},
        "user": {"ru": "Пришли текст заявления: кому адресовано (должность, ФИО), от кого, суть просьбы или требования, дата.",
                 "en": "Send the text of the statement: who it's addressed to (position, full name), from whom, essence of the request or demand, date.",
                 "de": "Sende den Text des Antrags: an wen gerichtet (Position, Name), von wem, Anliegen der Bitte oder Forderung, Datum.",
                 "ar": "أرسل نص الطلب: الموجّه إليه (المنصب، الاسم الكامل)، من مقدّمه، جوهر الطلب أو المطلب، التاريخ.",
                 "zh": "请发送申请文本：收件人（职位、姓名）、申请人、请求或要求的内容、日期。",
                 "es": "Envía el texto de la solicitud: a quién va dirigida (cargo, nombre completo), de quién, motivo de la petición o exigencia, fecha.",
                 "fr": "Envoie le texte de la demande : à qui elle est adressée (fonction, nom complet), de qui, objet de la demande, date."}},
    "proxy": {
        "ai": {"ru": "Напиши, на какие действия нужна доверенность. Укажи, если знаешь: ФИО и паспортные данные доверителя и представителя, срок действия.",
               "en": "Write what actions the power of attorney is for. Specify, if known: full names and ID details of the principal and the representative, validity period.",
               "de": "Schreibe, für welche Handlungen die Vollmacht benötigt wird. Gib an, falls bekannt: Namen und Ausweisdaten von Vollmachtgeber und Bevollmächtigtem, Gültigkeitsdauer.",
               "ar": "اكتب لأي إجراءات تحتاج التوكيل. حدد إن كنت تعرف: الاسم الكامل وبيانات الهوية للموكِّل والوكيل، مدة الصلاحية.",
               "zh": "请说明委托书用于哪些事项。如知道的话，请注明委托人和受托人的姓名及证件信息、有效期。",
               "es": "Escribe para qué acciones se necesita el poder. Indica, si lo sabes: nombre completo y datos de identificación del poderdante y el apoderado, plazo de validez.",
               "fr": "Écris pour quelles actions la procuration est nécessaire. Indique, si tu les connais : noms complets et pièces d'identité du mandant et du mandataire, durée de validité."},
        "user": {"ru": "Пришли данные для доверенности: ФИО и паспортные данные доверителя и представителя, точный перечень полномочий, срок действия, дату составления.",
                 "en": "Send the details for the power of attorney: full names and ID details of the principal and the representative, exact list of powers, validity period, date of drafting.",
                 "de": "Sende die Daten für die Vollmacht: Namen und Ausweisdaten von Vollmachtgeber und Bevollmächtigtem, genaue Liste der Befugnisse, Gültigkeitsdauer, Ausstellungsdatum.",
                 "ar": "أرسل بيانات التوكيل: الاسم الكامل وبيانات الهوية للموكِّل والوكيل، قائمة دقيقة بالصلاحيات، مدة الصلاحية، تاريخ التحرير.",
                 "zh": "请发送委托书所需信息：委托人和受托人的姓名及证件信息、准确的权限清单、有效期、签订日期。",
                 "es": "Envía los datos para el poder: nombre completo y datos de identificación del poderdante y el apoderado, listado exacto de facultades, plazo de validez, fecha de redacción.",
                 "fr": "Envoie les données de la procuration : noms complets et pièces d'identité du mandant et du mandataire, liste exacte des pouvoirs, durée de validité, date de rédaction."}},
    "loan": {
        "ai": {"ru": "Напиши сумму и условия займа. Укажи, если знаешь: ФИО сторон, срок возврата, проценты.",
               "en": "Write the amount and terms of the loan. Specify, if known: full names of the parties, repayment term, interest.",
               "de": "Schreibe den Betrag und die Bedingungen des Darlehens. Gib an, falls bekannt: Namen der Parteien, Rückzahlungsfrist, Zinsen.",
               "ar": "اكتب مبلغ وشروط القرض. حدد إن كنت تعرف: الاسم الكامل للطرفين، مدة السداد، الفائدة.",
               "zh": "请说明借款金额和条件。如知道的话，请注明双方姓名、还款期限、利息。",
               "es": "Escribe el importe y las condiciones del préstamo. Indica, si lo sabes: nombre completo de las partes, plazo de devolución, intereses.",
               "fr": "Écris le montant et les conditions du prêt. Indique, si tu les connais : noms complets des parties, délai de remboursement, intérêts."},
        "user": {"ru": "Пришли данные для расписки/договора займа: ФИО и паспортные данные заимодавца и заёмщика, сумму, срок возврата, проценты (если есть).",
                 "en": "Send the details for the promissory note/loan agreement: full names and ID details of the lender and the borrower, amount, repayment term, interest (if any).",
                 "de": "Sende die Daten für die Schuldschein/den Darlehensvertrag: Namen und Ausweisdaten von Darlehensgeber und Darlehensnehmer, Betrag, Rückzahlungsfrist, Zinsen (falls vorhanden).",
                 "ar": "أرسل بيانات سند الدين/عقد القرض: الاسم الكامل وبيانات الهوية للمُقرض والمقترض، المبلغ، مدة السداد، الفائدة (إن وجدت).",
                 "zh": "请发送借条/借款合同所需信息：出借人和借款人的姓名及证件信息、金额、还款期限、利息（如有）。",
                 "es": "Envía los datos para el pagaré/contrato de préstamo: nombre completo y datos de identificación del prestamista y el prestatario, importe, plazo de devolución, intereses (si los hay).",
                 "fr": "Envoie les données de la reconnaissance de dette/contrat de prêt : noms complets et pièces d'identité du prêteur et de l'emprunteur, montant, délai de remboursement, intérêts (le cas échéant)."}},
    "claim": {
        "ai": {"ru": "Напиши суть претензии и к кому она адресована. Укажи, если знаешь: от кого претензия, основание (договор/факт), требование.",
               "en": "Write the essence of the claim and to whom it's addressed. Specify, if known: from whom the claim is, the basis (contract/fact), the demand.",
               "de": "Schreibe das Anliegen der Beschwerde und an wen sie gerichtet ist. Gib an, falls bekannt: von wem die Beschwerde stammt, die Grundlage (Vertrag/Sachverhalt), die Forderung.",
               "ar": "اكتب جوهر المطالبة ولمن هي موجّهة. حدد إن كنت تعرف: مقدّم المطالبة، الأساس (عقد/واقعة)، المطلب.",
               "zh": "请说明索赔的内容以及收件人。如知道的话，请注明索赔方、依据（合同/事实）、诉求。",
               "es": "Escribe el motivo de la reclamación y a quién va dirigida. Indica, si lo sabes: quién la presenta, el fundamento (contrato/hecho), la exigencia.",
               "fr": "Écris l'objet de la réclamation et à qui elle est adressée. Indique, si tu les connais : de qui vient la réclamation, le fondement (contrat/fait), la demande."},
        "user": {"ru": "Пришли данные для претензии: кому адресована (ФИО/организация), от кого, суть нарушения, конкретное требование, срок ответа.",
                 "en": "Send the details for the claim: who it's addressed to (name/organization), from whom, essence of the violation, specific demand, response deadline.",
                 "de": "Sende die Daten für die Beschwerde: an wen gerichtet (Name/Organisation), von wem, Art der Verletzung, konkrete Forderung, Antwortfrist.",
                 "ar": "أرسل بيانات المطالبة: الموجّهة إليه (الاسم/المنظمة)، مقدّمها، جوهر المخالفة، المطلب المحدد، مهلة الرد.",
                 "zh": "请发送索赔函所需信息：收件人（姓名/单位）、索赔方、违约内容、具体诉求、回复期限。",
                 "es": "Envía los datos para la reclamación: a quién va dirigida (nombre/organización), de quién, motivo del incumplimiento, exigencia concreta, plazo de respuesta.",
                 "fr": "Envoie les données de la réclamation : à qui elle est adressée (nom/organisation), de qui, nature du manquement, demande précise, délai de réponse."}},
    "consent": {
        "ai": {"ru": "Напиши, кто едет и куда (например: согласие на выезд сына в Турцию). Укажи, если знаешь: ФИО и данные родителя, ФИО и дату рождения ребёнка, с кем и куда он выезжает, срок действия согласия.",
               "en": "Write who is traveling and where (for example: consent for a son's trip to Turkey). Specify, if known: parent's full name and details, child's full name and date of birth, with whom and where the child is traveling, validity period of the consent.",
               "de": "Schreibe, wer wohin reist (zum Beispiel: Zustimmung zur Ausreise des Sohnes in die Türkei). Gib an, falls bekannt: Name und Daten des Elternteils, Name und Geburtsdatum des Kindes, mit wem und wohin es reist, Gültigkeitsdauer der Zustimmung.",
               "ar": "اكتب من يسافر وإلى أين (مثال: موافقة على سفر الابن إلى تركيا). حدد إن كنت تعرف: اسم وبيانات ولي الأمر، اسم وتاريخ ميلاد الطفل، مع من وإلى أين يسافر، مدة صلاحية الموافقة.",
               "zh": "请说明谁去哪里旅行（例如：同意儿子前往土耳其）。如知道的话，请注明家长姓名及信息、孩子姓名及出生日期、与谁一同前往何地、同意书有效期。",
               "es": "Escribe quién viaja y adónde (por ejemplo: consentimiento para el viaje del hijo a Turquía). Indica, si lo sabes: nombre y datos del progenitor, nombre y fecha de nacimiento del menor, con quién y adónde viaja, plazo de validez del consentimiento.",
               "fr": "Écris qui voyage et où (par exemple : autorisation de voyage du fils en Turquie). Indique, si tu les connais : nom et coordonnées du parent, nom et date de naissance de l'enfant, avec qui et où il voyage, durée de validité de l'autorisation."},
        "user": {"ru": "Пришли данные для согласия на выезд ребёнка: ФИО родителя (доверителя), ФИО и дату рождения ребёнка, с кем именно и в какую страну выезжает ребёнок, на какой срок.",
                 "en": "Send the details for the child travel consent: parent's (principal's) full name, child's full name and date of birth, with whom exactly and to which country the child is traveling, for how long.",
                 "de": "Sende die Daten für die Zustimmung zur Ausreise des Kindes: Name des Elternteils (Vollmachtgebers), Name und Geburtsdatum des Kindes, mit wem genau und in welches Land es reist, für welchen Zeitraum.",
                 "ar": "أرسل بيانات موافقة سفر الطفل: اسم ولي الأمر (الموكِّل)، اسم وتاريخ ميلاد الطفل، مع من بالضبط وإلى أي بلد يسافر الطفل، ولأي مدة.",
                 "zh": "请发送儿童出境同意书所需信息：家长（委托人）姓名、孩子姓名及出生日期、具体与谁及前往哪个国家、期限多久。",
                 "es": "Envía los datos para el consentimiento de viaje del menor: nombre completo del progenitor (otorgante), nombre y fecha de nacimiento del menor, con quién exactamente y a qué país viaja, por cuánto tiempo.",
                 "fr": "Envoie les données de l'autorisation de voyage de l'enfant : nom complet du parent (mandant), nom et date de naissance de l'enfant, avec qui exactement et vers quel pays il voyage, pour quelle durée."}},
    "marriage_contract": {
        "ai": {"ru": "Напиши, какие имущественные вопросы супруги хотят закрепить в брачном договоре. Укажи, если знаешь: ФИО супругов, режим имущества (совместное/раздельное), конкретное имущество.",
               "en": "Write which property matters the spouses want to fix in the prenuptial agreement. Specify, if known: spouses' full names, property regime (joint/separate), specific property.",
               "de": "Schreibe, welche Vermögensfragen die Ehepartner im Ehevertrag regeln möchten. Gib an, falls bekannt: Namen der Ehepartner, Vermögensregelung (gemeinsam/getrennt), konkretes Vermögen.",
               "ar": "اكتب المسائل المالية التي يريد الزوجان تثبيتها في عقد الزواج. حدد إن كنت تعرف: الاسم الكامل للزوجين، نظام الملكية (مشتركة/منفصلة)، الممتلكات المحددة.",
               "zh": "请说明夫妻双方希望在婚前协议中约定的财产事项。如知道的话，请注明夫妻姓名、财产制度（共同/分别）、具体财产。",
               "es": "Escribe qué cuestiones patrimoniales quieren fijar los cónyuges en las capitulaciones matrimoniales. Indica, si lo sabes: nombre completo de los cónyuges, régimen de bienes (conjunto/separado), bienes concretos.",
               "fr": "Écris quelles questions patrimoniales les époux veulent fixer dans le contrat de mariage. Indique, si tu les connais : noms complets des époux, régime matrimonial (commun/séparé), biens précis."},
        "user": {"ru": "Пришли данные для брачного договора: ФИО обоих супругов, реквизиты свидетельства о браке, какое имущество и на каких условиях делится.",
                 "en": "Send the details for the prenuptial agreement: full names of both spouses, marriage certificate details, what property and on what terms is divided.",
                 "de": "Sende die Daten für den Ehevertrag: Namen beider Ehepartner, Angaben zur Heiratsurkunde, welches Vermögen und unter welchen Bedingungen aufgeteilt wird.",
                 "ar": "أرسل بيانات عقد الزواج: الاسم الكامل لكلا الزوجين، بيانات شهادة الزواج، الممتلكات وشروط تقسيمها.",
                 "zh": "请发送婚前协议所需信息：夫妻双方姓名、结婚证信息、财产及分配条件。",
                 "es": "Envía los datos para las capitulaciones matrimoniales: nombre completo de ambos cónyuges, datos del certificado de matrimonio, qué bienes y en qué condiciones se reparten.",
                 "fr": "Envoie les données du contrat de mariage : noms complets des deux époux, informations de l'acte de mariage, quels biens et selon quelles conditions sont partagés."}},
    "gift": {
        "ai": {"ru": "Напиши, что и кому дарится (например: дарение квартиры сыну). Укажи, если знаешь: ФИО дарителя и одаряемого, точное описание предмета дарения.",
               "en": "Write what is given to whom (for example: gifting an apartment to a son). Specify, if known: full names of the donor and recipient, exact description of the gift.",
               "de": "Schreibe, was wem geschenkt wird (zum Beispiel: Schenkung einer Wohnung an den Sohn). Gib an, falls bekannt: Namen von Schenker und Beschenktem, genaue Beschreibung des Geschenks.",
               "ar": "اكتب ما الذي يُهدى ولمن (مثال: هبة شقة للابن). حدد إن كنت تعرف: الاسم الكامل للواهب والموهوب له، وصفاً دقيقاً للهبة.",
               "zh": "请说明赠与的对象和内容（例如：将公寓赠与儿子）。如知道的话，请注明赠与人和受赠人姓名、赠与物的准确描述。",
               "es": "Escribe qué se dona y a quién (por ejemplo: donación de un piso al hijo). Indica, si lo sabes: nombre completo del donante y del donatario, descripción exacta del bien donado.",
               "fr": "Écris ce qui est donné à qui (par exemple : donation d'un appartement au fils). Indique, si tu les connais : noms complets du donateur et du donataire, description exacte du bien donné."},
        "user": {"ru": "Пришли данные для договора дарения: ФИО и паспортные данные дарителя и одаряемого, точное описание предмета дарения, документы-основания (если есть).",
                 "en": "Send the details for the gift agreement: full names and ID details of the donor and recipient, exact description of the gift, supporting documents (if any).",
                 "de": "Sende die Daten für den Schenkungsvertrag: Namen und Ausweisdaten von Schenker und Beschenktem, genaue Beschreibung des Geschenks, Grundlagendokumente (falls vorhanden).",
                 "ar": "أرسل بيانات عقد الهبة: الاسم الكامل وبيانات الهوية للواهب والموهوب له، وصفاً دقيقاً للهبة، المستندات الداعمة (إن وجدت).",
                 "zh": "请发送赠与合同所需信息：赠与人和受赠人的姓名及证件信息、赠与物的准确描述、相关证明文件（如有）。",
                 "es": "Envía los datos para el contrato de donación: nombre completo y datos de identificación del donante y el donatario, descripción exacta del bien donado, documentos de respaldo (si los hay).",
                 "fr": "Envoie les données de l'acte de donation : noms complets et pièces d'identité du donateur et du donataire, description exacte du bien donné, documents justificatifs (le cas échéant)."}},
    "lawsuit": {
        "ai": {"ru": "Напиши суть спора и в какой суд обращаешься. Укажи, если знаешь: ФИО истца и ответчика, обстоятельства дела, какие требования заявляешь.",
               "en": "Write the essence of the dispute and which court you're filing with. Specify, if known: full names of the plaintiff and defendant, circumstances of the case, what demands you're making.",
               "de": "Schreibe das Wesen des Streits und bei welchem Gericht du klagst. Gib an, falls bekannt: Namen von Kläger und Beklagtem, Umstände des Falls, welche Forderungen du stellst.",
               "ar": "اكتب جوهر النزاع والمحكمة التي تتوجه إليها. حدد إن كنت تعرف: الاسم الكامل للمدعي والمدعى عليه، ملابسات القضية، المطالب التي تقدمها.",
               "zh": "请说明纠纷内容以及向哪个法院起诉。如知道的话，请注明原告和被告姓名、案件情况、诉讼请求。",
               "es": "Escribe la esencia del litigio y ante qué juzgado presentas la demanda. Indica, si lo sabes: nombre completo del demandante y del demandado, circunstancias del caso, qué exigencias planteas.",
               "fr": "Écris l'objet du litige et devant quel tribunal tu portes l'affaire. Indique, si tu les connais : noms complets du demandeur et du défendeur, circonstances de l'affaire, quelles demandes tu formules."},
        "user": {"ru": "Пришли данные для искового заявления: наименование суда, ФИО/данные истца и ответчика, обстоятельства дела, исковые требования, цену иска (если есть).",
                 "en": "Send the details for the statement of claim: name of the court, names/details of the plaintiff and defendant, circumstances of the case, the demands, the claim amount (if any).",
                 "de": "Sende die Daten für die Klageschrift: Name des Gerichts, Namen/Daten von Kläger und Beklagtem, Umstände des Falls, Klageforderungen, Streitwert (falls vorhanden).",
                 "ar": "أرسل بيانات صحيفة الدعوى: اسم المحكمة، اسم/بيانات المدعي والمدعى عليه، ملابسات القضية، طلبات الدعوى، قيمة الدعوى (إن وجدت).",
                 "zh": "请发送起诉状所需信息：法院名称、原告/被告姓名及信息、案件情况、诉讼请求、诉讼标的额（如有）。",
                 "es": "Envía los datos para la demanda: nombre del juzgado, nombre/datos del demandante y del demandado, circunstancias del caso, pretensiones, cuantía de la demanda (si la hay).",
                 "fr": "Envoie les données de l'assignation : nom du tribunal, noms/coordonnées du demandeur et du défendeur, circonstances de l'affaire, demandes, montant du litige (le cas échéant)."}},
    "alimony": {
        "ai": {"ru": "Напиши, кто и на кого платит алименты, и на каких условиях (сумма, периодичность).",
               "en": "Write who pays child/spousal support to whom, and on what terms (amount, frequency).",
               "de": "Schreibe, wer an wen Unterhalt zahlt und zu welchen Bedingungen (Betrag, Häufigkeit).",
               "ar": "اكتب من يدفع النفقة ولمن، وبأي شروط (المبلغ، الدورية).",
               "zh": "请说明抚养费的支付方和接收方，以及条件（金额、支付周期）。",
               "es": "Escribe quién paga la pensión alimenticia y a quién, y en qué condiciones (importe, periodicidad).",
               "fr": "Écris qui verse une pension alimentaire à qui, et selon quelles conditions (montant, périodicité)."},
        "user": {"ru": "Пришли данные для соглашения об алиментах: ФИО плательщика и получателя, ФИО и дату рождения ребёнка, размер и периодичность выплат.",
                 "en": "Send the details for the support agreement: full names of the payer and recipient, child's full name and date of birth, amount and frequency of payments.",
                 "de": "Sende die Daten für die Unterhaltsvereinbarung: Namen von Zahler und Empfänger, Name und Geburtsdatum des Kindes, Höhe und Häufigkeit der Zahlungen.",
                 "ar": "أرسل بيانات اتفاقية النفقة: الاسم الكامل للدافع والمستفيد، اسم وتاريخ ميلاد الطفل، مقدار الدفعات ودوريتها.",
                 "zh": "请发送抚养费协议所需信息：支付方和接收方姓名、孩子姓名及出生日期、支付金额和周期。",
                 "es": "Envía los datos para el acuerdo de alimentos: nombre completo del pagador y del beneficiario, nombre y fecha de nacimiento del menor, importe y periodicidad de los pagos.",
                 "fr": "Envoie les données de l'accord de pension alimentaire : noms complets du payeur et du bénéficiaire, nom et date de naissance de l'enfant, montant et périodicité des versements."}},
    "services": {
        "ai": {"ru": "Напиши, какую услугу и кто оказывает. Укажи, если знаешь: заказчика и исполнителя, суть услуги, стоимость и сроки.",
               "en": "Write what service and who provides it. Specify, if known: client and contractor, essence of the service, cost and deadlines.",
               "de": "Schreibe, welche Dienstleistung und von wem erbracht wird. Gib an, falls bekannt: Auftraggeber und Auftragnehmer, Art der Dienstleistung, Kosten und Fristen.",
               "ar": "اكتب ما الخدمة ومن يقدّمها. حدد إن كنت تعرف: العميل والمنفّذ، جوهر الخدمة، التكلفة والمواعيد.",
               "zh": "请说明服务内容及提供方。如知道的话，请注明委托方和服务方、服务内容、费用和期限。",
               "es": "Escribe qué servicio y quién lo presta. Indica, si lo sabes: cliente y prestador, esencia del servicio, coste y plazos.",
               "fr": "Écris quel service et qui le fournit. Indique, si tu les connais : client et prestataire, nature du service, coût et délais."},
        "user": {"ru": "Пришли данные для договора оказания услуг: заказчик и исполнитель, точное описание услуги, стоимость, порядок оплаты, сроки оказания.",
                 "en": "Send the details for the services agreement: client and contractor, exact description of the service, cost, payment terms, deadlines.",
                 "de": "Sende die Daten für den Dienstleistungsvertrag: Auftraggeber und Auftragnehmer, genaue Beschreibung der Dienstleistung, Kosten, Zahlungsbedingungen, Fristen.",
                 "ar": "أرسل بيانات عقد تقديم الخدمات: العميل والمنفّذ، وصفاً دقيقاً للخدمة، التكلفة، شروط الدفع، مواعيد التنفيذ.",
                 "zh": "请发送服务合同所需信息：委托方和服务方、服务的准确描述、费用、付款方式、服务期限。",
                 "es": "Envía los datos para el contrato de prestación de servicios: cliente y prestador, descripción exacta del servicio, coste, forma de pago, plazos de ejecución.",
                 "fr": "Envoie les données du contrat de prestation de services : client et prestataire, description exacte du service, coût, modalités de paiement, délais d'exécution."}},
    "employment": {
        "ai": {"ru": "Напиши должность и условия работы. Укажи, если знаешь: работодателя, работника, оклад, дату начала работы, режим работы.",
               "en": "Write the position and working conditions. Specify, if known: employer, employee, salary, start date, work schedule.",
               "de": "Schreibe die Position und die Arbeitsbedingungen. Gib an, falls bekannt: Arbeitgeber, Arbeitnehmer, Gehalt, Beginn der Tätigkeit, Arbeitszeit.",
               "ar": "اكتب المنصب وشروط العمل. حدد إن كنت تعرف: صاحب العمل، الموظف، الراتب، تاريخ بدء العمل، نظام العمل.",
               "zh": "请说明职位和工作条件。如知道的话，请注明雇主、员工、薪资、入职日期、工作制度。",
               "es": "Escribe el puesto y las condiciones de trabajo. Indica, si lo sabes: empleador, empleado, salario, fecha de inicio, horario de trabajo.",
               "fr": "Écris le poste et les conditions de travail. Indique, si tu les connais : employeur, employé, salaire, date de début, horaires de travail."},
        "user": {"ru": "Пришли данные для трудового договора: работодатель и работник (ФИО, паспортные данные), должность, оклад, режим работы, дату начала работы.",
                 "en": "Send the details for the employment contract: employer and employee (full name, ID details), position, salary, work schedule, start date.",
                 "de": "Sende die Daten für den Arbeitsvertrag: Arbeitgeber und Arbeitnehmer (Name, Ausweisdaten), Position, Gehalt, Arbeitszeit, Beginn der Tätigkeit.",
                 "ar": "أرسل بيانات عقد العمل: صاحب العمل والموظف (الاسم الكامل، بيانات الهوية)، المنصب، الراتب، نظام العمل، تاريخ بدء العمل.",
                 "zh": "请发送劳动合同所需信息：雇主和员工（姓名、证件信息）、职位、薪资、工作制度、入职日期。",
                 "es": "Envía los datos para el contrato de trabajo: empleador y empleado (nombre completo, datos de identificación), puesto, salario, horario de trabajo, fecha de inicio.",
                 "fr": "Envoie les données du contrat de travail : employeur et employé (nom complet, pièce d'identité), poste, salaire, horaires de travail, date de début."}},
    "work_act": {
        "ai": {"ru": "Напиши, какие работы или услуги выполнены и по какому договору. Укажи, если знаешь: заказчика и исполнителя, перечень работ, стоимость.",
               "en": "Write what work or services were performed and under which agreement. Specify, if known: client and contractor, list of works, cost.",
               "de": "Schreibe, welche Arbeiten oder Leistungen erbracht wurden und aufgrund welchen Vertrags. Gib an, falls bekannt: Auftraggeber und Auftragnehmer, Liste der Arbeiten, Kosten.",
               "ar": "اكتب الأعمال أو الخدمات المنجزة وبموجب أي عقد. حدد إن كنت تعرف: العميل والمنفّذ، قائمة الأعمال، التكلفة.",
               "zh": "请说明已完成的工作或服务内容及所依据的合同。如知道的话，请注明委托方和服务方、工作清单、费用。",
               "es": "Escribe qué trabajos o servicios se realizaron y en virtud de qué contrato. Indica, si lo sabes: cliente y prestador, listado de trabajos, coste.",
               "fr": "Écris quels travaux ou services ont été réalisés et en vertu de quel contrat. Indique, si tu les connais : client et prestataire, liste des travaux, coût."},
        "user": {"ru": "Пришли данные для акта выполненных работ: номер и дату договора, заказчик и исполнитель, перечень выполненных работ/услуг с объёмом и стоимостью.",
                 "en": "Send the details for the completion act: contract number and date, client and contractor, list of completed works/services with volume and cost.",
                 "de": "Sende die Daten für das Leistungsabnahmeprotokoll: Vertragsnummer und -datum, Auftraggeber und Auftragnehmer, Liste der erbrachten Arbeiten/Leistungen mit Umfang und Kosten.",
                 "ar": "أرسل بيانات محضر إنجاز الأعمال: رقم وتاريخ العقد، العميل والمنفّذ، قائمة الأعمال/الخدمات المنجزة مع الحجم والتكلفة.",
                 "zh": "请发送完工验收单所需信息：合同编号和日期、委托方和服务方、已完成工作/服务清单（含工作量和费用）。",
                 "es": "Envía los datos para el acta de trabajos realizados: número y fecha del contrato, cliente y prestador, listado de trabajos/servicios realizados con volumen y coste.",
                 "fr": "Envoie les données du procès-verbal de réception des travaux : numéro et date du contrat, client et prestataire, liste des travaux/services réalisés avec quantité et coût."}},
    "supply": {
        "ai": {"ru": "Напиши, какой товар поставляется и кем. Укажи, если знаешь: поставщика и покупателя, товар, количество, цену, сроки поставки.",
               "en": "Write what goods are supplied and by whom. Specify, if known: supplier and buyer, goods, quantity, price, delivery deadlines.",
               "de": "Schreibe, welche Ware geliefert wird und von wem. Gib an, falls bekannt: Lieferant und Käufer, Ware, Menge, Preis, Lieferfristen.",
               "ar": "اكتب ما هي البضاعة الموردة ومن يوردها. حدد إن كنت تعرف: المورد والمشتري، البضاعة، الكمية، السعر، مواعيد التوريد.",
               "zh": "请说明供应的货物及供应方。如知道的话，请注明供应商和买方、货物、数量、价格、交货期限。",
               "es": "Escribe qué mercancía se suministra y quién la suministra. Indica, si lo sabes: proveedor y comprador, mercancía, cantidad, precio, plazos de entrega.",
               "fr": "Écris quelle marchandise est fournie et par qui. Indique, si tu les connais : fournisseur et acheteur, marchandise, quantité, prix, délais de livraison."},
        "user": {"ru": "Пришли данные для договора поставки: поставщик и покупатель, наименование и количество товара, цена, сроки и порядок поставки и оплаты.",
                 "en": "Send the details for the supply agreement: supplier and buyer, name and quantity of goods, price, delivery and payment deadlines and terms.",
                 "de": "Sende die Daten für den Liefervertrag: Lieferant und Käufer, Bezeichnung und Menge der Ware, Preis, Liefer- und Zahlungsfristen und -bedingungen.",
                 "ar": "أرسل بيانات عقد التوريد: المورد والمشتري، اسم وكمية البضاعة، السعر، مواعيد وشروط التوريد والدفع.",
                 "zh": "请发送供货合同所需信息：供应商和买方、货物名称和数量、价格、供货和付款期限及方式。",
                 "es": "Envía los datos para el contrato de suministro: proveedor y comprador, denominación y cantidad de la mercancía, precio, plazos y condiciones de entrega y pago.",
                 "fr": "Envoie les données du contrat de fourniture : fournisseur et acheteur, désignation et quantité de la marchandise, prix, délais et modalités de livraison et de paiement."}},
    "agency": {
        "ai": {"ru": "Напиши, какие действия агент совершает в интересах принципала. Укажи, если знаешь: стороны, суть поручения, вознаграждение агента.",
               "en": "Write what actions the agent performs on behalf of the principal. Specify, if known: parties, essence of the assignment, agent's fee.",
               "de": "Schreibe, welche Handlungen der Agent im Interesse des Prinzipals ausführt. Gib an, falls bekannt: Parteien, Kern des Auftrags, Vergütung des Agenten.",
               "ar": "اكتب الإجراءات التي ينفذها الوكيل لصالح الموكِّل. حدد إن كنت تعرف: الطرفين، جوهر التكليف، أجر الوكيل.",
               "zh": "请说明代理人代表委托人执行的行为。如知道的话，请注明双方、委托事项内容、代理费。",
               "es": "Escribe qué acciones realiza el agente en interés del principal. Indica, si lo sabes: las partes, esencia del encargo, remuneración del agente.",
               "fr": "Écris quelles actions l'agent réalise dans l'intérêt du mandant. Indique, si tu les connais : les parties, nature du mandat, rémunération de l'agent."},
        "user": {"ru": "Пришли данные для агентского договора: принципал и агент, точное описание поручаемых действий, размер и порядок выплаты вознаграждения.",
                 "en": "Send the details for the agency agreement: principal and agent, exact description of the assigned actions, amount and payment terms of the fee.",
                 "de": "Sende die Daten für den Agenturvertrag: Prinzipal und Agent, genaue Beschreibung der beauftragten Handlungen, Höhe und Zahlungsweise der Vergütung.",
                 "ar": "أرسل بيانات عقد الوكالة: الموكِّل والوكيل، وصفاً دقيقاً للإجراءات المكلَّف بها، مقدار الأجر وطريقة دفعه.",
                 "zh": "请发送代理合同所需信息：委托人和代理人、委托事项的准确描述、代理费金额和支付方式。",
                 "es": "Envía los datos para el contrato de agencia: principal y agente, descripción exacta de las acciones encomendadas, importe y forma de pago de la remuneración.",
                 "fr": "Envoie les données du contrat d'agence : mandant et agent, description exacte des actions confiées, montant et modalités de versement de la rémunération."}},
    "joint_activity": {
        "ai": {"ru": "Напиши цель совместной деятельности и вклад каждой стороны.",
               "en": "Write the purpose of the joint activity and each party's contribution.",
               "de": "Schreibe das Ziel der gemeinsamen Tätigkeit und den Beitrag jeder Partei.",
               "ar": "اكتب هدف النشاط المشترك ومساهمة كل طرف.",
               "zh": "请说明共同经营的目的以及各方的出资。",
               "es": "Escribe el objetivo de la actividad conjunta y la aportación de cada parte.",
               "fr": "Écris l'objectif de l'activité commune et l'apport de chaque partie."},
        "user": {"ru": "Пришли данные для договора о совместной деятельности: стороны, цель, вклад каждого участника, порядок распределения прибыли и расходов.",
                 "en": "Send the details for the joint activity agreement: parties, purpose, each participant's contribution, profit and expense distribution terms.",
                 "de": "Sende die Daten für den Vertrag über gemeinsame Tätigkeit: Parteien, Ziel, Beitrag jedes Teilnehmers, Aufteilung von Gewinn und Kosten.",
                 "ar": "أرسل بيانات عقد النشاط المشترك: الأطراف، الهدف، مساهمة كل مشارك، طريقة توزيع الأرباح والمصاريف.",
                 "zh": "请发送共同经营合同所需信息：各方、目的、各参与方的出资、利润和费用分配方式。",
                 "es": "Envía los datos para el contrato de actividad conjunta: partes, objetivo, aportación de cada participante, forma de reparto de beneficios y gastos.",
                 "fr": "Envoie les données du contrat d'activité commune : parties, objectif, apport de chaque participant, modalités de répartition des bénéfices et des dépenses."}},
    "nonresidential_rent": {
        "ai": {"ru": "Напиши, какое нежилое помещение сдаётся в аренду. Укажи, если знаешь: стороны, адрес и площадь, срок аренды, сумму.",
               "en": "Write what non-residential premises are being leased. Specify, if known: parties, address and area, lease term, amount.",
               "de": "Schreibe, welche Gewerberäume vermietet werden. Gib an, falls bekannt: Parteien, Adresse und Fläche, Mietdauer, Betrag.",
               "ar": "اكتب أي عقار غير سكني يُؤجَّر. حدد إن كنت تعرف: الطرفين، العنوان والمساحة، مدة الإيجار، المبلغ.",
               "zh": "请说明出租的非居住用房产。如知道的话，请注明双方、地址和面积、租期、金额。",
               "es": "Escribe qué local no residencial se alquila. Indica, si lo sabes: partes, dirección y superficie, plazo del alquiler, importe.",
               "fr": "Écris quel local à usage non résidentiel est loué. Indique, si tu les connais : parties, adresse et surface, durée du bail, montant."},
        "user": {"ru": "Пришли данные для аренды нежилого помещения: арендодатель и арендатор, точный адрес и площадь, срок аренды, размер и порядок оплаты.",
                 "en": "Send the details for the non-residential lease: lessor and lessee, exact address and area, lease term, amount and payment terms.",
                 "de": "Sende die Daten für die Gewerbemiete: Vermieter und Mieter, genaue Adresse und Fläche, Mietdauer, Höhe und Zahlungsweise.",
                 "ar": "أرسل بيانات إيجار العقار غير السكني: المؤجر والمستأجر، العنوان الدقيق والمساحة، مدة الإيجار، المبلغ وطريقة الدفع.",
                 "zh": "请发送非居住用房产租赁所需信息：出租方和承租方、准确地址和面积、租期、金额和付款方式。",
                 "es": "Envía los datos para el alquiler del local no residencial: arrendador y arrendatario, dirección exacta y superficie, plazo del alquiler, importe y forma de pago.",
                 "fr": "Envoie les données de la location du local non résidentiel : bailleur et locataire, adresse exacte et surface, durée du bail, montant et modalités de paiement."}},
    "cession": {
        "ai": {"ru": "Напиши, какое право требования и по какому обязательству уступается.",
               "en": "Write what claim right and under which obligation is being assigned.",
               "de": "Schreibe, welcher Forderungsanspruch und aufgrund welcher Verpflichtung abgetreten wird.",
               "ar": "اكتب أي حق مطالبة وبموجب أي التزام يُحال.",
               "zh": "请说明转让的债权及其所依据的义务。",
               "es": "Escribe qué derecho de crédito y en virtud de qué obligación se cede.",
               "fr": "Écris quel droit de créance et en vertu de quelle obligation est cédé."},
        "user": {"ru": "Пришли данные для договора цессии: цедент и цессионарий, реквизиты первоначального обязательства, сумма и объём уступаемого права, цена уступки.",
                 "en": "Send the details for the assignment agreement: assignor and assignee, details of the original obligation, amount and scope of the assigned right, assignment price.",
                 "de": "Sende die Daten für den Abtretungsvertrag: Zedent und Zessionar, Angaben zur ursprünglichen Verpflichtung, Höhe und Umfang des abgetretenen Rechts, Abtretungspreis.",
                 "ar": "أرسل بيانات عقد الحوالة: المُحيل والمُحال إليه، بيانات الالتزام الأصلي، مبلغ ونطاق الحق المُحال، سعر الحوالة.",
                 "zh": "请发送债权转让合同所需信息：转让方和受让方、原债务信息、转让权利的金额和范围、转让价格。",
                 "es": "Envía los datos para el contrato de cesión: cedente y cesionario, datos de la obligación original, importe y alcance del derecho cedido, precio de la cesión.",
                 "fr": "Envoie les données du contrat de cession : cédant et cessionnaire, informations sur l'obligation d'origine, montant et étendue du droit cédé, prix de la cession."}},
    "nda": {
        "ai": {"ru": "Напиши, какую конфиденциальную информацию нужно защитить и между кем.",
               "en": "Write what confidential information needs protecting and between whom.",
               "de": "Schreibe, welche vertraulichen Informationen geschützt werden müssen und zwischen wem.",
               "ar": "اكتب أي معلومات سرية يجب حمايتها وبين من.",
               "zh": "请说明需要保护的机密信息以及涉及双方。",
               "es": "Escribe qué información confidencial hay que proteger y entre quiénes.",
               "fr": "Écris quelle information confidentielle doit être protégée et entre qui."},
        "user": {"ru": "Пришли данные для соглашения о неразглашении: стороны, что считается конфиденциальной информацией, срок действия обязательств, ответственность за разглашение.",
                 "en": "Send the details for the NDA: parties, what's considered confidential information, duration of obligations, liability for disclosure.",
                 "de": "Sende die Daten für die Vertraulichkeitsvereinbarung: Parteien, was als vertrauliche Information gilt, Dauer der Verpflichtungen, Haftung bei Offenlegung.",
                 "ar": "أرسل بيانات اتفاقية عدم الإفشاء: الطرفين، ما يُعتبر معلومات سرية، مدة الالتزامات، المسؤولية عن الإفشاء.",
                 "zh": "请发送保密协议所需信息：各方、机密信息的界定、义务期限、泄露责任。",
                 "es": "Envía los datos para el acuerdo de confidencialidad: partes, qué se considera información confidencial, duración de las obligaciones, responsabilidad por divulgación.",
                 "fr": "Envoie les données de l'accord de confidentialité : parties, ce qui est considéré comme information confidentielle, durée des obligations, responsabilité en cas de divulgation."}},
    "self_employed": {
        "ai": {"ru": "Напиши, какую работу выполняет самозанятый и для кого.",
               "en": "Write what work the self-employed person performs and for whom.",
               "de": "Schreibe, welche Arbeit der Selbstständige ausführt und für wen.",
               "ar": "اكتب أي عمل ينفذه العامل المستقل ولمن.",
               "zh": "请说明自由职业者从事的工作及服务对象。",
               "es": "Escribe qué trabajo realiza el autónomo y para quién.",
               "fr": "Écris quel travail réalise le travailleur indépendant et pour qui."},
        "user": {"ru": "Пришли данные для договора с самозанятым: заказчик и исполнитель (ФИО, ИНН самозанятого), суть работ, стоимость, сроки, порядок оплаты.",
                 "en": "Send the details for the self-employed contract: client and contractor (full name, tax ID of the self-employed person), essence of the work, cost, deadlines, payment terms.",
                 "de": "Sende die Daten für den Vertrag mit dem Selbstständigen: Auftraggeber und Auftragnehmer (Name, Steuernummer des Selbstständigen), Art der Arbeiten, Kosten, Fristen, Zahlungsweise.",
                 "ar": "أرسل بيانات عقد العامل المستقل: العميل والمنفّذ (الاسم الكامل، الرقم الضريبي للعامل المستقل)، جوهر العمل، التكلفة، المواعيد، طريقة الدفع.",
                 "zh": "请发送与自由职业者签订合同所需信息：委托方和服务方（姓名、自由职业者税号）、工作内容、费用、期限、付款方式。",
                 "es": "Envía los datos para el contrato con el autónomo: cliente y prestador (nombre completo, NIF del autónomo), esencia de los trabajos, coste, plazos, forma de pago.",
                 "fr": "Envoie les données du contrat avec le travailleur indépendant : client et prestataire (nom complet, numéro fiscal de l'indépendant), nature des travaux, coût, délais, modalités de paiement."}},
    "warranty_letter": {
        "ai": {"ru": "Напиши, что именно гарантируется и кому адресовано письмо.",
               "en": "Write exactly what's guaranteed and to whom the letter is addressed.",
               "de": "Schreibe, was genau garantiert wird und an wen der Brief gerichtet ist.",
               "ar": "اكتب ما الذي يُضمن بالضبط ولمن موجّهة الرسالة.",
               "zh": "请说明具体保证的内容以及信函的收件人。",
               "es": "Escribe qué se garantiza exactamente y a quién va dirigida la carta.",
               "fr": "Écris ce qui est exactement garanti et à qui la lettre est adressée."},
        "user": {"ru": "Пришли данные для гарантийного письма: кому адресовано, кто гарантирует, суть гарантии, срок исполнения.",
                 "en": "Send the details for the letter of guarantee: who it's addressed to, who's guaranteeing, essence of the guarantee, fulfillment deadline.",
                 "de": "Sende die Daten für das Garantieschreiben: an wen gerichtet, wer garantiert, Inhalt der Garantie, Erfüllungsfrist.",
                 "ar": "أرسل بيانات خطاب الضمان: الموجّه إليه، الضامن، جوهر الضمان، مهلة التنفيذ.",
                 "zh": "请发送保证函所需信息：收件人、担保方、保证内容、履行期限。",
                 "es": "Envía los datos para la carta de garantía: a quién va dirigida, quién garantiza, contenido de la garantía, plazo de cumplimiento.",
                 "fr": "Envoie les données de la lettre de garantie : à qui elle est adressée, qui garantit, nature de la garantie, délai d'exécution."}},
}

WORD_KIND_HINTS = {k: {"ai": v["ai"]["ru"], "user": v["user"]["ru"]} for k, v in WORD_KIND_HINTS_I18N.items()}  # обратная совместимость


def word_hint(kind: str, mode: str, lang: str = "ru") -> str:
    variants = WORD_KIND_HINTS_I18N.get(kind, WORD_KIND_HINTS_I18N["doc"])[mode]
    return variants.get(lang) or variants["ru"]


@dp.message(Form.waiting_word_mode, F.text.in_(ALL_BTN_AI_GENERATE_LABELS))
async def word_mode_ai(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await state.update_data(mode="ai", user_text="")
    await m.answer(tr("msg_content_lang_prompt", lang), reply_markup=content_lang_kb(lang))
    await state.set_state(Form.waiting_word_content_lang)


@dp.message(Form.waiting_word_content_lang)
async def word_content_lang(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    code = resolve_content_lang((m.text or "").strip(), lang)
    if not code:
        await m.answer(tr("msg_pick_content_lang", lang), reply_markup=content_lang_kb(lang))
        return
    await state.update_data(content_lang=code)
    data = await state.get_data()
    await m.answer(word_hint(data.get("word_kind", "doc"), "ai", lang), reply_markup=cancel_kb(lang))
    await state.set_state(Form.waiting_word_topic)


@dp.message(Form.waiting_word_mode, F.text.in_(ALL_BTN_OWN_TEXT_LABELS))
async def word_mode_user(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await state.update_data(mode="user")
    data = await state.get_data()
    await m.answer(word_hint(data.get("word_kind", "doc"), "user", lang), reply_markup=cancel_kb(lang))
    await state.set_state(Form.waiting_word_text)


@dp.message(Form.waiting_word_mode, F.text.in_(ALL_BTN_TEMPLATE_LABELS))
async def word_mode_template(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    kind = data.get("word_kind", "doc")
    uid = m.from_user.id
    if kind in STUDY_KINDS:
        await m.answer(tr("msg_no_template_study", lang), reply_markup=mode_kb(False, lang=lang))
        return
    ok, reason = start_job(uid)
    if not ok:
        await m.answer(reason)
        return
    await m.answer(tr("msg_building_template", lang), reply_markup=cancel_kb(lang))
    try:
        u = get_user(uid)
        tpl_path = f"/tmp/tpl_{uid}.docx"
        try:
            meta = json.loads(META_SCHEMAS.get(kind, META_SCHEMAS["doc"]))
        except Exception:
            meta = {}
        build_word(tpl_path, "Документ", template_sections_for(kind), kind, meta)
        await m.answer_document(FSInputFile(tpl_path), caption=tr("msg_template_caption", lang))
        try:
            os.remove(tpl_path)
        except Exception:
            pass
        u["generations"] += 1
        u["history"].append(f"{datetime.now().strftime('%d.%m %H:%M')} — шаблон: {KIND_LABELS.get(kind, kind)}")
        note_success(uid)
        await m.answer(tr("msg_ready", lang) + "\n\n" + await signature_line(lang), reply_markup=main_kb(lang))
        await state.clear()
    finally:
        finish_job(uid)


@dp.message(Form.waiting_word_mode)
async def waiting_word_mode_fallback(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    show_template = data.get("word_kind", "doc") not in STUDY_KINDS
    await m.answer(tr("msg_didnt_understand", lang), reply_markup=mode_kb(show_template, lang=lang))


# Выбор объёма (короткий/средний/подробный) имеет смысл только там, где объём реально
# влияет на глубину раскрытия темы: реферат, курсовая, эссе, доклад. Для конспекта
# он не нужен — конспект по своей природе всегда сжатый и короткий. Для обычного
# документа тоже убран — там слишком неопределённый формат, чтобы объём был осмысленным
# параметром. Для юридических документов (договоров, доверенностей, актов и т.п.) объём
# определяется их обязательной структурой по ГК РФ, а не пожеланием пользователя.
WORD_SIZE_KINDS = {"referat", "report", "essay", "coursework"}


async def word_after_input(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    kind = data.get("word_kind", "doc")
    if kind in WORD_SIZE_KINDS:
        await m.answer(tr("msg_which_size", lang), reply_markup=word_size_kb(lang))
        await state.set_state(Form.waiting_word_size)
    else:
        await word_build_draft(m, state, data, "short")


@dp.message(Form.waiting_word_topic)
async def word_topic(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    text = m.text or ""
    if len(text) > 500:
        await m.answer(tr("msg_too_long_500", lang))
        return
    await state.update_data(topic=text)
    await word_after_input(m, state)


@dp.message(Form.waiting_word_text)
async def word_user_text(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    text = m.text or ""
    if len(text) < 30:
        await m.answer(tr("msg_too_little_text_doc", lang))
        return
    if len(text) > 4000:
        await m.answer(tr("msg_too_long_4000", lang))
        return
    await state.update_data(user_text=text, topic=text[:80].replace("\n", " "))
    await word_after_input(m, state)


async def word_build_draft(m: Message, state: FSMContext, data: dict, size: str):
    lang = user_lang(m.from_user.id)
    await state.update_data(word_size=size)
    await m.answer(tr("msg_building_draft", lang), reply_markup=cancel_kb(lang))
    size_map = {"short": "поверхностное раскрытие темы — только суть и ключевые моменты, без глубокого разбора деталей и подпунктов, но по-настоящему содержательно, объём определяй по теме, не режь искусственно", "long": "полное раскрытие темы — подробно, с деталями, подпунктами и глубоким разбором, объём определяй по теме"}
    kind = data.get("word_kind", "doc")
    kind_name = WORD_KIND_DESC.get(kind, "документ")
    lang_instr = grok_lang_instruction(content_gen_lang(data, lang) or "ru")
    prompt = f"""Собери черновик-план (не полный текст): {kind_name}.
Тема/данные: {data.get('topic')}
Текст пользователя: {data.get('user_text')}
Доп: {data.get('extra')}
Итоговый документ будет иметь такой объём: {size_map[size]}. Но сейчас нужен не сам документ,
а короткий план для согласования с пользователем.
Исправь ошибки. {"Пиши формальным юридическим/деловым языком, без канцелярита-воды." if kind not in STUDY_KINDS else ANTI_AI_DETECTOR_STYLE}
Важно: никогда не придумывай паспортные данные, суммы, даты или ФИО, которых нет в данных пользователя -
для них ставь пропуски [указать ...]. Если документ содержательный (реферат, коммерческое предложение) -
план должен отражать суть темы конкретно, а не общими фразами.
Обычным текстом, коротко: название и 4 пункта структуры (каждый пункт — одна строка, без раскрытия
содержания). Без JSON.{lang_instr}"""
    sample = await ask_grok(prompt)
    if grok_failed(sample):
        await m.answer(
            tr("msg_grok_error", lang),
            reply_markup=word_size_kb(lang) if kind in WORD_SIZE_KINDS else word_confirm_kb(lang)
        )
        return
    await state.update_data(sample=sample)
    await send_draft(m, tr("msg_draft_ready_doc", lang, sample=sample), title=data.get("topic") or kind_name, reply_markup=word_confirm_kb(lang))
    await state.set_state(Form.waiting_word_confirm)


@dp.message(Form.waiting_word_size)
async def word_size(m: Message, state: FSMContext):
    t = (m.text or "").lower()
    size = "long" if "полн" in t or "full" in t or "voll" in t or "完整" in t or "الكامل" in t or "completo" in t or "complet" in t else "short"
    data = await state.get_data()
    await word_build_draft(m, state, data, size)


@dp.message(Form.waiting_word_confirm, F.text.in_(ALL_BTN_CHANGE_QUERY_LABELS))
async def word_change(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    if data.get("mode") == "user":
        await m.answer(tr("msg_new_text", lang), reply_markup=cancel_kb(lang))
        await state.set_state(Form.waiting_word_text)
    else:
        await m.answer(tr("msg_new_topic", lang), reply_markup=cancel_kb(lang))
        await state.set_state(Form.waiting_word_topic)


@dp.message(Form.waiting_word_confirm, F.text.in_(ALL_BTN_ADD_INFO_LABELS))
async def word_add(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    if (await state.get_data()).get("extra_used", 0) >= 3:
        await m.answer(tr("msg_extra_limit", lang), reply_markup=word_confirm_kb(lang))
        return
    await m.answer(tr("msg_extra_prompt", lang), reply_markup=cancel_kb(lang))
    await state.set_state(Form.waiting_word_extra)


@dp.message(Form.waiting_word_extra)
async def word_extra(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    text = m.text or ""
    if len(text) > 800:
        await m.answer(tr("msg_too_long_800", lang))
        return
    extra = ((data.get("extra") or "") + "\n" + text).strip()
    await state.update_data(extra=extra, extra_used=data.get("extra_used", 0) + 1)
    await m.answer(tr("msg_updating_draft", lang))
    lang_instr = grok_lang_instruction(content_gen_lang(data, lang) or "ru")
    sample = await ask_grok(
        f"Обнови черновик документа.\nТема: {data.get('topic')}\nТекст: {data.get('user_text')}\nДоп: {extra}\nКороткий план. Без JSON.{lang_instr}"
    )
    if grok_failed(sample):
        await m.answer(
            tr("msg_draft_update_failed", lang),
            reply_markup=word_confirm_kb(lang)
        )
        await state.set_state(Form.waiting_word_confirm)
        return
    await state.update_data(sample=sample)
    await send_draft(m, tr("msg_draft_updated", lang, sample=sample), title=data.get("topic") or WORD_KIND_DESC.get(data.get("word_kind", "doc"), "документ"), reply_markup=word_confirm_kb(lang))
    await state.set_state(Form.waiting_word_confirm)


@dp.message(Form.waiting_word_confirm, F.text.in_(ALL_BTN_BUILD_DOC_LABELS))
async def word_build(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    data = await state.get_data()
    uid = m.from_user.id
    u = get_user(uid)
    ok, reason = start_job(uid)
    if not ok:
        await m.answer(reason)
        return
    await m.answer(tr("msg_building_word", lang))
    kind = data.get("word_kind", "doc")
    size = data.get("word_size", "short")
    lang_instr = grok_json_lang_instruction(content_gen_lang(data, lang) or "ru")
    size_map = {"short": "поверхностное раскрытие темы — только суть и ключевые моменты, без глубокого разбора деталей и подпунктов, но по-настоящему содержательно, объём определяй по теме, не режь искусственно", "long": "полное раскрытие темы — подробно, с деталями, подпунктами и глубоким разбором, объём определяй по теме"}
    kind_name = WORD_KIND_DESC.get(kind, "документ")
    # Схема мета-полей своя под каждый тип документа (модульный словарь META_SCHEMAS,
    # он же переиспользуется в режиме "Скачать шаблон") — раньше запрашивалась
    # всегда одна и та же (реферат/договор), из-за чего для заявления,
    # доверенности, КП поля вроде "кому"/"от кого" не заполнялись моделью
    # и в документе оставались заглушки, даже если пользователь их указал.
    meta_schema = META_SCHEMAS.get(kind, META_SCHEMAS["doc"])
    style_rule = ANTI_AI_DETECTOR_STYLE if kind in STUDY_KINDS else (
        "Пиши формальным юридическим/деловым языком, грамотно и точно по формулировкам ГК РФ, где применимо."
    )
    # Финальная сборка — это не короткий черновик-план, а весь текст документа целиком,
    # поэтому фиксированного лимита в 4000 токенов не хватало на курсовую/реферат с
    # "полным раскрытием темы" (там нужно 6000-10000+ слов), и модель тихо обрезала
    # содержание. Лимит теперь подбирается по типу документа и выбранной глубине.
    # Русский текст обычно занимает 1.5-2 токена на слово, плюс накладные расходы
    # на JSON-обёртку (кавычки, экранирование, ключи полей) - лимит с запасом,
    # чтобы модель физически не упёрлась в потолок на середине последнего раздела.
    if kind == "coursework":
        gen_max_tokens = 20000 if size == "long" else 8000
    elif kind == "referat":
        gen_max_tokens = 12000 if size == "long" else 6000
    elif kind in WORD_SIZE_KINDS:  # report, essay
        gen_max_tokens = 7000 if size == "long" else 4000
    else:
        gen_max_tokens = 6000
    # Расплывчатые формулировки вроде "несколько содержательных абзацев" модель
    # игнорирует и всё равно пишет коротко, даже с большим лимитом токенов -
    # нужны точные числовые ориентиры по объёму на раздел, иначе она работает
    # по привычке писать компактно, независимо от того, сколько токенов доступно.
    length_hint = ""
    if kind == "coursework":
        if size == "long":
            length_hint = "\nЭто полноценная курсовая работа для сдачи, суммарный объём всего документа — 6000-9000 слов. Каждый содержательный раздел (кроме титульного листа, содержания и списка литературы) должен быть НЕ МЕНЕЕ 600-900 слов (это примерно 4-6 полноценных абзацев с фактами, примерами, анализом, а не общими фразами) - пиши подробно и разворачивай мысль, а не сжимай её в 2-3 предложения."
        else:
            length_hint = "\nСуммарный объём документа — 2000-3000 слов. Каждый содержательный раздел (кроме титульного листа, содержания и списка литературы) — примерно 200-300 слов, по существу, без искусственного разжижения."
    elif kind == "referat":
        if size == "long":
            length_hint = "\nЭто полноценный реферат для сдачи, суммарный объём всего документа — 3000-5000 слов. Каждый содержательный раздел — НЕ МЕНЕЕ 400-600 слов (несколько развёрнутых абзацев с фактами и анализом), не сжимай в 2-3 предложения."
        else:
            length_hint = "\nСуммарный объём документа — 1200-1800 слов. Каждый содержательный раздел — примерно 150-250 слов, по существу."
    elif kind in WORD_SIZE_KINDS:  # report, essay
        if size == "long":
            length_hint = "\nСуммарный объём документа — 1500-2500 слов. Каждый содержательный раздел — примерно 250-400 слов, развёрнуто, с конкретикой."
        else:
            length_hint = "\nСуммарный объём документа — 600-1000 слов. Каждый содержательный раздел — примерно 100-180 слов, по существу, без воды."
    # Модель склонна писать заключение как краткий пересказ глав и не проверять
    # число источников - это отдельная, часто игнорируемая инструкция, поэтому
    # прописываем её явно, а не полагаемся на общее описание вида документа.
    structure_hint = ""
    if kind in ("coursework", "referat"):
        structure_hint = (
            "\nПервым разделом в sections обязательно должен идти раздел «Содержание» - "
            "перечисли в его content все остальные разделы и подразделы документа по одному на строке "
            "(без номеров страниц, они не нужны в черновике). Не пропускай этот раздел.\n"
            "Заключение пиши как самостоятельные выводы по каждой задаче из введения "
            "и рекомендации по теме - не пересказывай содержание глав.\n"
            "Раздел со списком литературы должен содержать не менее 5 источников "
            "(автор/название/год/издательство или ссылка), оформленных по ГОСТ.\n"
            "ОБЯЗАТЕЛЬНО добавь хотя бы одну таблицу (ключ \"table\", см. схему ниже) в практическую/аналитическую "
            "главу - это часть требований к курсовой/реферату, а не опция. Собери в неё конкретные цифры по теме: "
            "результаты опроса, статистику, сравнение показателей до/после, данные по годам - придумай "
            "правдоподобные иллюстративные значения, если реальных нет под рукой (как и с остальным текстом "
            "черновика). Это требование не выполняется текстовым описанием цифр вместо таблицы."
        )
    raw = await ask_grok(f"""Собери Word: {kind_name}.
Данные: {data.get('topic')}
Текст пользователя: {data.get('user_text')}
Доп: {data.get('extra')}
Объём: {size_map[size]}{length_hint}{structure_hint}
{style_rule}
Если в данных пользователя есть даты, ФИО, паспортные данные, суммы, названия сторон - подставь их в meta и в текст
точно как есть, ничего не меняя и не придумывая.
Никогда не выдумывай паспортные данные, суммы, даты или ФИО, которых нет в данных пользователя -
для недостающих данных ставь [указать ...].
Для содержательных документов (реферат, коммерческое предложение, заявление) разделы должны раскрывать
тему конкретно и по существу, без воды и общих фраз.
Не создавай в sections отдельный раздел «Титульный лист» - обложка (организация, учебное заведение, ФИО
автора и руководителя, город, год) формируется отдельно из title и meta и так уже попадёт в документ;
первым разделом в sections должно идти «Введение» (или «Содержание», если оно есть в структуре документа).
Если в разделе уместна таблица (перечень товаров/имущества с ценой и количеством в договоре/акте, статистика
или практические данные в курсовой/реферате) - добавь в объект этого раздела ключ "table":
{{"headers":["...","..."],"rows":[["...","..."],["...","..."]]}}. Не выдумывай реальные официальные цифры
(даты, номера, суммы) для юридических документов - только то, что дал пользователь; для курсовых/рефератов
данные в таблице могут быть иллюстративными для черновика, как и остальной текст. Таблицы уместны не в каждом
разделе - добавляй только там, где это реально яснее текста, не в каждый раздел подряд.
Только JSON:
{{"title":"...","meta":{meta_schema},"sections":[{{"title":"Введение","content":"абзац1\n\nабзац2"}},{{"title":"...","content":"...","table":{{"headers":["...","..."],"rows":[["...","..."]]}}}}]}}{lang_instr}""", max_tokens=gen_max_tokens)
    try:
        content = extract_json(raw)
        if not isinstance(content.get("sections"), list) or not content["sections"]:
            raise ValueError("В ответе модели нет разделов документа")
    except Exception as e:
        print("Word JSON parse error:", e)
        finish_job(uid)
        await m.answer(tr("msg_no_text", lang, btn=WORD_KIND_LABELS.get(kind, {}).get(lang) or tr("msg_which_doc", lang)), reply_markup=main_kb(lang))
        await state.clear()
        return

    # Числовые ориентиры по объёму в промпте модель выполняет непоследовательно -
    # может уложиться в цель по одному разделу и заметно недобрать по другим,
    # так что итоговый документ всё равно выходит в разы короче нужного. Вместо
    # того чтобы полагаться на то, что промпт сработает с первого раза, меряем
    # фактический объём и, если он далеко от цели, просим модель дописать черновик
    # подробнее - раздел за разделом, с явной цифрой по каждому разделу отдельно
    # (общая просьба "дописать подробнее" на практике не даёт нужного прироста).
    min_total = WORD_MIN_TOTAL_WORDS.get((kind, size))
    section_min = WORD_SECTION_MIN_WORDS.get((kind, size))
    if min_total:
        actual_words = sum(len((b.get("content") or "").split()) for b in content.get("sections", []))
        attempts = 0
        while actual_words < min_total * 0.7 and attempts < 2:
            attempts += 1
            msg = tr("msg_writing_more", lang) if attempts == 1 else tr("msg_writing_a_bit_more", lang)
            await m.answer(msg)
            orig_sections = content.get("sections", [])
            targets_text = "\n".join(
                f"{i}. «{(b.get('title') or '').strip()}» — сейчас {len((b.get('content') or '').split())} слов, "
                + ("оставь как есть." if any(k in (b.get('title') or '').lower() for k in _WORD_SECTION_SKIP_EXPAND)
                   else f"нужно не менее {section_min} слов.")
                for i, b in enumerate(orig_sections, 1)
            )
            expand_raw = await ask_grok(f"""Вот черновик документа в JSON, он слишком короткий: сейчас примерно {actual_words} слов, а нужно суммарно не менее {min_total}.
Раздел за разделом, вот что нужно дописать:
{targets_text}
В документе должно остаться РОВНО {len(orig_sections)} разделов, с теми же заголовками, в том же порядке -
не удаляй, не объединяй, не переименовывай и не убирай ни один раздел (включая "Содержание" и заголовки глав
вроде "Глава 1"), только доработай содержимое (content) там, где указана недостача.
Добавляй конкретные факты, примеры, аргументы, анализ, детали по теме - без воды и без повторов одной мысли разными словами.
Это критично: не пересказывай своими словами то, что уже написано в разделе, и не добавляй общие рассуждения
не по существу - это одновременно снижает уникальность текста (антиплагиат видит шаблонные обороты как
совпадения с другими работами) и делает текст более узнаваемым как сгенерированный. Расширяй раздел только
за счёт нового содержания: дополнительные примеры, цифры, источники, конкретные детали по теме, которых
раньше не было.
{style_rule}
Верни ТОЛЬКО JSON той же структуры, без пояснений:
{{"title":"...","meta":{meta_schema},"sections":[{{"title":"...","content":"..."}}]}}{lang_instr}
Черновик:
{json.dumps(content, ensure_ascii=False)}""", max_tokens=gen_max_tokens)
            try:
                expanded = extract_json(expand_raw)
                expanded_sections = expanded.get("sections")
                expanded_words = sum(len((b.get("content") or "").split()) for b in expanded_sections or [])
                # Принимаем результат только если он не только длиннее, но и не потерял
                # разделы (модель иногда "экономит" объём, просто слив несколько
                # разделов в один или выкинув служебные заголовки при переписывании).
                if (isinstance(expanded_sections, list) and expanded_sections
                        and len(expanded_sections) >= len(orig_sections)
                        and expanded_words > actual_words):
                    content = expanded
                    actual_words = expanded_words
                else:
                    break  # результат хуже или структурно испорчен - не продолжаем попытки
            except Exception as e:
                print("Word expand parse error:", e)
                break

    try:
        docx_path = f"/tmp/doc_{uid}.docx"
        build_word(
            docx_path,
            content.get("title", "Документ"),
            content.get("sections", []),
            kind,
            content.get("meta") or {}
        )

        pdf_path = f"/tmp/doc_{uid}.pdf"
        pdf = canvas.Canvas(pdf_path, pagesize=A4)
        w, h = A4
        fn, fb = register_pdf_fonts()

        pdf.setFont(fb, 14)
        pdf.drawString(40, h - 50, content.get("title", "")[:65])
        y = h - 80
        for s in content.get("sections", []):
            if y < 70:
                pdf.showPage()
                y = h - 50
            pdf.setFont(fb, 11)
            pdf.drawString(40, y, (s.get("title") or "")[:70])
            y -= 16
            pdf.setFont(fn, 9)
            for line in wrap_lines((s.get("content") or ""), fn, 9, w - 80):
                if y < 50:
                    pdf.showPage()
                    y = h - 50
                    pdf.setFont(fn, 9)
                pdf.drawString(40, y, line)
                y -= 12
            y -= 10
        pdf.save()

        fname = safe_filename(content.get("title"), fallback=WORD_KIND_DESC.get(kind, "Документ"), author=(content.get("meta") or {}).get("author"))
        await m.answer_document(FSInputFile(docx_path, filename=f"{fname}.docx"), caption=tr("msg_word_caption", lang))
        await m.answer_document(FSInputFile(pdf_path, filename=f"{fname}.pdf"), caption=tr("msg_pdf_caption", lang))
        u["generations"] += 1
        u["history"].append(f"{datetime.now().strftime('%d.%m %H:%M')} — {content.get('title')}")
        note_success(uid)
        for p in (docx_path, pdf_path):
            try:
                os.remove(p)
            except Exception:
                pass
        final_msg = tr("msg_doc_ready", lang)
        if kind not in STUDY_KINDS:
            final_msg += tr("msg_check_gaps", lang)
        final_msg += "\n\n" + await signature_line(lang)
        await m.answer(final_msg, reply_markup=main_kb(lang))
        await state.clear()
    finally:
        finish_job(uid)


@dp.message(Form.waiting_word_confirm)
async def waiting_word_confirm_fallback(m: Message, state: FSMContext):
    lang = user_lang(m.from_user.id)
    await m.answer(
        tr("msg_didnt_understand", lang),
        reply_markup=word_confirm_kb(lang)
    )


@dp.message(F.text.in_(ALL_BTN_HISTORY_LABELS))
async def history(m: Message):
    lang = user_lang(m.from_user.id)
    u = get_user(m.from_user.id)
    await m.answer(tr("msg_history_empty", lang) if not u["history"] else tr("msg_history_title", lang) + "\n".join(u["history"][-10:]))


@dp.message(F.text.in_(ALL_BTN_PLAN_LABELS))
async def my_plan(m: Message):
    lang = user_lang(m.from_user.id)
    u = get_user(m.from_user.id)
    limit = PLAN_LIMITS.get(u["plan"], 15)
    await m.answer(tr("msg_plan_info", lang, used=u['generations'], limit=limit, left=max(0, limit - u['generations'])))


@dp.message(Command("grant"))
async def grant(m: Message):
    if m.from_user.id not in ADMIN_IDS:
        return
    try:
        uid = int(m.text.split()[1])
        u = get_user(uid)
        u["plan"] = "premium"
        u["generations"] = 0
        save_users()
        await m.answer(f"Доступ выдан пользователю {uid}, счётчик генераций сброшен")
    except (IndexError, ValueError):
        await m.answer("Формат: /grant user_id")


@dp.errors()
async def global_error_handler(event):
    """Подстраховка на случай непредвиденных ошибок вне двух основных
    обработчиков (те уже защищены через try/finally выше и сами снимают
    busy). Здесь — просто не даём боту молчать и на всякий случай ещё раз
    снимаем блокировку, если получится определить пользователя."""
    print("Необработанная ошибка:", repr(event.exception))
    uid = None
    update = getattr(event, "update", None)
    msg = getattr(update, "message", None) or getattr(update, "callback_query", None)
    if msg is not None:
        user = getattr(msg, "from_user", None)
        uid = getattr(user, "id", None)
    if uid is not None:
        try:
            finish_job(uid)
            await bot.send_message(uid, "Что-то пошло не так. Попробуй ещё раз или напиши /cancel.")
        except Exception as e:
            print("Не смог уведомить пользователя об ошибке:", e)
    return True


async def main():
    print("Бот запущен")
    print("REPLICATE TOKEN:", "YES" if REPLICATE_API_TOKEN else "NO")
    await dp.start_polling(bot)


if __name__ == "__main__":
    asyncio.run(main())
